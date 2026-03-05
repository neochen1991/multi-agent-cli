from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

OUT_DIR = Path('/Users/neochen/multi-agent-cli_v2/output')
PPT_PATH = OUT_DIR / 'presentation.pptx'

slides = [
    {
        'title': '多 Agent RCA 平台',
        'subtitle': 'Code Wiki 到 PPT 的系统化讲解',
        'bullets': [
            '主题：代码架构、调度机制、Skill 能力与扩展方法',
            '受众：新加入项目的研发/SRE/测试同学',
            '目标：20 分钟内建立可落地的代码认知地图',
        ],
    },
    {
        'title': '1. 项目目标：可控、可解释、可回放',
        'bullets': [
            '这是“多 Agent 协作排障系统”，不是单轮对话机器人',
            'LLM 负责推理，系统负责流程与治理',
            '可控：主 Agent 先下命令，专家后执行',
            '可解释：工具调用与 Skill 命中可审计',
            '可回放：会话事件、断点与结论可恢复',
        ],
    },
    {
        'title': '2. 代码分层架构（从前端到运行时）',
        'bullets': [
            'Frontend：Incident / Settings 页面承接交互与配置',
            'API：/debates、/settings/tooling、WS 实时流',
            'Services：debate_service、agent_tool_context_service、agent_skill_service',
            'Runtime：LangGraph Orchestrator + Nodes + Routing + Execution',
            'Persistence：session store、lineage、tooling config',
        ],
    },
    {
        'title': '3. 端到端执行链路（一次会话）',
        'bullets': [
            '用户在 Incident 发起分析',
            'DebateService 做资产采集与上下文压缩',
            'Orchestrator 启动 LangGraph 执行',
            'Agent 执行：命令 -> 工具/Skill -> LLM -> 结构化输出',
            '事件流实时推送前端，最终落地报告',
        ],
    },
    {
        'title': '4. LangGraph 图：核心节点拓扑',
        'bullets': [
            '主链：init_session -> round_start -> supervisor_decide',
            '执行节点：analysis_parallel / collaboration / speak:agent',
            '收敛链：round_evaluate -> (round_start | finalize)',
            '关键状态：next_step、continue_next_round',
            '协作节点受 DEBATE_ENABLE_COLLABORATION 开关控制',
        ],
    },
    {
        'title': '5. 调度机制：HybridRouter 怎么做决定',
        'bullets': [
            'Stage 1：Seeded（主 Agent 预置步骤）',
            'Stage 2：Consensus shortcut（Judge 高置信收敛）',
            'Stage 3：覆盖后强制 Critic/Rebuttal/Judge',
            'Stage 4：预算保护，步数超限回退规则路由',
            'Stage 5/6：动态 LLM 路由 + 异常兜底 guardrail',
        ],
    },
    {
        'title': '6. Agent 体系：13 个角色分工',
        'bullets': [
            'Coordination：ProblemAnalysisAgent',
            'Analysis：Log/Domain/Code/Database/Metrics/Change/Runbook/RuleSuggestion',
            'Critique：Critic -> Rebuttal',
            'Judgment：Judge',
            'Verification：Verification',
        ],
    },
    {
        'title': '7. 单个 Agent 的运行流水线（最重要）',
        'bullets': [
            '接收命令：task/focus/expected_output/use_tool/skill_hints',
            '命令门禁：command_gate 判定是否允许工具/Skill',
            '工具上下文：日志、代码、数据库、指标等',
            'Skill 上下文：hints 或文本匹配命中',
            'LLM 调用 + 结构化归一化 + 状态回写',
        ],
    },
    {
        'title': '8. Skill 能力：配置入口与生效路径',
        'bullets': [
            '配置模型：enabled/skills_dir/max_skills/max_skill_chars/allowed_agents',
            '前端入口：/settings -> Agent Skill Router 配置卡片',
            '后端接口：GET/PUT /api/v1/settings/tooling',
            '持久化：tooling_config.json',
            '运行时：tooling_service.get_config() -> skill select',
        ],
    },
    {
        'title': '9. Skill 选择逻辑：如何命中',
        'bullets': [
            '优先级 1：命令显式 skill_hints',
            '优先级 2：运行时自动补 hints（按 agent + incident 信号）',
            '优先级 3：文本匹配打分（task/focus/expected_output + triggers）',
            '前置条件：enabled + allowed + has_command + allow_tool',
            '产出：skill_context + agent_skill_router 审计',
        ],
    },
    {
        'title': '10. 可靠性治理：系统如何避免卡死与乱跑',
        'bullets': [
            'LLM 调用治理：超时计划、重试、队列控制、fallback turn',
            '调度治理：规则引擎防重复、防超预算、防低信号循环',
            '终态治理：会话必须进入 completed/failed/cancelled',
            '结论治理：可配置拒绝占位结论门禁',
            '可恢复治理：events/sessions/tasks + WS resume/snapshot',
        ],
    },
    {
        'title': '11. 审计与可观测：如何追踪一次分析',
        'bullets': [
            '关键事件：agent_command_issued、tool_context_prepared、tool_io、feedback、supervisor_decision',
            '可追踪：谁下命令、谁调工具、命中什么 Skill、如何收敛',
            'Incident 页面可直接查看 Skill 路由与工具审计卡片',
        ],
    },
    {
        'title': '12. 如何扩展新 Agent（工程路径）',
        'bullets': [
            '更新协议文档（agent-catalog / protocol-contracts）',
            'specs.py 增 AgentSpec，builder/routing 增映射',
            'agent_tool_context_service 增上下文分支',
            '新增对应 Skill 并放开 allowed_agents',
            '补测试与回归验证',
        ],
    },
    {
        'title': '13. 如何扩展新 Skill（工程路径）',
        'bullets': [
            '创建 backend/skills/skill-name/SKILL.md',
            '写 front matter：name/description/triggers/agents',
            '写 Goal / Checklist / Output Contract',
            '命令中加 skill_hints 验证命中',
            '在 Incident 页面确认 agent_skill_router 审计',
        ],
    },
    {
        'title': '14. 总结与落地建议',
        'bullets': [
            '四层流水线：主 Agent 指挥 -> 工具/Skill -> LLM -> 系统治理',
            '新人建议先读：langgraph_runtime + agent_tool_context_service',
            '再读：agent_skill_service + settings 配置链路',
            '最后结合真实 session 事件流做一次端到端追踪',
        ],
    },
]


def parse_notes(notes_text: str):
    notes_map = {}
    current = None
    buf = []
    for raw in notes_text.splitlines():
        line = raw.strip()
        if line.startswith('## Slide '):
            if current is not None:
                notes_map[current] = '\n'.join([x for x in buf if x]).strip()
            try:
                current = int(line.replace('## Slide ', '').strip())
            except ValueError:
                current = None
            buf = []
            continue
        if current is not None:
            buf.append(raw)
    if current is not None:
        notes_map[current] = '\n'.join([x for x in buf if x]).strip()
    return notes_map


def style_title(shape):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.name = 'PingFang SC'
    p.font.color.rgb = RGBColor(31, 41, 55)


def style_bullets(shape):
    tf = shape.text_frame
    for i, p in enumerate(tf.paragraphs):
        if i == 0 and not p.text.strip():
            continue
        p.level = 0
        p.font.size = Pt(20)
        p.font.name = 'PingFang SC'
        p.font.color.rgb = RGBColor(31, 41, 55)


def add_slide(prs, title, bullets, note):
    layout = prs.slide_layouts[1]  # title and content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
    style_bullets(body)

    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note


if __name__ == '__main__':
    notes_text = (OUT_DIR / 'notes.md').read_text(encoding='utf-8')
    notes_map = parse_notes(notes_text)

    prs = Presentation()
    prs.slide_width = 13_333_333
    prs.slide_height = 7_500_000

    for idx, item in enumerate(slides, start=1):
        add_slide(
            prs,
            title=item['title'],
            bullets=item['bullets'],
            note=notes_map.get(idx, ''),
        )

    prs.save(PPT_PATH)
    print(f'generated: {PPT_PATH}')
