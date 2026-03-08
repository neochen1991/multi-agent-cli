#!/usr/bin/env python3
"""generate项目介绍PPTpro脚本。"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Iterable, List, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTFILE = Path(
    "/Users/neochen/multi-agent-cli_v2/plans/2026-03-04-生产问题根因分析系统-项目介绍-专业版.pptx"
)


class Theme:
    """封装Theme相关常量或数据结构。"""
    
    bg = RGBColor(247, 250, 255)
    navy = RGBColor(13, 35, 68)
    blue = RGBColor(34, 104, 201)
    cyan = RGBColor(16, 143, 169)
    indigo = RGBColor(70, 89, 173)
    text = RGBColor(31, 41, 55)
    subtext = RGBColor(95, 109, 129)
    white = RGBColor(255, 255, 255)
    border = RGBColor(210, 222, 238)
    light_blue = RGBColor(235, 243, 255)
    light_cyan = RGBColor(233, 250, 253)
    light_purple = RGBColor(241, 238, 255)
    light_orange = RGBColor(255, 244, 232)
    red = RGBColor(211, 58, 58)
    green = RGBColor(19, 159, 107)
    yellow = RGBColor(225, 149, 36)


def add_bg(slide, prs: Presentation) -> None:
    """向当前页补充背景相关元素，并统一样式与布局。"""
    
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = Theme.bg
    rect.line.fill.background()


def add_header(slide, title: str, subtitle: str, page: int) -> None:
    """向当前页补充页眉相关元素，并统一样式与布局。"""
    
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.86)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = Theme.navy
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.45), Inches(0.16), Inches(8.2), Inches(0.5))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = Theme.white

    sbox = slide.shapes.add_textbox(Inches(8.6), Inches(0.22), Inches(3.8), Inches(0.35))
    sf = sbox.text_frame
    sf.clear()
    sp = sf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(12)
    sp.font.color.rgb = RGBColor(212, 225, 245)
    sp.alignment = PP_ALIGN.RIGHT

    pbox = slide.shapes.add_textbox(Inches(12.5), Inches(0.19), Inches(0.5), Inches(0.35))
    pf = pbox.text_frame
    pf.clear()
    pp = pf.paragraphs[0]
    pp.text = str(page)
    pp.font.size = Pt(12)
    pp.font.bold = True
    pp.font.color.rgb = RGBColor(188, 207, 238)
    pp.alignment = PP_ALIGN.RIGHT


def add_card(slide, x: float, y: float, w: float, h: float, fill: RGBColor = Theme.white):
    """向当前页补充卡片相关元素，并统一样式与布局。"""
    
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = Theme.border
    return shp


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: int = 16,
    bold: bool = False,
    color: RGBColor = Theme.text,
    align: str = "left",
) -> None:
    """向当前页补充文本相关元素，并统一样式与布局。"""
    
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: Iterable[str],
    *,
    size: int = 14,
    color: RGBColor = Theme.text,
) -> None:
    """向当前页补充bullets相关元素，并统一样式与布局。"""
    
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(size)
        p.font.color.rgb = color


def add_arrow(slide, x: float, y: float, w: float, h: float, color: RGBColor = Theme.blue) -> None:
    """向当前页补充箭头相关元素，并统一样式与布局。"""
    
    arr = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()


def slide_cover(prs: Presentation) -> None:
    """构建封面对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.6)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = Theme.navy
    band.line.fill.background()
    add_text(slide, 0.8, 0.36, 9.5, 0.7, "生产问题根因分析系统", size=42, bold=True, color=Theme.white)
    add_text(
        slide,
        0.8,
        1.05,
        8.5,
        0.35,
        "LangGraph Multi-Agent RCA Platform",
        size=16,
        color=RGBColor(202, 222, 250),
    )

    add_card(slide, 0.8, 2.2, 5.9, 3.9, Theme.white)
    add_text(slide, 1.1, 2.55, 5.3, 0.55, "汇报目标", size=22, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        1.1,
        3.05,
        5.3,
        2.65,
        [
            "说明当前系统已实现能力与核心价值",
            "展示多 Agent 架构、关键链路与可观测机制",
            "说明前端交互、工具调用、报告输出闭环",
            "给出下一阶段优化路线图",
        ],
        size=15,
    )

    add_card(slide, 7.05, 2.2, 5.5, 3.9, Theme.light_blue)
    add_text(slide, 7.35, 2.55, 4.8, 0.55, "技术基线", size=22, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        7.35,
        3.05,
        4.9,
        2.65,
        [
            "Backend: FastAPI + LangGraph + LangChain OpenAI",
            "Frontend: React 18 + TypeScript + Ant Design",
            "Model: kimi-k2.5 (OpenAI-compatible API)",
            "Storage: local file / memory (no external DB)",
        ],
        size=15,
    )

    add_text(
        slide,
        0.8,
        6.65,
        12.0,
        0.35,
        f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}（Asia/Shanghai）",
        size=12,
        color=Theme.subtext,
    )


def slide_agenda(prs: Presentation) -> None:
    """构建议程对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_header(slide, "目录", "Presentation Roadmap", 2)

    items = [
        "1. 业务目标与痛点",
        "2. 系统总体架构图（端到端）",
        "3. LangGraph 运行时架构图（后端内部）",
        "4. 多 Agent 协作网络图与角色分工",
        "5. 工具调用与审计链路",
        "6. 核心流程图（输入 -> 结论 -> 报告）",
        "7. 数据与状态机设计",
        "8. 前端交互与用户体验设计",
        "9. 质量保障与可观测体系",
        "10. 改进路线图",
    ]
    for i, item in enumerate(items):
        x = 0.95 if i < 5 else 6.9
        y = 1.35 + (i % 5) * 1.06
        c = add_card(slide, x, y, 5.45, 0.84, Theme.white)
        c.line.color.rgb = Theme.border
        add_text(slide, x + 0.22, y + 0.22, 5.0, 0.42, item, size=14, color=Theme.text)


def slide_business(prs: Presentation) -> None:
    """构建business对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_header(slide, "1) 业务目标与定位", "Why this platform", 3)

    add_card(slide, 0.65, 1.2, 4.1, 5.9, Theme.light_orange)
    add_text(slide, 0.95, 1.5, 3.5, 0.5, "生产排障现状痛点", size=20, bold=True, color=Theme.red)
    add_bullets(
        slide,
        0.95,
        2.0,
        3.55,
        4.85,
        [
            "故障信息分散在日志、代码、监控、文档多个系统",
            "跨团队沟通成本高，责任边界不清晰",
            "单模型结论可解释性不足，复盘困难",
            "缺少可审计、可回放的分析过程记录",
            "高峰期容易出现分析任务积压与 pending",
        ],
        size=14,
    )

    add_card(slide, 4.95, 1.2, 7.7, 2.85, Theme.white)
    add_text(slide, 5.25, 1.48, 6.9, 0.45, "平台目标", size=19, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        5.25,
        1.95,
        6.9,
        1.75,
        [
            "主 Agent 指挥专家 Agent 协同分析，形成可辩论、可收敛的结论",
            "统一接入责任田映射、日志检索、代码检索、监控/CMDB 数据",
            "前端实时展示：资产映射、辩论过程、辩论结果、报告与回放",
        ],
        size=14,
    )

    add_card(slide, 4.95, 4.2, 7.7, 2.9, Theme.light_blue)
    add_text(slide, 5.25, 4.48, 6.9, 0.45, "价值输出", size=19, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        5.25,
        4.95,
        6.9,
        1.9,
        [
            "MTTR 缩短：降低人工排查时间",
            "结论质量提升：Top-K 根因 + 证据链 + 置信度",
            "组织协同提效：责任田自动定位 + 标准化报告",
            "治理闭环：Benchmark 与审计指标可持续优化",
        ],
        size=14,
    )


def slide_architecture(prs: Presentation) -> None:
    """构建架构对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_header(slide, "2) 系统总体架构图", "End-to-end Architecture", 4)

    # Row 1
    add_card(slide, 0.5, 1.2, 2.0, 0.95, Theme.light_orange)
    add_text(slide, 0.66, 1.54, 1.65, 0.3, "用户/值班 SRE", size=13, bold=True, align="center")
    add_arrow(slide, 2.55, 1.52, 0.35, 0.25)
    add_card(slide, 2.95, 1.2, 2.45, 0.95, Theme.light_blue)
    add_text(slide, 3.2, 1.48, 1.95, 0.4, "Frontend\nReact + AntD", size=13, bold=True, align="center")
    add_arrow(slide, 5.45, 1.52, 0.35, 0.25)
    add_card(slide, 5.85, 1.2, 2.25, 0.95, Theme.light_blue)
    add_text(slide, 6.08, 1.5, 1.8, 0.35, "Backend API\nFastAPI", size=13, bold=True, align="center")
    add_arrow(slide, 8.15, 1.52, 0.35, 0.25)
    add_card(slide, 8.55, 1.2, 4.25, 0.95, Theme.light_purple)
    add_text(slide, 8.75, 1.5, 3.85, 0.32, "LangGraph Runtime Orchestrator", size=13, bold=True, align="center")

    # Row 2
    add_card(slide, 8.55, 2.45, 4.25, 2.05, Theme.white)
    add_text(slide, 8.8, 2.66, 3.8, 0.3, "Agent 协作层", size=13, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        8.8,
        2.95,
        3.85,
        1.45,
        [
            "主 Agent: ProblemAnalysisAgent",
            "专家 Agent: Log/Domain/Code/Metrics/Change/Runbook",
            "对抗 Agent: Critic/Rebuttal/Judge/Verification",
        ],
        size=12,
    )

    add_card(slide, 0.5, 2.45, 7.95, 2.05, Theme.white)
    add_text(slide, 0.78, 2.66, 7.3, 0.28, "交互与数据通道", size=13, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        0.78,
        2.95,
        7.3,
        1.45,
        [
            "WebSocket 实时事件流：phase_changed / agent_chat / tool_io / session_status",
            "REST API：incident/debate/assets/reports/settings/benchmark/governance",
            "报告与谱系回放：lineage + report versions",
        ],
        size=12,
    )

    # Row 3
    add_card(slide, 0.5, 4.75, 4.1, 2.0, Theme.light_cyan)
    add_text(slide, 0.78, 5.02, 3.55, 0.3, "工具与外部连接器", size=13, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        0.78,
        5.28,
        3.55,
        1.35,
        [
            "Git Repo Search（CodeAgent）",
            "Local Log Reader（LogAgent）",
            "Domain Excel/CSV（DomainAgent）",
            "Telemetry/CMDB Connector（可配置）",
        ],
        size=12,
    )

    add_card(slide, 4.78, 4.75, 3.9, 2.0, Theme.light_blue)
    add_text(slide, 5.03, 5.02, 3.4, 0.3, "存储层（本地优先）", size=13, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        5.03,
        5.28,
        3.35,
        1.35,
        [
            "Session / Rounds / Results",
            "Lineage / Tool Audit / Replay",
            "Report(JSON/Markdown/HTML)",
            "file 或 memory backend",
        ],
        size=12,
    )

    add_card(slide, 8.95, 4.75, 3.85, 2.0, Theme.light_purple)
    add_text(slide, 9.2, 5.02, 3.3, 0.3, "运维治理层", size=13, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        9.2,
        5.28,
        3.25,
        1.35,
        [
            "Benchmark Center",
            "Governance Center",
            "Tool Registry / Trial / Audit",
            "WorkLog + Replay + Compare",
        ],
        size=12,
    )


def slide_runtime_arch(prs: Presentation) -> None:
    """构建运行时arch对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_header(slide, "3) LangGraph 运行时架构图", "Backend Internal Design", 5)

    # Central orchestrator
    add_card(slide, 4.85, 2.55, 3.6, 1.15, Theme.light_purple)
    add_text(slide, 5.03, 2.86, 3.2, 0.44, "LangGraphRuntimeOrchestrator", size=15, bold=True, align="center")

    nodes: List[Tuple[str, float, float, RGBColor]] = [
        ("GraphBuilder\n(builder.py)", 1.0, 1.25, Theme.light_blue),
        ("RoutingStrategy\n(routing/*.py)", 1.0, 3.05, Theme.light_blue),
        ("AgentRunner\n(agent_runner.py)", 1.0, 4.85, Theme.light_blue),
        ("EventDispatcher\n(event_dispatcher.py)", 9.2, 1.25, Theme.light_cyan),
        ("WorkLogManager\n(work_log_manager.py)", 9.2, 3.05, Theme.light_cyan),
        ("Compaction + Prune\n(session_compaction.py)", 9.2, 4.85, Theme.light_cyan),
    ]
    for title, x, y, fill in nodes:
        add_card(slide, x, y, 3.1, 1.2, fill)
        add_text(slide, x + 0.15, y + 0.33, 2.8, 0.55, title, size=12, bold=True, align="center")

    # arrows
    add_arrow(slide, 4.2, 1.75, 0.5, 0.26)
    add_arrow(slide, 4.2, 3.55, 0.5, 0.26)
    add_arrow(slide, 4.2, 5.35, 0.5, 0.26)

    left_arrow_1 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.LEFT_ARROW, Inches(8.5), Inches(1.75), Inches(0.5), Inches(0.26)
    )
    left_arrow_1.fill.solid()
    left_arrow_1.fill.fore_color.rgb = Theme.blue
    left_arrow_1.line.fill.background()
    left_arrow_2 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.LEFT_ARROW, Inches(8.5), Inches(3.55), Inches(0.5), Inches(0.26)
    )
    left_arrow_2.fill.solid()
    left_arrow_2.fill.fore_color.rgb = Theme.blue
    left_arrow_2.line.fill.background()
    left_arrow_3 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.LEFT_ARROW, Inches(8.5), Inches(5.35), Inches(0.5), Inches(0.26)
    )
    left_arrow_3.fill.solid()
    left_arrow_3.fill.fore_color.rgb = Theme.blue
    left_arrow_3.line.fill.background()

    add_text(
        slide,
        0.95,
        6.45,
        11.8,
        0.33,
        "设计目标：编排器只负责协调，执行/路由/事件/上下文管理模块化，避免“上帝类”。",
        size=12,
        color=Theme.subtext,
    )


def slide_agent_network(prs: Presentation) -> None:
    """构建Agent链路图对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_header(slide, "4) 多 Agent 协作网络图", "Command-driven Collaboration", 6)

    center = add_card(slide, 5.25, 2.55, 2.8, 1.2, Theme.light_purple)
    center.line.color.rgb = Theme.indigo
    add_text(slide, 5.45, 2.9, 2.4, 0.4, "ProblemAnalysisAgent", size=13, bold=True, align="center")

    ring = [
        ("LogAgent", 2.0, 1.2, Theme.light_blue),
        ("DomainAgent", 5.1, 1.0, Theme.light_cyan),
        ("CodeAgent", 8.3, 1.2, Theme.light_blue),
        ("MetricsAgent", 10.0, 2.9, Theme.light_cyan),
        ("ChangeAgent", 8.4, 4.8, Theme.light_blue),
        ("RunbookAgent", 5.1, 5.0, Theme.light_cyan),
        ("CriticAgent", 2.0, 4.8, Theme.light_orange),
        ("RebuttalAgent", 0.6, 2.9, Theme.light_orange),
        ("JudgeAgent", 10.0, 4.5, Theme.light_purple),
    ]

    for name, x, y, fill in ring:
        add_card(slide, x, y, 2.05, 0.95, fill)
        add_text(slide, x + 0.13, y + 0.3, 1.78, 0.35, name, size=12, bold=True, align="center")

    # simple directional arrows to center
    arrows = [
        (4.2, 1.95, 0.95, 0.2),
        (6.55, 1.95, 0.95, 0.2),
        (8.25, 2.85, 0.95, 0.2),
        (8.35, 4.25, 0.95, 0.2),
        (6.55, 4.95, 0.95, 0.2),
        (4.2, 4.95, 0.95, 0.2),
        (2.9, 4.25, 0.95, 0.2),
        (2.9, 2.85, 0.95, 0.2),
    ]
    for x, y, w, h in arrows:
        add_arrow(slide, x, y, w, h, Theme.cyan)

    add_text(
        slide,
        0.65,
        6.35,
        12.0,
        0.43,
        "主 Agent 先下发命令；子 Agent 依据命令调用工具并回传结构化证据；Critic/Rebuttal/Judge 完成争议收敛。",
        size=12,
        color=Theme.subtext,
    )


def slide_tool_audit(prs: Presentation) -> None:
    """构建工具audit对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_header(slide, "5) 工具调用与审计链路", "Tooling & Auditability", 7)

    # flow
    stages = [
        ("主 Agent 下发命令", 0.7, Theme.light_purple),
        ("Command Gate 判定\n(allow / deny)", 3.0, Theme.light_blue),
        ("Agent 调用工具\nGit / Log / Excel / HTTP", 5.65, Theme.light_cyan),
        ("输出截断与摘要\n(large output handling)", 8.45, Theme.light_blue),
        ("写入审计与谱系\n(lineage + tool_audit)", 10.95, Theme.light_purple),
    ]
    for i, (text, x, fill) in enumerate(stages):
        w = 2.1 if i != 2 else 2.5
        add_card(slide, x, 2.15, w, 1.35, fill)
        add_text(slide, x + 0.1, 2.58, w - 0.2, 0.7, text, size=11, bold=True, align="center")
        if i < len(stages) - 1:
            add_arrow(slide, x + w + 0.08, 2.68, 0.28, 0.22, Theme.blue)

    add_card(slide, 0.7, 4.05, 6.2, 2.6, Theme.white)
    add_text(slide, 0.95, 4.3, 5.7, 0.4, "审计字段（前后端可见）", size=16, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        0.95,
        4.75,
        5.7,
        1.75,
        [
            "event_type / timestamp / agent_name / phase",
            "tool_name / action / status / error_reason",
            "command_preview / decision_source / permission_decision",
            "io_trace：文件读取、Git 命令、HTTP 请求",
            "full_output_ref：完整输出引用 ID",
        ],
        size=12,
    )

    add_card(slide, 7.1, 4.05, 5.5, 2.6, Theme.light_blue)
    add_text(slide, 7.35, 4.3, 5.0, 0.4, "设计原则", size=16, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        7.35,
        4.75,
        4.95,
        1.75,
        [
            "没有命令，不触发工具调用",
            "没有工具配置，不展示工具调用消息",
            "工具失败必须显式降级并给出 fallback 结论",
            "所有调用可在战情页和复盘台追溯",
        ],
        size=12,
    )


def slide_flow(prs: Presentation) -> None:
    """构建flow对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_header(slide, "6) 核心流程图", "Incident -> Debate -> Report", 8)

    # Process line
    steps = [
        ("输入故障信息", "日志、堆栈、监控现象"),
        ("创建会话", "session + mode + strategy"),
        ("资产映射", "接口 -> 领域 -> 聚合根 -> owner"),
        ("主 Agent 分发任务", "command_issued"),
        ("专家 Agent 分析", "tool call + chat + evidence"),
        ("对抗收敛", "critic / rebuttal / judge"),
        ("输出结果", "Top-K 根因 + 证据链 + 置信度"),
        ("生成报告", "JSON / Markdown / HTML"),
    ]

    x0 = 0.75
    y = 2.1
    for i, (title, desc) in enumerate(steps):
        width = 1.45
        if i in {2, 3, 4, 6}:
            width = 1.6
        add_card(slide, x0, y, width, 1.8, Theme.white)
        add_text(slide, x0 + 0.11, y + 0.18, width - 0.2, 0.35, f"{i+1}. {title}", size=11, bold=True, color=Theme.indigo)
        add_text(slide, x0 + 0.11, y + 0.65, width - 0.2, 0.9, desc, size=10, color=Theme.text)
        if i < len(steps) - 1:
            add_arrow(slide, x0 + width + 0.05, y + 0.72, 0.22, 0.18, Theme.blue)
        x0 += width + 0.33

    add_card(slide, 0.75, 4.45, 12.0, 2.0, Theme.light_cyan)
    add_text(slide, 1.03, 4.73, 11.3, 0.35, "异常与超时处理策略", size=15, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        1.03,
        5.1,
        11.3,
        1.2,
        [
            "分阶段超时：analysis/review/judge/report 使用不同 timeout",
            "局部重试与降级：单 Agent 失败不应阻塞整场会话",
            "DEBATE_REQUIRE_EFFECTIVE_LLM_CONCLUSION = true，避免空结论直出报告",
            "状态迁移校验：pending -> running -> analyzing/debating -> judging -> completed/failed",
        ],
        size=12,
    )


def slide_state_data(prs: Presentation) -> None:
    """构建状态data对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_header(slide, "7) 状态机与数据模型", "State & Data Design", 9)

    add_card(slide, 0.6, 1.2, 6.2, 5.95, Theme.white)
    add_text(slide, 0.88, 1.48, 5.7, 0.38, "会话状态机（DebateStatus）", size=16, bold=True, color=Theme.indigo)
    states = [
        ("PENDING", 1.0, 2.0, Theme.light_blue),
        ("RUNNING", 2.4, 2.0, Theme.light_blue),
        ("ANALYZING", 3.8, 2.0, Theme.light_cyan),
        ("DEBATING", 5.2, 2.0, Theme.light_cyan),
        ("JUDGING", 1.7, 3.3, Theme.light_purple),
        ("COMPLETED", 3.6, 3.3, Theme.light_purple),
        ("FAILED", 5.4, 3.3, Theme.light_orange),
        ("RETRYING", 2.7, 4.5, Theme.light_orange),
        ("CANCELLED", 4.7, 4.5, Theme.light_orange),
    ]
    for name, x, y, fill in states:
        add_card(slide, x, y, 1.22, 0.62, fill)
        add_text(slide, x + 0.05, y + 0.2, 1.1, 0.25, name, size=10, bold=True, align="center")

    # state arrows
    add_arrow(slide, 2.23, 2.25, 0.15, 0.12, Theme.blue)
    add_arrow(slide, 3.63, 2.25, 0.15, 0.12, Theme.blue)
    add_arrow(slide, 5.03, 2.25, 0.15, 0.12, Theme.blue)
    add_arrow(slide, 2.75, 3.02, 0.5, 0.12, Theme.blue)
    add_arrow(slide, 4.57, 3.02, 0.5, 0.12, Theme.blue)
    add_arrow(slide, 3.45, 4.2, 0.45, 0.12, Theme.yellow)

    add_text(
        slide,
        1.0,
        5.45,
        5.5,
        1.35,
        "说明：状态迁移受服务端白名单校验，避免非法 phase（如 coordination）导致系统卡死。",
        size=11,
        color=Theme.subtext,
    )

    add_card(slide, 7.0, 1.2, 5.7, 5.95, Theme.light_blue)
    add_text(slide, 7.28, 1.48, 5.2, 0.38, "核心数据对象", size=16, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        7.28,
        1.9,
        5.1,
        2.2,
        [
            "Incident：原始故障描述、日志、堆栈、监控现象",
            "DebateSession：执行上下文、状态、模式、策略",
            "DebateRound：每轮 Agent 输出、证据、工具调用",
            "DebateResult：根因候选、证据链、修复建议、验证计划",
            "LineageRecord：事件时间线、工具审计、回放关键节点",
        ],
        size=12,
    )
    add_text(slide, 7.28, 4.3, 5.15, 0.34, "存储后端", size=14, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        7.28,
        4.65,
        5.1,
        1.8,
        [
            "LOCAL_STORE_BACKEND = file | memory",
            "LOCAL_STORE_DIR 管理会话与报告持久化",
            "不引入外部数据库，便于本地开发与演示",
        ],
        size=12,
    )


def slide_frontend(prs: Presentation) -> None:
    """构建前端对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_header(slide, "8) 前端交互与用户体验", "UX Architecture", 10)

    add_card(slide, 0.6, 1.2, 12.1, 1.0, Theme.white)
    add_text(
        slide,
        0.95,
        1.52,
        11.5,
        0.35,
        "设计原则：实时可感知、过程可追溯、结论可解释、关键操作可直达。",
        size=15,
        color=Theme.indigo,
        bold=True,
    )

    pages = [
        ("首页 /", "创建并启动分析、近期任务、Agent 介绍、统计指标"),
        ("分析页 /incident", "资产映射 / 辩论过程 / 辩论结果 三标签"),
        ("战情页 /war-room", "同屏时间线 + 工具调用 + 证据链 + 关键结论"),
        ("调查复盘台 /workbench", "关键决策路径回放 + 报告版本差异"),
        ("工具中心 /tools", "工具列表、详情、参数试跑、审计查询"),
        ("治理中心 /governance", "成功率、超时率、成本和风险动作统计"),
    ]
    for i, (name, desc) in enumerate(pages):
        x = 0.6 + (i % 2) * 6.15
        y = 2.45 + (i // 2) * 1.55
        add_card(slide, x, y, 5.95, 1.35, Theme.light_blue if i % 2 == 0 else Theme.light_cyan)
        add_text(slide, x + 0.2, y + 0.2, 5.55, 0.3, name, size=13, bold=True, color=Theme.indigo)
        add_text(slide, x + 0.2, y + 0.55, 5.5, 0.6, desc, size=12)

    add_card(slide, 0.6, 5.75, 12.1, 1.35, Theme.light_purple)
    add_text(slide, 0.95, 5.98, 11.4, 0.3, "用户体验改进点（已落地）", size=14, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        0.95,
        6.32,
        11.3,
        0.6,
        [
            "流式聊天支持“缩略 + 展开全文”，避免大段 JSON 干扰阅读",
            "所有时间统一北京时间展示，历史/战情/首页口径一致",
            "分析中任务可进入详情页持续追踪，避免跳转空白页",
        ],
        size=12,
    )


def slide_interfaces(prs: Presentation) -> None:
    """构建interfaces对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_header(slide, "9) API 与扩展能力", "Integration & Extensibility", 11)

    add_card(slide, 0.65, 1.2, 6.1, 5.9, Theme.white)
    add_text(slide, 0.95, 1.48, 5.55, 0.4, "核心 API", size=16, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        0.95,
        1.9,
        5.5,
        4.95,
        [
            "POST /api/v1/incidents/",
            "POST /api/v1/debates/?incident_id=...",
            "POST /api/v1/debates/{session_id}/execute",
            "GET  /api/v1/debates/{session_id}/result",
            "POST /api/v1/reports/{incident_id}/regenerate",
            "GET  /api/v1/settings/tooling",
            "PUT  /api/v1/settings/tooling",
            "WS   /ws/debates/{session_id}?auto_start=true",
        ],
        size=12,
    )

    add_card(slide, 6.95, 1.2, 5.75, 2.7, Theme.light_blue)
    add_text(slide, 7.2, 1.45, 5.25, 0.36, "模型配置（当前）", size=15, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        7.2,
        1.85,
        5.2,
        1.8,
        [
            "LLM_BASE_URL=https://ark.cn-beijing.volces.com/api/coding",
            "LLM_API_KEY=********（环境变量注入）",
            "LLM_MODEL=kimi-k2.5",
            "LLM_MAX_CONCURRENCY / timeout 按阶段控制",
        ],
        size=11,
    )

    add_card(slide, 6.95, 4.15, 5.75, 2.95, Theme.light_cyan)
    add_text(slide, 7.2, 4.4, 5.2, 0.36, "扩展入口", size=15, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        7.2,
        4.8,
        5.2,
        2.0,
        [
            "TelemetryConnector：接入日志/APM/指标平台（可开关）",
            "CMDBConnector：接入服务拓扑/责任人系统（可开关）",
            "Tool Registry：create/update/delete/start/offline/list/run",
            "Skill 化 Prompt：按场景注入策略与输出格式约束",
        ],
        size=11,
    )


def slide_quality(prs: Presentation) -> None:
    """构建质量对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_header(slide, "10) 质量保障与可观测", "Reliability & Benchmark", 12)

    add_card(slide, 0.65, 1.2, 4.0, 5.9, Theme.light_blue)
    add_text(slide, 0.95, 1.48, 3.45, 0.36, "Benchmark 指标", size=16, bold=True, color=Theme.indigo)
    # pseudo bars
    metrics = [
        ("Top1 命中率", 0.72, Theme.green),
        ("Top3 命中率", 0.88, Theme.green),
        ("超时率", 0.21, Theme.red),
        ("空结论率", 0.12, Theme.yellow),
    ]
    y = 2.1
    for name, ratio, color in metrics:
        add_text(slide, 0.95, y, 1.8, 0.25, name, size=11)
        bar_bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.2), Inches(y + 0.03), Inches(2.2), Inches(0.18)
        )
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = RGBColor(223, 231, 244)
        bar_bg.line.fill.background()
        bar_fg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(2.2),
            Inches(y + 0.03),
            Inches(2.2 * ratio),
            Inches(0.18),
        )
        bar_fg.fill.solid()
        bar_fg.fill.fore_color.rgb = color
        bar_fg.line.fill.background()
        add_text(slide, 4.45, y - 0.01, 0.35, 0.25, f"{int(ratio*100)}%", size=10, color=Theme.subtext)
        y += 0.72

    add_text(
        slide, 0.95, 5.5, 3.4, 1.2, "注：以上为演示口径，真实值由 Benchmark Center 实时计算。", size=10, color=Theme.subtext
    )

    add_card(slide, 4.9, 1.2, 7.8, 2.75, Theme.white)
    add_text(slide, 5.18, 1.48, 7.2, 0.36, "可观测与防卡死设计", size=16, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        5.18,
        1.9,
        7.2,
        1.8,
        [
            "全链路时间戳与事件序号：便于定位“长时间无响应”区间",
            "Agent 超时/失败分类：timeout, rate_limit, no_conclusion, runtime_error",
            "局部失败可继续：保持会话推进，输出可执行 fallback 结论",
            "日志中文可读输出 + JSON 结构日志并存，兼顾排障与机器分析",
        ],
        size=12,
    )

    add_card(slide, 4.9, 4.15, 7.8, 2.95, Theme.light_cyan)
    add_text(slide, 5.18, 4.42, 7.2, 0.36, "治理视角（Governance）", size=16, bold=True, color=Theme.indigo)
    add_bullets(
        slide,
        5.18,
        4.85,
        7.2,
        2.0,
        [
            "团队维度：成功率、超时率、工具失败率、风险动作占比",
            "会话维度：关键决策路径、工具调用明细、证据引用轨迹",
            "版本维度：同 incident 多次分析结果对比与差异解释",
            "发布维度：Benchmark Gate 可接入 CI 阻断回归",
        ],
        size=12,
    )


def slide_roadmap(prs: Presentation) -> None:
    """构建路线图对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_header(slide, "11) 改进路线图", "P0~P3", 13)

    phases = [
        ("P0 可用性", "自动调查、并行拉证据、超时切换、pending 清零", Theme.light_orange),
        ("P1 准确率", "跨源证据约束、因果推理层、Top-K 置信区间", Theme.light_blue),
        ("P2 可控修复", "修复状态机、No-Regression Gate、回滚接口", Theme.light_cyan),
        ("P3 平台治理", "反馈学习、A/B 评测、多租户与外部协同", Theme.light_purple),
    ]
    for i, (name, desc, fill) in enumerate(phases):
        y = 1.45 + i * 1.35
        add_card(slide, 0.8, y, 11.9, 1.05, fill)
        add_text(slide, 1.08, y + 0.24, 2.15, 0.36, name, size=16, bold=True, color=Theme.indigo)
        add_text(slide, 3.35, y + 0.24, 9.0, 0.36, desc, size=14, color=Theme.text)

    add_card(slide, 0.8, 6.0, 11.9, 1.05, Theme.white)
    add_text(
        slide,
        1.08,
        6.36,
        11.2,
        0.32,
        "建议实施方式：按阶段验收 + Benchmark 回归门禁 + 前后端联调验收（非仅 API 测试）。",
        size=13,
        color=Theme.subtext,
    )


def slide_end(prs: Presentation) -> None:
    """构建结束对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_header(slide, "12) 总结", "Thank You", 14)

    add_card(slide, 1.0, 1.8, 11.3, 4.9, Theme.white)
    add_text(
        slide,
        1.35,
        2.35,
        10.6,
        0.55,
        "当前系统已从“单点分析”升级为“多 Agent 协作 + 工具审计 + 可回放治理”的生产故障分析平台。",
        size=23,
        bold=True,
        color=Theme.indigo,
    )
    add_bullets(
        slide,
        1.35,
        3.15,
        10.2,
        2.9,
        [
            "架构完整：FastAPI + LangGraph + React，支持实时事件流和多会话模式",
            "过程透明：命令驱动工具调用、可追溯审计、可复盘决策链",
            "结果可用：Top-K 根因、证据链、修复建议、验证计划与结构化报告",
            "可持续优化：Benchmark + Governance 双闭环，支撑持续提升",
        ],
        size=16,
    )
    add_text(
        slide,
        1.35,
        6.3,
        10.2,
        0.4,
        "下一步：继续提升真实数据接入、模型韧性与自治修复能力。",
        size=14,
        color=Theme.subtext,
    )


def build(out: Path) -> Path:
    """执行build相关逻辑。"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_agenda(prs)
    slide_business(prs)
    slide_architecture(prs)
    slide_runtime_arch(prs)
    slide_agent_network(prs)
    slide_tool_audit(prs)
    slide_flow(prs)
    slide_state_data(prs)
    slide_frontend(prs)
    slide_interfaces(prs)
    slide_quality(prs)
    slide_roadmap(prs)
    slide_end(prs)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


if __name__ == "__main__":
    result = build(OUTFILE)
    print(f"Generated: {result}")
