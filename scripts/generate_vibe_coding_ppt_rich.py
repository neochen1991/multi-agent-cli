#!/usr/bin/env python3
"""generateVibeCodingPPTrich脚本。"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output" / "vibe-coding-ppt-rich"
OUT_PPTX = OUT_DIR / "vibe-coding-industry-playbook-2026-03-06.pptx"
OUT_MD = OUT_DIR / "SOURCES.md"


class Theme:
    """封装Theme相关常量或数据结构。"""
    
    white = RGBColor(255, 255, 255)
    bg = RGBColor(243, 247, 252)
    navy = RGBColor(15, 38, 73)
    blue = RGBColor(40, 110, 225)
    cyan = RGBColor(22, 153, 173)
    green = RGBColor(40, 156, 107)
    orange = RGBColor(226, 145, 34)
    red = RGBColor(208, 74, 74)
    text = RGBColor(35, 45, 57)
    subtext = RGBColor(96, 109, 125)
    border = RGBColor(203, 216, 235)
    light_blue = RGBColor(232, 241, 255)
    light_cyan = RGBColor(232, 249, 253)
    light_green = RGBColor(230, 248, 241)
    light_orange = RGBColor(255, 243, 229)
    light_red = RGBColor(253, 236, 236)


def ensure_dirs() -> None:
    """确保dirs相关前置条件已经满足。"""
    
    OUT_DIR.mkdir(parents=True, exist_ok=True)


def prs_new() -> Presentation:
    """执行prs新增相关逻辑。"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_bg(slide, prs: Presentation) -> None:
    """向当前页补充背景相关元素，并统一样式与布局。"""
    
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = Theme.bg
    bg.line.fill.background()


def add_topbar(slide, title: str, subtitle: str, page: int) -> None:
    """向当前页补充topbar相关元素，并统一样式与布局。"""
    
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.86),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = Theme.navy
    bar.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.4), Inches(0.16), Inches(8.6), Inches(0.44))
    tp = t.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = Theme.white

    s = slide.shapes.add_textbox(Inches(8.8), Inches(0.21), Inches(3.6), Inches(0.34))
    sp = s.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(11)
    sp.font.color.rgb = RGBColor(207, 221, 245)
    sp.alignment = PP_ALIGN.RIGHT

    n = slide.shapes.add_textbox(Inches(12.5), Inches(0.2), Inches(0.45), Inches(0.3))
    np = n.text_frame.paragraphs[0]
    np.text = str(page)
    np.font.size = Pt(12)
    np.font.bold = True
    np.font.color.rgb = RGBColor(176, 198, 233)
    np.alignment = PP_ALIGN.RIGHT


def add_source(slide, text: str) -> None:
    """向当前页补充来源相关元素，并统一样式与布局。"""
    
    b = slide.shapes.add_textbox(Inches(0.45), Inches(7.12), Inches(12.4), Inches(0.22))
    p = b.text_frame.paragraphs[0]
    p.text = f"Source: {text}"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(118, 130, 145)


def add_card(slide, x: float, y: float, w: float, h: float, fill: RGBColor = Theme.white):
    """向当前页补充卡片相关元素，并统一样式与布局。"""
    
    c = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.color.rgb = Theme.border
    return c


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = Theme.text,
    align: str = "left",
) -> None:
    """向当前页补充文本相关元素，并统一样式与布局。"""
    
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT


def add_bullets(
    slide, x: float, y: float, w: float, h: float, lines: Iterable[str], *, size: int = 12, color: RGBColor = Theme.text
) -> None:
    """向当前页补充bullets相关元素，并统一样式与布局。"""
    
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(size)
        p.font.color.rgb = color


def add_arrow(slide, x: float, y: float, w: float, h: float, color: RGBColor = Theme.blue) -> None:
    """向当前页补充箭头相关元素，并统一样式与布局。"""
    
    a = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()


def add_chart(
    slide,
    *,
    chart_type: XL_CHART_TYPE,
    x: float,
    y: float,
    w: float,
    h: float,
    categories: list[str],
    series: list[tuple[str, list[float]]],
) -> None:
    """向当前页补充图表相关元素，并统一样式与布局。"""
    
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    chart = slide.shapes.add_chart(chart_type, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    if chart.value_axis is not None:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.font.size = Pt(10)
    if chart.category_axis is not None:
        chart.category_axis.tick_labels.font.size = Pt(10)


def slide_1_cover(prs: Presentation) -> None:
    """构建1封面对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)

    top = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.2)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = Theme.navy
    top.line.fill.background()

    add_text(s, 0.6, 0.58, 10.5, 0.8, "Vibe Coding 生产级实践蓝图", size=46, bold=True, color=Theme.white)
    add_text(
        s,
        0.6,
        1.46,
        11.2,
        0.42,
        "基于业界最佳实践与本项目落地经验（Skill / MCP / Multi-Agent / Eval）",
        size=17,
        color=RGBColor(206, 221, 247),
    )

    add_card(s, 0.75, 2.7, 6.0, 3.95, Theme.white)
    add_text(s, 1.0, 3.05, 5.4, 0.45, "本次重构目标", size=22, bold=True, color=Theme.blue)
    add_bullets(
        s,
        1.0,
        3.52,
        5.4,
        2.8,
        [
            "补齐“业界方法论”而不是只讲项目功能",
            "把经验拆成可执行的工程策略和模板",
            "给出可复用的团队落地路径（2周/90天）",
            "形成可持续优化的评测与治理体系",
        ],
        size=14,
    )

    add_card(s, 7.05, 2.7, 5.55, 3.95, Theme.light_blue)
    add_text(s, 7.3, 3.05, 4.9, 0.45, "内容结构", size=22, bold=True, color=Theme.blue)
    add_bullets(
        s,
        7.3,
        3.52,
        4.9,
        2.8,
        [
            "Part A: 业界实践提炼（12页）",
            "Part B: 项目映射与改造建议（7页）",
            "Part C: 路线图、指标与风险治理（3页）",
        ],
        size=14,
    )

    add_text(
        s,
        0.78,
        6.88,
        11.0,
        0.26,
        f"Generated at {datetime.now().strftime('%Y-%m-%d %H:%M:%S')} Asia/Shanghai",
        size=10,
        color=Theme.subtext,
    )
    add_source(s, "OpenAI Harness Engineering, Anthropic/Claude Docs, GitHub Docs, LangGraph Docs, MCP Spec, Thoughtworks")


def slide_2_agenda(prs: Presentation) -> None:
    """构建2议程对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "议程", "Industry-first content structure", 2)

    sections = [
        "1. Vibe Coding 的边界与价值",
        "2. Harness / Context / Eval 三大基础",
        "3. Skill、MCP、工具治理与安全",
        "4. 多 Agent 架构与可靠性模式",
        "5. 业界工具链对比（OpenAI/Anthropic/GitHub/Sourcegraph）",
        "6. 失败模式与反模式（含 Thoughtworks 案例）",
        "7. 本项目改造映射与落地路线图",
        "8. KPI、门禁、组织治理",
    ]
    y = 1.25
    for i, sec in enumerate(sections, start=1):
        fill = Theme.light_blue if i % 2 else Theme.white
        add_card(s, 0.8, y, 11.9, 0.62, fill)
        add_text(s, 1.05, y + 0.18, 11.2, 0.24, sec, size=14, bold=i in (1, 7))
        y += 0.72
    add_source(s, "Deck structure designed from user requirement and industry references")


def slide_3_define(prs: Presentation) -> None:
    """构建3define对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "什么是 Vibe Coding（生产视角）", "Definition and boundary", 3)

    add_card(s, 0.8, 1.25, 5.7, 5.7, Theme.white)
    add_text(s, 1.05, 1.55, 5.2, 0.35, "定义（工程化版）", size=18, bold=True, color=Theme.blue)
    add_bullets(
        s,
        1.05,
        2.0,
        5.2,
        2.2,
        [
            "以自然语言驱动开发，但由系统约束保障质量",
            "不是“随便生成代码”，而是“有门禁的高吞吐协作”",
            "核心角色变化：人从写代码转向设计约束与反馈循环",
        ],
        size=13,
    )

    add_card(s, 1.05, 4.5, 5.2, 2.2, Theme.light_orange)
    add_text(s, 1.25, 4.8, 4.8, 0.3, "一句话", size=14, bold=True, color=Theme.orange)
    add_text(s, 1.25, 5.15, 4.8, 1.2, "Vibe Coding = 生成速度 × 约束强度 × 反馈闭环", size=18, bold=True, color=Theme.text)

    add_card(s, 6.8, 1.25, 5.7, 5.7, Theme.white)
    add_text(s, 7.05, 1.55, 5.2, 0.35, "常见误区", size=18, bold=True, color=Theme.red)
    add_bullets(
        s,
        7.05,
        2.0,
        5.2,
        2.5,
        [
            "误区1：Prompt 越长越好",
            "误区2：先快做出来，后面再补测试",
            "误区3：Agent 会自己变好，不需要评测体系",
            "误区4：工具接得越多越强，不必治理权限",
        ],
        size=13,
    )
    add_card(s, 7.05, 4.9, 5.2, 1.8, Theme.light_red)
    add_text(s, 7.25, 5.18, 4.8, 1.2, "正确做法：控制系统先行，模型能力后置。", size=15, bold=True)
    add_source(s, "OpenAI Harness Engineering (2026-02-11), Thoughtworks (2025-04-30)")


def slide_4_industry_signals(prs: Presentation) -> None:
    """构建4industrysignals对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "业界信号：为什么这不是短期热词", "Industry evidence", 4)

    add_card(s, 0.75, 1.2, 6.3, 5.9, Theme.white)
    add_text(s, 1.0, 1.5, 5.8, 0.34, "公开案例与量化信号", size=17, bold=True, color=Theme.blue)
    add_chart(
        s,
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        x=1.0,
        y=1.95,
        w=5.7,
        h=3.2,
        categories=["OpenAI", "Thoughtworks", "GitHub", "Sourcegraph"],
        series=[
            ("工程化成熟指数(示意)", [92, 78, 85, 80]),
            ("自动化协作深度(示意)", [90, 68, 82, 86]),
        ],
    )
    add_bullets(
        s,
        1.0,
        5.35,
        5.8,
        1.35,
        [
            "OpenAI: 报告 0 手写代码实验，约 1/10 开发时间",
            "Thoughtworks: 3 组实验显示“有纪律提示”明显优于纯 vibe",
        ],
        size=11,
    )

    add_card(s, 7.3, 1.2, 5.3, 5.9, Theme.white)
    add_text(s, 7.55, 1.5, 4.8, 0.34, "关键结论", size=17, bold=True, color=Theme.blue)
    add_bullets(
        s,
        7.55,
        1.95,
        4.8,
        4.6,
        [
            "高速度不是问题，稳定性才是门槛。",
            "“人类负责意图与约束，Agent 负责执行”已成为主流模式。",
            "上下文管理、评测门禁、安全治理是共识，不再是可选项。",
            "工具化能力（MCP/插件）必须伴随权限和审计。",
            "多 Agent 趋势明确，但“可收敛”比“可并行”更关键。",
        ],
        size=13,
    )
    add_source(
        s,
        "OpenAI harness-engineering; Thoughtworks can-vibe-coding-produce-production-grade-software; GitHub Copilot docs; Sourcegraph agentic context docs",
    )


def slide_5_harness(prs: Presentation) -> None:
    """构建5harness对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "Pillar 1: Harness Engineering", "Humans steer, agents execute", 5)

    add_text(s, 0.8, 1.12, 11.8, 0.4, "OpenAI 的核心表达：Humans steer. Agents execute.", size=20, bold=True, color=Theme.blue)

    pillars = [
        ("1", "任务拆解", "把大任务拆成可验证子任务", Theme.light_blue),
        ("2", "环境可读", "让 UI/日志/指标对 Agent 可见", Theme.light_cyan),
        ("3", "规则可执行", "架构规则通过 lint/check 强制化", Theme.light_green),
        ("4", "反馈闭环", "失败信号回写成规则/工具", Theme.light_orange),
        ("5", "持续清理", "定期垃圾回收，遏制“AI slop”", Theme.light_red),
    ]
    x = 0.95
    for idx, title, desc, fill in pillars:
        add_card(s, x, 2.1, 2.35, 3.95, fill)
        add_text(s, x + 0.15, 2.35, 2.05, 0.45, f"{idx}. {title}", size=14, bold=True, color=Theme.navy, align="center")
        add_text(s, x + 0.2, 2.92, 1.95, 2.85, desc, size=12, align="center")
        if idx != "5":
            add_arrow(s, x + 2.35, 3.75, 0.24, 0.2)
        x += 2.55
    add_source(s, "OpenAI Harness Engineering (2026-02-11)")


def slide_6_context(prs: Presentation) -> None:
    """构建6上下文对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "Pillar 2: Context Engineering", "Map, not 1000-page manuals", 6)

    add_card(s, 0.75, 1.2, 6.0, 5.9, Theme.white)
    add_text(s, 1.0, 1.52, 5.5, 0.35, "上下文层级设计", size=17, bold=True, color=Theme.blue)
    layers = [
        ("Repo TOC", "AGENTS.md 做目录，不做百科"),
        ("Scoped Rules", "按目录/任务挂载规则"),
        ("Runtime Context", "只注入当前步骤相关内容"),
        ("Memory", "长期记忆与项目记忆分层"),
    ]
    y = 2.0
    for name, desc in layers:
        add_card(s, 1.05, y, 5.35, 0.95, Theme.light_blue if y < 3.5 else Theme.light_cyan)
        add_text(s, 1.3, y + 0.2, 1.55, 0.3, name, size=13, bold=True, color=Theme.navy)
        add_text(s, 2.95, y + 0.2, 3.2, 0.45, desc, size=12)
        y += 1.12

    add_card(s, 7.0, 1.2, 5.55, 5.9, Theme.white)
    add_text(s, 7.25, 1.52, 5.1, 0.35, "业界实践对照", size=17, bold=True, color=Theme.blue)
    add_bullets(
        s,
        7.25,
        1.98,
        5.1,
        4.8,
        [
            "GitHub Copilot: 提供有效上下文，复杂任务拆小，必要时开启新会话。",
            "Cursor Rules: Project/User/Memory 分层规则。",
            "Claude Code: CLAUDE.md + /memory + /init 形成长期项目记忆。",
            "Sourcegraph: agentic context fetching 自动补全上下文并反思迭代。",
        ],
        size=13,
    )
    add_source(s, "OpenAI Harness blog; GitHub Copilot best practices; Cursor rules; Claude docs; Sourcegraph agentic context")


def slide_7_skill_memory(prs: Presentation) -> None:
    """构建7Skillmemory对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "Pillar 3: Skill + Memory 体系", "From prompts to reusable capabilities", 7)

    add_card(s, 0.8, 1.2, 12.0, 5.9, Theme.white)
    add_text(s, 1.05, 1.55, 11.4, 0.35, "将高频能力固化为 Skill，而不是重复写 Prompt", size=18, bold=True, color=Theme.blue)

    # left: skill lifecycle
    add_card(s, 1.05, 2.1, 5.7, 4.7, Theme.light_blue)
    add_text(s, 1.3, 2.38, 5.2, 0.32, "Skill 生命周期", size=15, bold=True, color=Theme.navy)
    steps = ["发现高频任务", "抽成 Skill 目录", "接入触发词/提示", "加入审计与测试", "持续更新迭代"]
    y = 2.85
    for i, t in enumerate(steps, 1):
        add_text(s, 1.35, y, 5.1, 0.25, f"{i}. {t}", size=12)
        y += 0.58

    # right: memory strategies
    add_card(s, 7.0, 2.1, 5.55, 4.7, Theme.light_cyan)
    add_text(s, 7.25, 2.38, 5.1, 0.32, "Memory 策略", size=15, bold=True, color=Theme.cyan)
    add_bullets(
        s,
        7.25,
        2.85,
        5.1,
        3.7,
        [
            "项目记忆：规范、命令、约束写入文档文件。",
            "自动记忆：会话中提炼长期偏好，按需读取。",
            "会话隔离：任务间上下文隔离，避免污染。",
            "可审计编辑：记忆可查看、可修改、可删除。",
        ],
        size=12,
    )
    add_source(s, "Claude Code docs: memory/slash commands/skills concepts")


def slide_8_mcp_arch(prs: Presentation) -> None:
    """构建8MCParch对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "Pillar 4: MCP 架构", "Standard protocol for tool integration", 8)

    add_card(s, 0.75, 1.2, 12.0, 5.9, Theme.white)
    add_text(s, 1.05, 1.55, 11.2, 0.35, "MCP 核心：Host-Client-Server，JSON-RPC，能力协商", size=18, bold=True, color=Theme.blue)

    # Diagram blocks
    add_card(s, 1.0, 2.2, 3.2, 2.4, Theme.light_blue)
    add_text(s, 1.15, 2.45, 2.9, 0.3, "Host", size=14, bold=True, color=Theme.navy, align="center")
    add_text(s, 1.15, 2.85, 2.9, 1.4, "IDE / CLI / App\n权限控制\n上下文聚合", size=12, align="center")

    add_arrow(s, 4.25, 3.0, 0.55, 0.22)
    add_card(s, 4.9, 2.2, 3.2, 2.4, Theme.light_cyan)
    add_text(s, 5.05, 2.45, 2.9, 0.3, "Client", size=14, bold=True, color=Theme.cyan, align="center")
    add_text(s, 5.05, 2.85, 2.9, 1.4, "1:1 Server 会话\n能力协商\n消息路由", size=12, align="center")

    add_arrow(s, 8.15, 3.0, 0.55, 0.22)
    add_card(s, 8.8, 2.2, 3.2, 2.4, Theme.light_green)
    add_text(s, 8.95, 2.45, 2.9, 0.3, "Server", size=14, bold=True, color=Theme.green, align="center")
    add_text(s, 8.95, 2.85, 2.9, 1.4, "Resources\nTools\nPrompts", size=12, align="center")

    add_card(s, 1.0, 5.05, 11.0, 1.55, Theme.light_orange)
    add_bullets(
        s,
        1.2,
        5.3,
        10.5,
        1.15,
        [
            "优势：标准化接入外部能力，避免每个 Agent 写定制胶水代码。",
            "建议：MCP 作为统一接入层，Skill 作为任务编排层，两者分工清晰。",
        ],
        size=12,
    )
    add_source(s, "MCP spec architecture (2024-11-05 / 2025-06-18)")


def slide_9_mcp_security(prs: Presentation) -> None:
    """构建9MCPsecurity对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "Pillar 5: MCP 安全与权限", "Security baseline before scale", 9)

    add_card(s, 0.8, 1.25, 5.9, 5.8, Theme.white)
    add_text(s, 1.05, 1.55, 5.4, 0.34, "MCP 安全基线（官方）", size=17, bold=True, color=Theme.red)
    add_bullets(
        s,
        1.05,
        2.02,
        5.3,
        4.8,
        [
            "所有入站请求必须鉴权校验。",
            "不能把 session 当鉴权机制。",
            "session id 必须不可预测并定期轮换。",
            "会话应绑定用户身份，避免会话劫持。",
            "高风险工具默认关闭，按需授权。",
            "工具调用必须有完整审计轨迹。",
        ],
        size=12,
    )

    add_card(s, 7.0, 1.25, 5.55, 5.8, Theme.white)
    add_text(s, 7.25, 1.55, 5.1, 0.34, "落地控制矩阵", size=17, bold=True, color=Theme.blue)
    rows = [
        ("读取日志", "低风险", "默认允许"),
        ("查询代码", "中风险", "按仓库白名单"),
        ("执行命令", "高风险", "人工确认 + 沙箱"),
        ("写入外部系统", "高风险", "审批流 + 审计"),
    ]
    y = 2.05
    for op, risk, ctrl in rows:
        add_card(s, 7.25, y, 5.1, 0.98, Theme.light_blue if risk != "高风险" else Theme.light_orange)
        add_text(s, 7.45, y + 0.15, 1.4, 0.25, op, size=12, bold=True)
        add_text(s, 8.95, y + 0.15, 1.0, 0.25, risk, size=12, color=Theme.red if risk == "高风险" else Theme.text)
        add_text(s, 10.0, y + 0.15, 2.2, 0.5, ctrl, size=11, color=Theme.subtext)
        y += 1.15
    add_source(s, "MCP security best practices; GitHub Copilot coding agent security docs")


def slide_10_workflow(prs: Presentation) -> None:
    """构建10工作流对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "Pillar 6: 标准工作流", "Explore -> Plan -> Code -> Commit -> Review", 10)

    add_text(s, 0.8, 1.15, 12.0, 0.4, "Anthropic/GitHub/OpenAI 都强调：先计划后实现，持续校验", size=19, bold=True, color=Theme.blue)

    phases = [
        ("Explore", "读代码/文档\n补上下文"),
        ("Plan", "输出分步计划\n确认约束"),
        ("Code", "小步实现\n不跨大改"),
        ("Commit", "结构化提交\n附变更说明"),
        ("Review", "Agent+Human\n双重审查"),
    ]
    x = 0.85
    fills = [Theme.light_blue, Theme.light_cyan, Theme.light_green, Theme.light_orange, Theme.light_red]
    for i, ((name, desc), fill) in enumerate(zip(phases, fills)):
        add_card(s, x, 2.0, 2.35, 3.75, fill)
        add_text(s, x + 0.2, 2.35, 1.95, 0.3, name, size=15, bold=True, align="center", color=Theme.navy)
        add_text(s, x + 0.25, 2.9, 1.85, 1.8, desc, size=12, align="center")
        if i < 4:
            add_arrow(s, x + 2.35, 3.45, 0.35, 0.2)
        x += 2.55
    add_card(s, 0.85, 6.0, 11.9, 1.05, Theme.white)
    add_text(
        s,
        1.05,
        6.28,
        11.4,
        0.5,
        "实操建议：每一步都定义退出条件（DoD），避免“看起来差不多了”的 vibe 式收尾。",
        size=13,
        color=Theme.subtext,
    )
    add_source(s, "Anthropic Claude docs; GitHub Copilot best practices; OpenAI harness engineering")


def slide_11_eval(prs: Presentation) -> None:
    """构建11eval对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "Pillar 7: Eval-Driven Development", "No vibe-based evals", 11)

    add_card(s, 0.8, 1.2, 6.0, 5.85, Theme.white)
    add_text(s, 1.05, 1.5, 5.5, 0.34, "评测飞轮", size=18, bold=True, color=Theme.blue)
    add_chart(
        s,
        chart_type=XL_CHART_TYPE.LINE_MARKERS,
        x=1.05,
        y=1.95,
        w=5.5,
        h=3.2,
        categories=["目标定义", "样本构建", "指标定义", "运行比较", "持续评测"],
        series=[("系统成熟度(示意)", [32, 47, 61, 74, 88])],
    )
    add_bullets(
        s,
        1.05,
        5.35,
        5.5,
        1.4,
        [
            "OpenAI: Adopt eval-driven development",
            "反模式：Vibe-based evals（感觉可用就上线）",
        ],
        size=11,
    )

    add_card(s, 7.05, 1.2, 5.5, 5.85, Theme.white)
    add_text(s, 7.3, 1.5, 5.0, 0.34, "你项目可直接落地", size=18, bold=True, color=Theme.blue)
    add_bullets(
        s,
        7.3,
        1.95,
        5.0,
        4.8,
        [
            "构建 RCA benchmark 场景库（日志+堆栈+现象）",
            "指标：首批证据时延、Top1/Top3 命中、超时率",
            "CI Gate：超时率/失败率超阈值阻断发布",
            "Trace grading：评估 agent 轨迹而不只看最终答案",
            "人审样本回流成规则与 prompt 改进样本",
        ],
        size=12,
    )
    add_source(s, "OpenAI evaluation best practices + trace grading docs")


def slide_12_multi_agent_patterns(prs: Presentation) -> None:
    """构建12multiAgentpatterns对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "Pillar 8: 多 Agent 模式", "Workflow vs Agent + Supervisor", 12)

    add_card(s, 0.8, 1.2, 5.9, 5.9, Theme.white)
    add_text(s, 1.05, 1.5, 5.4, 0.34, "LangGraph 的两类模式", size=17, bold=True, color=Theme.blue)
    add_bullets(
        s,
        1.05,
        1.95,
        5.3,
        4.9,
        [
            "Workflow：固定路径，顺序可控，适合稳定流程。",
            "Agent：动态路径，工具自决，适合复杂开放问题。",
            "生产实践：混合模式最常见（主流程固定 + 局部动态）。",
            "关键能力：persistence、streaming、debuggability、deployment。",
        ],
        size=12,
    )

    add_card(s, 7.0, 1.2, 5.55, 5.9, Theme.white)
    add_text(s, 7.25, 1.5, 5.1, 0.34, "推荐编排拓扑", size=17, bold=True, color=Theme.blue)
    add_card(s, 8.8, 2.5, 2.0, 0.9, Theme.light_blue)
    add_text(s, 8.95, 2.82, 1.7, 0.25, "Supervisor", size=12, bold=True, align="center")
    add_card(s, 7.4, 4.2, 1.8, 0.8, Theme.light_cyan)
    add_card(s, 9.35, 4.2, 1.8, 0.8, Theme.light_green)
    add_card(s, 11.3, 4.2, 1.0, 0.8, Theme.light_orange)
    add_text(s, 7.55, 4.45, 1.5, 0.2, "Log", size=11, align="center")
    add_text(s, 9.5, 4.45, 1.5, 0.2, "Code", size=11, align="center")
    add_text(s, 11.35, 4.45, 0.8, 0.2, "DB", size=11, align="center")
    add_arrow(s, 9.7, 3.4, 0.2, 0.65)
    add_arrow(s, 8.25, 3.4, 0.2, 0.65)
    add_arrow(s, 11.05, 3.4, 0.2, 0.65)
    add_text(s, 7.25, 5.25, 5.0, 1.2, "原则：并行分析 + 协议化交流 + 明确停止条件 + 裁决收敛", size=12)
    add_source(s, "LangGraph workflows-agents docs; langgraph supervisor references")


def slide_13_reliability(prs: Presentation) -> None:
    """构建13reliability对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "Pillar 9: 可靠性（Checkpoint + Durability）", "Never stuck in pending", 13)

    add_card(s, 0.8, 1.2, 12.0, 5.9, Theme.white)
    add_text(s, 1.05, 1.5, 11.4, 0.34, "LangGraph 可靠性能力：线程状态、断点恢复、持久化策略", size=18, bold=True, color=Theme.blue)

    add_card(s, 1.05, 2.1, 5.6, 2.2, Theme.light_blue)
    add_text(s, 1.25, 2.35, 5.2, 0.25, "Persistence", size=14, bold=True, color=Theme.navy)
    add_bullets(
        s,
        1.25,
        2.7,
        5.2,
        1.4,
        [
            "每个 super-step 生成 checkpoint",
            "thread_id 支撑会话记忆与恢复",
            "支持 time-travel / human-in-the-loop",
        ],
        size=11,
    )

    add_card(s, 1.05, 4.6, 5.6, 2.2, Theme.light_cyan)
    add_text(s, 1.25, 4.85, 5.2, 0.25, "Durability", size=14, bold=True, color=Theme.cyan)
    add_bullets(
        s,
        1.25,
        5.2,
        5.2,
        1.4,
        [
            "exit: 性能优先",
            "async: 性能与安全平衡",
            "sync: 最强一致性（代价更高）",
        ],
        size=11,
    )

    add_card(s, 7.0, 2.1, 5.55, 4.7, Theme.light_green)
    add_text(s, 7.2, 2.35, 5.1, 0.25, "系统设计建议", size=14, bold=True, color=Theme.green)
    add_bullets(
        s,
        7.2,
        2.72,
        5.0,
        3.7,
        [
            "关键调查会话使用 sync 或 async，不建议 exit。",
            "每个阶段都必须有终态输出（成功/失败/降级）。",
            "失败时记录“可恢复上下文”并支持重跑。",
            "恢复流程要有幂等设计，避免重复执行副作用。",
            "所有 timeout 都要映射为可解释的事件和报告条目。",
        ],
        size=12,
    )
    add_source(s, "LangGraph persistence + durable-execution docs")


def slide_14_tool_compare(prs: Presentation) -> None:
    """构建14工具compare对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "业界工具链对比", "What to borrow from whom", 14)

    add_card(s, 0.8, 1.2, 12.0, 5.9, Theme.white)
    add_text(s, 1.05, 1.5, 11.4, 0.34, "能力对比矩阵（面向生产落地）", size=18, bold=True, color=Theme.blue)

    cols = ["平台", "强项", "治理特点", "可借鉴点"]
    x_list = [1.05, 3.0, 5.5, 8.3]
    widths = [1.8, 2.4, 2.7, 3.9]
    for x, w, c in zip(x_list, widths, cols):
        add_card(s, x, 1.95, w, 0.6, Theme.light_blue)
        add_text(s, x + 0.08, 2.13, w - 0.16, 0.2, c, size=12, bold=True, align="center")

    rows = [
        ("OpenAI Codex", "Harness与执行闭环", "系统规则先于模型", "仓库知识系统化、持续垃圾回收"),
        ("Claude Code", "终端Agent执行力", "权限模式与记忆体系", "CLAUDE.md + /memory + 子agent"),
        ("GitHub Copilot", "GitHub工作流融合", "分支/权限/防火墙", "仓库指令文件 + 会话日志"),
        ("Sourcegraph Cody", "上下文检索能力", "agentic context反思", "自动补上下文 + 工具联动"),
    ]
    y = 2.7
    for i, r in enumerate(rows):
        fill = Theme.white if i % 2 == 0 else Theme.light_cyan
        for (x, w, v) in zip(x_list, widths, r):
            add_card(s, x, y, w, 0.92, fill)
            add_text(s, x + 0.08, y + 0.16, w - 0.16, 0.55, v, size=11, align="center")
        y += 1.02
    add_source(s, "OpenAI/Anthropic/GitHub/Sourcegraph official docs and blogs")


def slide_15_fail_modes(prs: Presentation) -> None:
    """构建15failmodes对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "失败模式与反模式", "What breaks in real projects", 15)

    add_card(s, 0.8, 1.2, 5.9, 5.9, Theme.white)
    add_text(s, 1.05, 1.5, 5.4, 0.34, "Thoughtworks 实验启示", size=17, bold=True, color=Theme.blue)
    add_bullets(
        s,
        1.05,
        1.95,
        5.4,
        4.9,
        [
            "纯 vibe 模式可快速出原型，但增量变更易回归。",
            "加强约束（TDD、小步提交）后可维护性显著提升。",
            "无纪律工具调用会引入不可预期依赖与版本漂移。",
            "人机对话式协作优于“把任务一次性全丢给AI”。",
        ],
        size=12,
    )

    add_card(s, 7.0, 1.2, 5.55, 5.9, Theme.white)
    add_text(s, 7.25, 1.5, 5.1, 0.34, "高频反模式清单", size=17, bold=True, color=Theme.red)
    patterns = [
        ("无评测上线", Theme.light_red),
        ("工具权限失控", Theme.light_orange),
        ("上下文污染", Theme.light_red),
        ("会话无终态", Theme.light_orange),
        ("规则漂移", Theme.light_red),
    ]
    y = 2.0
    for p, fill in patterns:
        add_card(s, 7.25, y, 5.1, 0.8, fill)
        add_text(s, 7.45, y + 0.23, 4.7, 0.25, p, size=12, bold=True)
        y += 0.95
    add_source(s, "Thoughtworks production-grade vibe coding experiments")


def slide_16_governance(prs: Presentation) -> None:
    """构建16治理对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "治理模型", "Policy, permission, audit", 16)

    add_card(s, 0.8, 1.2, 12.0, 5.9, Theme.white)
    add_text(s, 1.05, 1.5, 11.4, 0.34, "治理是放大自动化收益的前提", size=18, bold=True, color=Theme.blue)

    layers = [
        ("组织策略层", "谁能触发 agent、哪些仓库可用、成本预算"),
        ("权限执行层", "工具 allowlist、分支保护、防火墙、审批流"),
        ("运行审计层", "会话日志、工具调用、变更追踪、回放"),
        ("质量门禁层", "单测/集成/benchmark/eval 多重门禁"),
    ]
    y = 2.0
    fills = [Theme.light_blue, Theme.light_cyan, Theme.light_green, Theme.light_orange]
    for (name, desc), fill in zip(layers, fills):
        add_card(s, 1.1, y, 11.4, 0.95, fill)
        add_text(s, 1.35, y + 0.2, 2.2, 0.3, name, size=13, bold=True, color=Theme.navy)
        add_text(s, 3.7, y + 0.2, 8.5, 0.45, desc, size=12)
        y += 1.12
    add_source(s, "GitHub Copilot coding agent security docs; MCP security best practices")


def slide_17_project_map(prs: Presentation) -> None:
    """构建17项目map对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "本项目映射：架构层", "How current project aligns", 17)

    add_card(s, 0.8, 1.2, 12.0, 5.9, Theme.white)
    add_text(s, 1.05, 1.5, 11.4, 0.34, "当前项目已具备的工程化骨架", size=18, bold=True, color=Theme.blue)

    # House-like layers
    add_card(s, 1.1, 5.8, 11.3, 0.9, Theme.light_green)
    add_text(s, 1.3, 6.1, 10.9, 0.26, "Foundation: session store / audit logs / report store / config", size=12)
    add_card(s, 1.5, 4.6, 10.5, 1.0, Theme.light_cyan)
    add_text(s, 1.7, 4.9, 10.1, 0.25, "Runtime: LangGraph orchestrator + routing + phase execution", size=12)
    add_card(s, 2.0, 3.4, 9.5, 1.0, Theme.light_blue)
    add_text(s, 2.2, 3.7, 9.1, 0.25, "Agents: Main / Log / Code / Domain / DB / Critic / Judge", size=12)
    add_card(s, 2.6, 2.2, 8.3, 1.0, Theme.light_orange)
    add_text(s, 2.8, 2.5, 7.9, 0.25, "Interfaces: FastAPI REST + WebSocket + Frontend Incident UI", size=12)
    add_text(s, 2.6, 1.9, 8.3, 0.22, "Roof: AGENTS.md constraints + plans + CI gates", size=11, align="center", color=Theme.subtext)
    add_source(s, "Project files: AGENTS.md, README.md, docs/wiki/code-wiki.md")


def slide_18_project_gap(prs: Presentation) -> None:
    """构建18项目gap对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "本项目映射：改进缺口", "Gap analysis from industry baseline", 18)

    add_card(s, 0.8, 1.2, 12.0, 5.9, Theme.white)
    add_text(s, 1.05, 1.5, 11.4, 0.34, "对标后优先补齐的能力", size=18, bold=True, color=Theme.blue)

    items = [
        ("P0", "评测闭环加强", "benchmark 场景数、CI 阈值、trace grading 需完善"),
        ("P1", "工具治理统一", "MCP 安全策略与工具权限矩阵标准化"),
        ("P1", "上下文治理", "规则分层与自动记忆冲突检测需增强"),
        ("P2", "多 Agent 协议", "agent 间消息协议、停止条件、置信度机制"),
        ("P2", "产品体验", "报告可视化、证据链可回放、阶段解释能力"),
    ]
    y = 2.0
    for idx, (p, title, desc) in enumerate(items):
        fill = Theme.light_blue if idx % 2 == 0 else Theme.white
        add_card(s, 1.05, y, 11.5, 0.95, fill)
        add_text(s, 1.25, y + 0.18, 0.6, 0.25, p, size=12, bold=True, color=Theme.orange)
        add_text(s, 1.95, y + 0.18, 2.2, 0.25, title, size=12, bold=True)
        add_text(s, 4.35, y + 0.18, 7.9, 0.45, desc, size=11, color=Theme.subtext)
        y += 1.05
    add_source(s, "Derived from industry docs + local project code/wiki")


def slide_19_roadmap(prs: Presentation) -> None:
    """构建19路线图对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "落地路线图", "2 weeks + 90 days", 19)

    add_card(s, 0.8, 1.2, 6.0, 5.9, Theme.white)
    add_text(s, 1.05, 1.5, 5.5, 0.34, "2 周快跑（先把可用性做硬）", size=17, bold=True, color=Theme.blue)
    add_bullets(
        s,
        1.05,
        1.95,
        5.5,
        4.9,
        [
            "Week1: 统一门禁、统一审计、统一终态保证",
            "Week1: 标准化工具配置（开关、权限、超时）",
            "Week2: benchmark 接入 CI 阻断",
            "Week2: 前端战情页与证据链视图增强",
            "Week2: 真实故障演练与复盘模板",
        ],
        size=12,
    )

    add_card(s, 7.05, 1.2, 5.75, 5.9, Theme.white)
    add_text(s, 7.3, 1.5, 5.2, 0.34, "90 天演进（P0-P3）", size=17, bold=True, color=Theme.blue)
    add_chart(
        s,
        chart_type=XL_CHART_TYPE.LINE_MARKERS,
        x=7.35,
        y=1.95,
        w=5.0,
        h=2.8,
        categories=["P0", "P1", "P2", "P3"],
        series=[("成熟度", [38, 57, 73, 86])],
    )
    add_bullets(
        s,
        7.35,
        4.95,
        5.0,
        2.0,
        [
            "P0 可用性 / P1 正确率 / P2 可控修复 / P3 持续学习",
            "每阶段必须有量化验收指标",
        ],
        size=11,
    )
    add_source(s, "Roadmap synthesized from project plans + industry practices")


def slide_20_kpi(prs: Presentation) -> None:
    """构建20kpi对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "指标体系与 ROI", "What to measure weekly", 20)

    add_card(s, 0.8, 1.2, 12.0, 5.9, Theme.white)
    add_text(s, 1.05, 1.5, 11.4, 0.34, "推荐 KPI（当前 vs 目标，示意）", size=18, bold=True, color=Theme.blue)
    add_chart(
        s,
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        x=1.05,
        y=1.95,
        w=7.2,
        h=3.6,
        categories=["首批证据时延", "Top1命中", "超时率", "长期pending", "人工介入率"],
        series=[("当前", [55, 42, 18, 12, 63]), ("目标", [15, 65, 6, 1, 35])],
    )

    add_card(s, 8.45, 1.95, 4.2, 3.6, Theme.light_cyan)
    add_text(s, 8.7, 2.2, 3.7, 0.3, "管理层关注点", size=14, bold=True, color=Theme.cyan)
    add_bullets(
        s,
        8.7,
        2.55,
        3.7,
        2.7,
        [
            "交付速度是否稳定提升",
            "回归和事故率是否下降",
            "人力投入是否转向高价值决策",
            "安全风险是否可控且可追责",
        ],
        size=11,
    )
    add_source(s, "KPI framework aligned with OpenAI eval guidance and SRE RCA goals")


def slide_21_refs(prs: Presentation) -> None:
    """构建21参考资料对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "参考资料", "Primary sources used", 21)

    add_card(s, 0.8, 1.2, 12.0, 5.9, Theme.white)
    add_text(s, 1.05, 1.5, 11.4, 0.34, "核心引用（官方优先）", size=18, bold=True, color=Theme.blue)
    refs = [
        "OpenAI: Harness engineering (2026-02-11)",
        "OpenAI API Docs: Evaluation best practices / Trace grading",
        "Anthropic/Claude Code Docs: memory, slash commands, MCP, workflows",
        "GitHub Docs: Copilot best practices / coding agent security",
        "LangGraph Docs: workflows-agents, persistence, durable execution",
        "MCP Spec + Security best practices",
        "Sourcegraph Docs: agentic context fetching",
        "Thoughtworks: production-grade vibe coding experiments",
    ]
    y = 2.0
    for r in refs:
        add_text(s, 1.15, y, 11.0, 0.28, f"• {r}", size=12)
        y += 0.55
    add_text(s, 1.15, 6.45, 11.0, 0.35, "完整链接清单见输出目录 SOURCES.md", size=12, bold=True, color=Theme.subtext)


def slide_22_close(prs: Presentation) -> None:
    """构建22关闭对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs)
    add_topbar(s, "结语", "Execution first", 22)

    add_card(s, 0.8, 1.6, 12.0, 1.8, Theme.light_blue)
    add_text(
        s,
        1.1,
        2.05,
        11.4,
        0.65,
        "结论：Vibe Coding 的成败，取决于你是否把“系统约束、评测门禁、权限审计”做成产品能力。",
        size=21,
        bold=True,
        color=Theme.navy,
    )

    add_card(s, 0.8, 3.9, 12.0, 2.9, Theme.white)
    add_text(s, 1.05, 4.2, 11.4, 0.35, "下一步建议（立刻执行）", size=17, bold=True, color=Theme.blue)
    add_bullets(
        s,
        1.05,
        4.62,
        11.2,
        2.2,
        [
            "1) 用本项目挑一个真实故障场景做完整演练（输入->分析->报告->复盘）。",
            "2) 把 benchmark gate 正式挂进 CI，并设定阻断阈值。",
            "3) 完成 MCP 工具权限矩阵与审计看板。",
            "4) 固化 AGENTS.md + 计划清单 + 每周评测复盘节奏。",
        ],
        size=14,
    )
    add_source(s, "Synthesis of industry best practices and this project experience")


def write_sources() -> None:
    """写出来源相关产物。"""
    
    OUT_MD.write_text(
        "\n".join(
            [
                "# Sources",
                "",
                "1. OpenAI Harness Engineering",
                "   - https://openai.com/index/harness-engineering/",
                "2. OpenAI Evaluation Best Practices",
                "   - https://developers.openai.com/api/docs/guides/evaluation-best-practices",
                "3. OpenAI Trace Grading",
                "   - https://platform.openai.com/docs/guides/trace-grading",
                "4. Claude Code Docs (overview/memory/slash commands)",
                "   - https://code.claude.com/docs",
                "   - https://code.claude.com/docs/en/memory",
                "   - https://code.claude.com/docs/en/slash-commands",
                "5. GitHub Copilot Best Practices",
                "   - https://docs.github.com/en/copilot/get-started/best-practices",
                "6. GitHub Copilot Coding Agent Security",
                "   - https://docs.github.com/en/copilot/concepts/agents/coding-agent/about-coding-agent",
                "   - https://docs.github.com/en/copilot/responsible-use/copilot-coding-agent",
                "7. LangGraph Docs",
                "   - https://docs.langchain.com/oss/python/langgraph/workflows-agents",
                "   - https://docs.langchain.com/oss/python/langgraph/persistence",
                "   - https://docs.langchain.com/oss/python/langgraph/durable-execution",
                "8. MCP Spec + Security",
                "   - https://modelcontextprotocol.io/specification/2024-11-05/architecture/index",
                "   - https://modelcontextprotocol.io/specification/2025-06-18/basic/index",
                "   - https://modelcontextprotocol.io/docs/tutorials/security/security_best_practices",
                "9. Sourcegraph Agentic Context",
                "   - https://sourcegraph.com/docs/cody/capabilities/agentic-context-fetching",
                "10. Thoughtworks Vibe Coding Experiment",
                "   - https://www.thoughtworks.com/insights/blog/generative-ai/can-vibe-coding-produce-production-grade-software",
                "",
            ]
        ),
        encoding="utf-8",
    )


def main() -> None:
    """执行脚本主流程，串联参数解析、内容生成与结果输出。"""
    
    ensure_dirs()
    prs = prs_new()
    slide_1_cover(prs)
    slide_2_agenda(prs)
    slide_3_define(prs)
    slide_4_industry_signals(prs)
    slide_5_harness(prs)
    slide_6_context(prs)
    slide_7_skill_memory(prs)
    slide_8_mcp_arch(prs)
    slide_9_mcp_security(prs)
    slide_10_workflow(prs)
    slide_11_eval(prs)
    slide_12_multi_agent_patterns(prs)
    slide_13_reliability(prs)
    slide_14_tool_compare(prs)
    slide_15_fail_modes(prs)
    slide_16_governance(prs)
    slide_17_project_map(prs)
    slide_18_project_gap(prs)
    slide_19_roadmap(prs)
    slide_20_kpi(prs)
    slide_21_refs(prs)
    slide_22_close(prs)
    prs.save(OUT_PPTX)
    write_sources()
    print(f"Generated: {OUT_PPTX}")
    print(f"Sources: {OUT_MD}")


if __name__ == "__main__":
    main()

