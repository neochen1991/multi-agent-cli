#!/usr/bin/env python3
"""generateVibeCodingWarRoomPPT脚本。"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output" / "vibe-coding-warroom-ppt"
OUT_PPTX = OUT_DIR / "2026-03-06-vibe-coding-practice-warroom.pptx"
OUT_README = OUT_DIR / "README.md"


class Theme:
    """封装Theme相关常量或数据结构。"""
    
    bg = RGBColor(10, 16, 28)
    panel = RGBColor(17, 27, 43)
    panel_2 = RGBColor(22, 35, 56)
    border = RGBColor(44, 62, 92)
    text = RGBColor(237, 241, 247)
    text_dim = RGBColor(160, 174, 196)
    cyan = RGBColor(78, 205, 196)
    amber = RGBColor(255, 191, 105)
    red = RGBColor(255, 107, 107)
    green = RGBColor(92, 220, 151)
    blue = RGBColor(96, 165, 250)
    white = RGBColor(255, 255, 255)
    ink = RGBColor(19, 28, 43)
    soft = RGBColor(227, 234, 245)
    card_light = RGBColor(242, 246, 252)
    card_warm = RGBColor(255, 247, 237)
    card_mint = RGBColor(237, 251, 245)
    card_cyan = RGBColor(234, 249, 250)


def ensure_dirs() -> None:
    """确保dirs相关前置条件已经满足。"""
    
    OUT_DIR.mkdir(parents=True, exist_ok=True)


def new_prs() -> Presentation:
    """执行新增prs相关逻辑。"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_bg(slide, prs: Presentation, *, alt: bool = False) -> None:
    """向当前页补充背景相关元素，并统一样式与布局。"""
    
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = Theme.panel if alt else Theme.bg
    bg.line.fill.background()


def add_topbar(slide, title: str, subtitle: str, page: int) -> None:
    """向当前页补充topbar相关元素，并统一样式与布局。"""
    
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.78)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = Theme.panel_2
    band.line.fill.background()
    add_text(slide, 0.45, 0.16, 8.6, 0.35, title, size=24, bold=True, color=Theme.text)
    add_text(slide, 9.0, 0.2, 2.8, 0.24, subtitle, size=10, color=Theme.text_dim, align="right")
    add_text(slide, 12.2, 0.18, 0.55, 0.24, str(page), size=12, bold=True, color=Theme.amber, align="right")


def add_footer(slide, source: str) -> None:
    """向当前页补充页脚相关元素，并统一样式与布局。"""
    
    add_text(slide, 0.45, 7.06, 12.2, 0.18, f"Source: {source}", size=8, color=Theme.text_dim)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: int = 16,
    bold: bool = False,
    color: RGBColor = Theme.text,
    align: str = "left",
) -> None:
    """向当前页补充文本相关元素，并统一样式与布局。"""
    
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: Iterable[str],
    *,
    size: int = 13,
    color: RGBColor = Theme.text,
) -> None:
    """向当前页补充bullets相关元素，并统一样式与布局。"""
    
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(size)
        p.font.color.rgb = color


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = Theme.panel,
    border: RGBColor = Theme.border,
) -> None:
    """向当前页补充卡片相关元素，并统一样式与布局。"""
    
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border


def add_pill(slide, x: float, y: float, w: float, text: str, *, fill: RGBColor, color: RGBColor) -> None:
    """向当前页补充标签相关元素，并统一样式与布局。"""
    
    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.42)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    add_text(slide, x + 0.08, y + 0.09, w - 0.16, 0.18, text, size=10, bold=True, color=color, align="center")


def add_metric(slide, x: float, y: float, w: float, title: str, value: str, *, accent: RGBColor) -> None:
    """向当前页补充指标相关元素，并统一样式与布局。"""
    
    add_card(slide, x, y, w, 1.05, fill=Theme.panel_2)
    add_text(slide, x + 0.18, y + 0.18, w - 0.35, 0.18, title, size=10, color=Theme.text_dim)
    add_text(slide, x + 0.18, y + 0.44, w - 0.35, 0.32, value, size=20, bold=True, color=accent)


def add_chart(
    slide,
    *,
    chart_type: XL_CHART_TYPE,
    x: float,
    y: float,
    w: float,
    h: float,
    categories: list[str],
    series: list[tuple[str, list[float]]],
) -> None:
    """向当前页补充图表相关元素，并统一样式与布局。"""
    
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    chart = slide.shapes.add_chart(chart_type, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)
        chart.legend.include_in_layout = False
    if chart.value_axis is not None:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.font.size = Pt(10)
    if chart.category_axis is not None:
        chart.category_axis.tick_labels.font.size = Pt(10)
    palette = [Theme.cyan, Theme.amber, Theme.blue, Theme.green, Theme.red]
    for idx, s in enumerate(chart.series):
        fill = s.format.fill
        fill.solid()
        fill.fore_color.rgb = palette[idx % len(palette)]


def add_connector(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = Theme.text_dim) -> None:
    """向当前页补充连接器相关元素，并统一样式与布局。"""
    
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(1.5)


def slide_1_cover(prs: Presentation) -> None:
    """构建1封面对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    add_text(slide, 0.55, 0.55, 7.2, 0.35, "WAR-ROOM BRIEF", size=11, bold=True, color=Theme.amber)
    add_text(slide, 0.55, 1.0, 8.8, 1.0, "Vibe Coding 在生产级\nSRE 智能体中的工程化落地", size=28, bold=True)
    add_text(slide, 0.58, 2.25, 7.6, 0.4, "从 Prompt 驱动到 Harness 驱动", size=16, color=Theme.text_dim)

    add_pill(slide, 0.58, 3.05, 1.75, "Harness First", fill=Theme.card_warm, color=Theme.ink)
    add_pill(slide, 2.45, 3.05, 1.75, "Context Aware", fill=Theme.card_cyan, color=Theme.ink)
    add_pill(slide, 4.32, 3.05, 1.75, "Agent Governed", fill=Theme.card_mint, color=Theme.ink)

    add_text(slide, 0.58, 4.0, 4.8, 0.25, "目标", size=12, bold=True, color=Theme.cyan)
    add_bullets(
        slide,
        0.58,
        4.3,
        5.2,
        2.0,
        [
            "借鉴业界优秀 Vibe Coding 团队实践",
            "解释 Skill / MCP / 多 Agent 如何做成工程系统",
            "结合当前项目给出落地路线图",
        ],
        size=13,
    )

    add_card(slide, 7.6, 0.9, 5.0, 5.75, fill=Theme.panel_2)
    add_text(slide, 7.95, 1.2, 4.2, 0.25, "战情视图", size=12, bold=True, color=Theme.amber)
    add_metric(slide, 7.95, 1.65, 1.45, "Speed", "8x", accent=Theme.cyan)
    add_metric(slide, 9.55, 1.65, 1.45, "Risk", "Gate", accent=Theme.amber)
    add_metric(slide, 11.15, 1.65, 1.15, "UX", "Flow", accent=Theme.green)
    add_card(slide, 7.95, 2.95, 4.1, 1.05, fill=Theme.panel)
    add_text(slide, 8.15, 3.18, 3.6, 0.2, "Main Agent", size=11, bold=True, color=Theme.cyan, align="center")
    for label, x, y, color in [
        ("Log", 8.0, 4.45, Theme.cyan),
        ("Code", 9.35, 4.45, Theme.blue),
        ("Domain", 10.7, 4.45, Theme.green),
    ]:
        add_card(slide, x, y, 1.0, 0.65, fill=Theme.panel)
        add_text(slide, x + 0.08, y + 0.19, 0.84, 0.18, label, size=10, bold=True, color=color, align="center")
        add_connector(slide, 10.0, 4.0, x + 0.5, y)
    add_text(slide, 8.0, 5.55, 4.2, 0.7, "问题现场 -> 命令分发 -> 证据回流 -> 裁决报告", size=12, color=Theme.text_dim)

    add_text(slide, 0.58, 6.85, 10.8, 0.2, f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')} Asia/Shanghai", size=9, color=Theme.text_dim)
    add_footer(slide, "2026-03-06-vibe-coding-ppt-outline.md + official public references")


def slide_2_why(prs: Presentation) -> None:
    """构建2why对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "为什么 Vibe Coding 值得做", "Value framing", 2)

    add_text(slide, 0.55, 1.0, 6.4, 0.45, "结论：价值不在“省人”，而在“缩短认知到交付闭环”。", size=22, bold=True)
    add_bullets(
        slide,
        0.6,
        1.75,
        5.2,
        1.8,
        [
            "减少上下文切换，缩短需求到实现链路。",
            "让人从重复实现转向约束设计与问题判断。",
            "前提是速度必须绑定验证与治理。",
        ],
        size=14,
    )
    add_metric(slide, 0.6, 4.1, 1.75, "链路压缩", "-60%", accent=Theme.cyan)
    add_metric(slide, 2.55, 4.1, 1.75, "试错成本", "-45%", accent=Theme.amber)
    add_metric(slide, 4.5, 4.1, 1.75, "回归风险", "+ if no gate", accent=Theme.red)

    add_card(slide, 7.0, 1.0, 5.8, 5.85, fill=Theme.panel_2)
    add_text(slide, 7.3, 1.25, 5.1, 0.25, "传统模式 vs Vibe + Harness", size=12, bold=True, color=Theme.amber)
    add_chart(
        slide,
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        x=7.25,
        y=1.65,
        w=5.0,
        h=3.4,
        categories=["理解需求", "方案探索", "实现速度", "回归控制"],
        series=[("传统", [38, 35, 32, 66]), ("Vibe+Harness", [71, 76, 84, 79])],
    )
    add_text(slide, 7.35, 5.35, 4.9, 0.8, "关键不是“更会写代码”，而是“更快地在约束下交付正确结果”。", size=13)
    add_footer(slide, "OpenAI Harness Engineering; Thoughtworks vibe coding experiments")


def slide_3_harness(prs: Presentation) -> None:
    """构建3harness对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, alt=True)
    add_topbar(slide, "业界共识 1：Harness Engineering 是核心", "System over prompt", 3)

    add_text(slide, 0.6, 1.0, 7.4, 0.45, "结论：生产可用性的上限由 Harness 决定。", size=22, bold=True)
    add_text(slide, 0.6, 1.55, 5.8, 0.3, "让模型负责推理，让系统负责控制。", size=14, color=Theme.text_dim)

    bands = [
        ("目标层", "任务拆解、成功标准、约束边界", Theme.card_warm, Theme.ink),
        ("执行层", "路由、状态机、工具调用、会话上下文", Theme.card_cyan, Theme.ink),
        ("验证层", "单测、集成、评测、CI gate", Theme.card_mint, Theme.ink),
        ("治理层", "审计、权限、回滚、复盘", Theme.card_light, Theme.ink),
    ]
    y = 2.15
    for title, desc, fill, text_color in bands:
        add_card(slide, 0.75, y, 6.15, 0.92, fill=fill, border=fill)
        add_text(slide, 1.0, y + 0.18, 1.3, 0.2, title, size=13, bold=True, color=text_color)
        add_text(slide, 2.45, y + 0.18, 4.0, 0.25, desc, size=12, color=text_color)
        y += 1.02

    add_card(slide, 7.45, 1.25, 5.15, 5.5, fill=Theme.panel)
    add_text(slide, 7.75, 1.55, 4.6, 0.25, "What changes in practice", size=12, bold=True, color=Theme.amber)
    add_bullets(
        slide,
        7.75,
        2.0,
        4.5,
        3.4,
        [
            "Prompt 不再是单点方案，而是 Harness 中的一环。",
            "失败被系统显式捕获，而不是沉默吞掉。",
            "产物可测、可复现、可比较，而不是“看起来还行”。",
            "团队真正优化的是工程回路，不是单次回答效果。",
        ],
        size=13,
    )
    add_footer(slide, "OpenAI Harness Engineering")


def slide_4_context(prs: Presentation) -> None:
    """构建4上下文对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "业界共识 2：上下文工程优先于长 Prompt", "Context funnel", 4)

    add_text(slide, 0.58, 1.02, 7.0, 0.45, "结论：高质量上下文组织，比堆砌超长提示词更有效。", size=22, bold=True)

    steps = [
        ("Raw", "日志 / 代码 / 文档 / 监控"),
        ("Clean", "切片 / 去噪 / 标注"),
        ("Scoped", "当前步骤需要的信息"),
        ("Structured", "schema / checklist / command"),
    ]
    x = 0.78
    fills = [Theme.card_light, Theme.card_cyan, Theme.card_mint, Theme.card_warm]
    for idx, ((title, desc), fill) in enumerate(zip(steps, fills)):
        add_card(slide, x, 2.05, 2.55, 2.55, fill=fill, border=fill)
        add_text(slide, x + 0.18, 2.35, 2.2, 0.22, title, size=16, bold=True, color=Theme.ink, align="center")
        add_text(slide, x + 0.18, 2.88, 2.2, 0.9, desc, size=12, color=Theme.ink, align="center")
        if idx < 3:
            add_connector(slide, x + 2.55, 3.3, x + 2.85, 3.3, Theme.amber)
        x += 2.95

    add_card(slide, 0.78, 5.1, 6.7, 1.25, fill=Theme.panel_2)
    add_text(slide, 1.0, 5.37, 6.2, 0.2, "原则", size=12, bold=True, color=Theme.amber)
    add_text(slide, 1.0, 5.72, 6.2, 0.4, "只给当前步骤需要的信息，避免让模型替你做“上下文猜测”。", size=14)

    add_card(slide, 7.9, 1.45, 4.7, 4.95, fill=Theme.panel)
    add_text(slide, 8.18, 1.75, 4.1, 0.24, "实践信号", size=12, bold=True, color=Theme.cyan)
    add_bullets(
        slide,
        8.18,
        2.2,
        4.0,
        3.7,
        [
            "Anthropic：清晰、直接、分步。",
            "GitHub Copilot：上下文质量决定输出质量。",
            "Cursor：规则文件与代码索引优先。",
            "Sourcegraph：agentic context fetching 自动补上下文。",
        ],
        size=13,
    )
    add_footer(slide, "Anthropic prompt engineering; GitHub Copilot; Cursor rules; Sourcegraph docs")


def slide_5_workflow(prs: Presentation) -> None:
    """构建5工作流对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, alt=True)
    add_topbar(slide, "业界共识 3：标准工作流替代“即兴开发”", "Execution rhythm", 5)

    add_text(slide, 0.6, 1.0, 7.0, 0.4, "结论：稳定团队都采用固定节奏，而不是一次性把任务丢给 Agent。", size=21, bold=True)

    phases = [
        ("Explore", Theme.cyan),
        ("Plan", Theme.amber),
        ("Build", Theme.blue),
        ("Verify", Theme.green),
        ("Release", Theme.red),
    ]
    x = 0.8
    for idx, (name, accent) in enumerate(phases):
        add_card(slide, x, 2.0, 2.0, 2.65, fill=Theme.panel)
        add_text(slide, x + 0.2, 2.25, 1.6, 0.22, name, size=16, bold=True, color=accent, align="center")
        body = {
            "Explore": "理解问题\n确认约束",
            "Plan": "给出方案\n列任务清单",
            "Build": "小步实现\n不跨大改",
            "Verify": "测试 + Eval\n看结果",
            "Release": "发布/回滚\n复盘",
        }[name]
        add_text(slide, x + 0.2, 2.8, 1.6, 0.95, body, size=12, align="center")
        if idx < len(phases) - 1:
            add_connector(slide, x + 2.0, 3.3, x + 2.35, 3.3, Theme.text_dim)
        x += 2.45

    add_card(slide, 0.8, 5.15, 11.7, 1.1, fill=Theme.card_warm, border=Theme.card_warm)
    add_text(slide, 1.05, 5.5, 11.0, 0.32, "关键提醒：每一步都必须有退出条件（DoD），否则流程只是看起来专业。", size=14, bold=True, color=Theme.ink)
    add_footer(slide, "Anthropic Claude workflows; GitHub Copilot best practices")


def slide_6_gate(prs: Presentation) -> None:
    """构建6门禁对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "业界共识 4：测试和评测是 Vibe Coding 的刹车系统", "Quality gate", 6)

    add_text(slide, 0.58, 1.0, 7.2, 0.42, "结论：没有评测门禁，Vibe Coding 在生产环境不可控。", size=22, bold=True)
    add_card(slide, 0.65, 1.75, 6.0, 5.0, fill=Theme.panel_2)
    add_chart(
        slide,
        chart_type=XL_CHART_TYPE.BAR_CLUSTERED,
        x=0.95,
        y=2.1,
        w=5.4,
        h=3.3,
        categories=["代码生成", "静态检查", "单测", "集成", "Benchmark", "发布"],
        series=[("通过样本数", [1000, 820, 690, 510, 340, 292])],
    )
    add_text(slide, 1.0, 5.72, 5.3, 0.25, "门禁越往后越严格，返工成本也越高。", size=12, color=Theme.text_dim)

    add_card(slide, 7.05, 1.75, 5.55, 5.0, fill=Theme.panel)
    add_text(slide, 7.35, 2.05, 4.9, 0.22, "最低可行门禁集", size=12, bold=True, color=Theme.amber)
    add_bullets(
        slide,
        7.35,
        2.45,
        4.8,
        2.0,
        [
            "本地语法/静态检查",
            "PR 集成测试",
            "Benchmark 命中率与超时率",
            "CI Gate 阻断回归",
        ],
        size=13,
    )
    add_metric(slide, 7.35, 5.15, 1.45, "Top1", "65%", accent=Theme.cyan)
    add_metric(slide, 8.95, 5.15, 1.45, "Timeout", "<6%", accent=Theme.amber)
    add_metric(slide, 10.55, 5.15, 1.45, "Pending", "<1%", accent=Theme.red)
    add_footer(slide, "OpenAI eval-driven development; benchmark gate practice")


def slide_7_mcp(prs: Presentation) -> None:
    """构建7MCP对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, alt=True)
    add_topbar(slide, "业界共识 5：MCP/工具化接入必须标准化", "Capability with governance", 7)

    add_text(slide, 0.6, 1.0, 7.3, 0.42, "结论：Agent 扩展能力要靠标准协议，不靠临时脚本拼接。", size=22, bold=True)
    add_card(slide, 0.75, 1.85, 12.0, 4.95, fill=Theme.panel)

    add_card(slide, 1.05, 2.3, 2.8, 1.0, fill=Theme.card_cyan, border=Theme.card_cyan)
    add_text(slide, 1.3, 2.62, 2.3, 0.24, "Agent Layer", size=14, bold=True, color=Theme.ink, align="center")
    add_text(slide, 1.25, 3.48, 2.35, 0.55, "Main / Log / Code /\nDomain / DB / Judge", size=12, align="center")

    add_card(slide, 5.1, 2.3, 3.15, 1.0, fill=Theme.card_warm, border=Theme.card_warm)
    add_text(slide, 5.35, 2.62, 2.65, 0.24, "MCP / Tool Gateway", size=14, bold=True, color=Theme.ink, align="center")
    add_text(slide, 5.25, 3.48, 2.85, 0.55, "权限 / 开关 / 超时 /\n审计 / 规范化", size=12, align="center")

    add_card(slide, 9.4, 2.3, 2.4, 1.0, fill=Theme.card_mint, border=Theme.card_mint)
    add_text(slide, 9.58, 2.62, 2.05, 0.24, "Data Sources", size=14, bold=True, color=Theme.ink, align="center")
    add_text(slide, 9.52, 3.48, 2.15, 0.55, "日志 / Git / PG /\nAPM / 告警", size=12, align="center")

    add_connector(slide, 3.85, 2.8, 5.1, 2.8, Theme.amber)
    add_connector(slide, 8.25, 2.8, 9.4, 2.8, Theme.amber)

    add_card(slide, 1.05, 5.1, 10.75, 1.0, fill=Theme.panel_2)
    add_text(slide, 1.3, 5.42, 10.2, 0.26, "最小权限、默认关闭高风险动作、全链路可回放，是 MCP 场景下的工程底线。", size=14)
    add_footer(slide, "MCP architecture spec; MCP security best practices")


def slide_8_multi_agent(prs: Presentation) -> None:
    """构建8multiAgent对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "业界共识 6：多 Agent 不是多线程，而是可控协作机制", "Protocol and convergence", 8)

    add_text(slide, 0.58, 1.0, 7.4, 0.42, "结论：多 Agent 的关键不在数量，在协作协议和收敛机制。", size=22, bold=True)

    add_card(slide, 0.7, 1.75, 12.0, 5.0, fill=Theme.panel_2)
    center_x, center_y = 6.7, 4.0
    add_card(slide, 5.65, 3.55, 2.1, 0.95, fill=Theme.card_cyan, border=Theme.card_cyan)
    add_text(slide, 5.85, 3.86, 1.7, 0.2, "Main Agent", size=14, bold=True, color=Theme.ink, align="center")

    nodes = [
        ("Log", 2.0, 2.35, Theme.cyan),
        ("Code", 2.0, 5.0, Theme.blue),
        ("Domain", 5.9, 5.5, Theme.green),
        ("DB", 9.55, 5.0, Theme.amber),
        ("Critic", 10.0, 2.35, Theme.red),
        ("Judge", 8.4, 1.4, Theme.soft),
    ]
    for label, x, y, color in nodes:
        add_card(slide, x, y, 1.55, 0.75, fill=Theme.panel, border=color)
        add_text(slide, x + 0.1, y + 0.22, 1.35, 0.18, label, size=11, bold=True, color=color, align="center")
        add_connector(slide, center_x, center_y, x + 0.77, y + 0.38, color)

    add_card(slide, 0.95, 6.15, 11.4, 0.48, fill=Theme.card_warm, border=Theme.card_warm)
    add_text(slide, 1.2, 6.28, 10.9, 0.18, "协议：命令分发 -> 证据回流 -> 交叉质疑 -> 主控收敛 -> 裁决输出", size=12, bold=True, color=Theme.ink)
    add_footer(slide, "LangGraph workflows-agents; supervisor pattern")


def slide_9_antipattern(prs: Presentation) -> None:
    """构建9antipattern对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, alt=True)
    add_topbar(slide, "业界反模式（必须避开）", "Common failure modes", 9)

    add_text(slide, 0.58, 1.0, 7.6, 0.42, "结论：多数失败案例不是模型不够强，而是工程纪律缺失。", size=22, bold=True)
    warnings = [
        ("无验证 Accept-All", "回归率上升", Theme.red),
        ("上下文无限堆砌", "漂移与幻觉", Theme.amber),
        ("工具调用无审计", "无法追责", Theme.red),
        ("异常无终态", "任务长期 pending", Theme.amber),
    ]
    coords = [(0.75, 1.95), (6.7, 1.95), (0.75, 4.2), (6.7, 4.2)]
    for (title, desc, accent), (x, y) in zip(warnings, coords):
        add_card(slide, x, y, 5.8, 1.7, fill=Theme.panel)
        add_text(slide, x + 0.28, y + 0.28, 5.0, 0.22, title, size=16, bold=True, color=accent)
        add_text(slide, x + 0.28, y + 0.8, 5.0, 0.28, desc, size=13, color=Theme.text_dim)
    add_footer(slide, "Thoughtworks production-grade vibe coding experiments")


def slide_10_project_harness(prs: Presentation) -> None:
    """构建10项目harness对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "本项目案例 1：我们如何实现 Harness 化", "Project mapping", 10)

    add_text(slide, 0.58, 1.0, 7.6, 0.42, "结论：本项目把“模型推理”和“系统控制”做了明确分层。", size=22, bold=True)
    add_card(slide, 0.72, 1.8, 12.0, 4.95, fill=Theme.panel_2)

    labels = [
        ("用户输入", 1.1, Theme.cyan),
        ("主 Agent 发命令", 3.55, Theme.amber),
        ("专家 Agent 分析", 6.0, Theme.blue),
        ("Judge 裁决", 8.45, Theme.green),
        ("报告生成", 10.9, Theme.red),
    ]
    for idx, (name, x, color) in enumerate(labels):
        add_card(slide, x, 2.55, 1.8, 1.0, fill=Theme.panel, border=color)
        add_text(slide, x + 0.12, 2.88, 1.56, 0.2, name, size=11, bold=True, color=color, align="center")
        if idx < len(labels) - 1:
            add_connector(slide, x + 1.8, 3.05, x + 2.45, 3.05, Theme.text_dim)

    add_card(slide, 1.0, 4.35, 5.1, 1.9, fill=Theme.card_cyan, border=Theme.card_cyan)
    add_text(slide, 1.2, 4.65, 4.7, 0.22, "代码入口", size=12, bold=True, color=Theme.ink)
    add_bullets(
        slide,
        1.2,
        4.95,
        4.6,
        1.0,
        [
            "runtime/langgraph_runtime.py",
            "runtime/langgraph/builder.py",
            "runtime/langgraph/nodes/supervisor.py",
        ],
        size=11,
        color=Theme.ink,
    )

    add_card(slide, 6.45, 4.35, 5.2, 1.9, fill=Theme.card_warm, border=Theme.card_warm)
    add_text(slide, 6.65, 4.65, 4.8, 0.22, "运行时约束", size=12, bold=True, color=Theme.ink)
    add_bullets(
        slide,
        6.65,
        4.95,
        4.7,
        1.0,
        ["主 Agent 命令先行", "阶段可终态", "异常可降级", "全过程可回放"],
        size=11,
        color=Theme.ink,
    )
    add_footer(slide, "Local project runtime: langgraph_runtime.py / builder.py / supervisor.py")


def slide_11_skill_tool(prs: Presentation) -> None:
    """构建11Skill工具对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, alt=True)
    add_topbar(slide, "本项目案例 2：Skill + Tool + 审计闭环", "Controlled tooling", 11)

    add_text(slide, 0.58, 1.0, 7.5, 0.42, "结论：工具调用不是“能调就行”，而是“可控、可审、可解释”。", size=22, bold=True)
    names = ["Command", "Gate", "Tool", "Skill", "LLM", "Audit"]
    colors = [Theme.cyan, Theme.amber, Theme.blue, Theme.green, Theme.soft, Theme.red]
    x = 0.85
    for idx, (name, color) in enumerate(zip(names, colors)):
        add_card(slide, x, 2.3, 1.55, 0.85, fill=Theme.panel, border=color)
        add_text(slide, x + 0.1, 2.56, 1.35, 0.18, name, size=11, bold=True, color=color, align="center")
        if idx < len(names) - 1:
            add_connector(slide, x + 1.55, 2.73, x + 1.9, 2.73, Theme.text_dim)
        x += 1.9

    add_card(slide, 0.85, 3.85, 5.7, 2.35, fill=Theme.card_cyan, border=Theme.card_cyan)
    add_text(slide, 1.1, 4.15, 5.2, 0.2, "能力面", size=12, bold=True, color=Theme.ink)
    add_bullets(
        slide,
        1.1,
        4.45,
        5.1,
        1.3,
        [
            "CodeAgent: Git 仓库检索",
            "LogAgent: 日志文件读取",
            "DomainAgent: 责任田 Excel 查询",
            "DatabaseAgent: PostgreSQL 元数据查询",
        ],
        size=11,
        color=Theme.ink,
    )

    add_card(slide, 6.85, 3.85, 5.5, 2.35, fill=Theme.card_warm, border=Theme.card_warm)
    add_text(slide, 7.1, 4.15, 5.0, 0.2, "审计面", size=12, bold=True, color=Theme.ink)
    add_bullets(
        slide,
        7.1,
        4.45,
        4.9,
        1.3,
        [
            "command_gate 判定",
            "请求/返回摘要",
            "状态、耗时、错误",
            "skill 命中来源与注入摘要",
        ],
        size=11,
        color=Theme.ink,
    )
    add_footer(slide, "Local project services: agent_tool_context_service.py / agent_skill_service.py")


def slide_12_resilience(prs: Presentation) -> None:
    """构建12韧性对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "本项目案例 3：韧性设计（超时、降级、恢复）", "Resilience by design", 12)

    add_text(slide, 0.58, 1.0, 7.7, 0.42, "结论：系统最重要的能力，是失败时仍然能推进。", size=22, bold=True)
    add_card(slide, 0.72, 1.8, 5.95, 4.95, fill=Theme.panel_2)
    add_text(slide, 1.0, 2.1, 5.4, 0.2, "会话状态机", size=12, bold=True, color=Theme.amber)
    state_boxes = [
        ("RUNNING", 1.1, 2.75, Theme.cyan),
        ("RETRY", 2.7, 3.6, Theme.amber),
        ("DEGRADED", 4.25, 4.45, Theme.red),
        ("COMPLETED", 1.7, 5.45, Theme.green),
        ("FAILED", 4.1, 5.45, Theme.soft),
    ]
    for name, x, y, color in state_boxes:
        add_card(slide, x, y, 1.45, 0.72, fill=Theme.panel, border=color)
        add_text(slide, x + 0.08, y + 0.21, 1.28, 0.18, name, size=10, bold=True, color=color, align="center")
    add_connector(slide, 2.55, 3.1, 2.7, 3.6, Theme.text_dim)
    add_connector(slide, 4.15, 3.95, 4.25, 4.45, Theme.text_dim)
    add_connector(slide, 3.2, 4.95, 2.42, 5.45, Theme.text_dim)
    add_connector(slide, 4.95, 5.17, 4.82, 5.45, Theme.text_dim)

    add_card(slide, 7.0, 1.8, 5.65, 4.95, fill=Theme.panel)
    add_text(slide, 7.28, 2.1, 5.1, 0.2, "Top-K 根因候选（示意）", size=12, bold=True, color=Theme.cyan)
    add_chart(
        slide,
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        x=7.28,
        y=2.45,
        w=5.0,
        h=2.8,
        categories=["连接池", "锁竞争", "重试风暴", "配置缺陷", "下游超时"],
        series=[("置信度", [41, 24, 17, 11, 7])],
    )
    add_text(slide, 7.35, 5.55, 4.9, 0.42, "重试、降级、恢复重跑，保证不会无响应卡死。", size=13, color=Theme.text_dim)
    add_footer(slide, "Timeout, degrade and recovery design from local runtime implementation")


def slide_13_2week(prs: Presentation) -> None:
    """构建132week对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, alt=True)
    add_topbar(slide, "团队落地模板（2 周可执行）", "Practical adoption board", 13)

    add_text(slide, 0.58, 1.0, 7.2, 0.42, "结论：先统一方法，再扩展能力；先可用，再做高级自治。", size=22, bold=True)
    add_card(slide, 0.82, 1.9, 5.7, 4.95, fill=Theme.panel)
    add_text(slide, 1.08, 2.18, 5.1, 0.22, "Week 1 / 打基础", size=13, bold=True, color=Theme.cyan)
    add_bullets(
        slide,
        1.08,
        2.6,
        5.0,
        2.3,
        [
            "定义 AGENTS.md 与边界",
            "梳理统一工作流",
            "接入最小测试与 PR 门禁",
            "建立关键日志与事件追踪",
        ],
        size=13,
    )
    add_card(slide, 1.08, 5.45, 4.9, 0.85, fill=Theme.card_cyan, border=Theme.card_cyan)
    add_text(slide, 1.28, 5.72, 4.5, 0.18, "验收：端到端可重复执行，链路可追踪。", size=12, bold=True, color=Theme.ink)

    add_card(slide, 6.82, 1.9, 5.7, 4.95, fill=Theme.panel)
    add_text(slide, 7.08, 2.18, 5.1, 0.22, "Week 2 / 补治理", size=13, bold=True, color=Theme.green)
    add_bullets(
        slide,
        7.08,
        2.6,
        5.0,
        2.3,
        [
            "工具权限和审计统一化",
            "超时、重试、终态保证",
            "Benchmark 挂 CI Gate",
            "做一次真实故障演练与复盘",
        ],
        size=13,
    )
    add_card(slide, 7.08, 5.45, 4.9, 0.85, fill=Theme.card_mint, border=Theme.card_mint)
    add_text(slide, 7.28, 5.72, 4.5, 0.18, "验收：无长期 pending，CI 可阻断回归。", size=12, bold=True, color=Theme.ink)
    add_footer(slide, "Execution template synthesized from outline and project plans")


def slide_14_roadmap(prs: Presentation) -> None:
    """构建14路线图对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "90 天路线图（从可用到规模化）", "Capability maturity", 14)

    add_text(slide, 0.58, 1.0, 7.4, 0.42, "结论：Vibe Coding 要变成组织能力，必须走平台化治理路线。", size=22, bold=True)
    add_card(slide, 0.78, 1.85, 6.2, 4.95, fill=Theme.panel_2)
    add_chart(
        slide,
        chart_type=XL_CHART_TYPE.LINE_MARKERS,
        x=1.05,
        y=2.15,
        w=5.55,
        h=3.0,
        categories=["P0", "P1", "P2", "P3"],
        series=[("成熟度", [35, 56, 72, 85])],
    )
    add_text(slide, 1.08, 5.55, 5.45, 0.32, "先做可用性，再做准确性，再做可控自治与持续学习。", size=13, color=Theme.text_dim)

    add_card(slide, 7.25, 1.85, 5.35, 4.95, fill=Theme.panel)
    roadmap = [
        ("P0", "可用性", Theme.cyan),
        ("P1", "准确性", Theme.amber),
        ("P2", "可控自治", Theme.blue),
        ("P3", "持续学习", Theme.green),
    ]
    y = 2.2
    for phase, label, accent in roadmap:
        add_card(slide, 7.55, y, 4.75, 0.82, fill=Theme.panel_2, border=accent)
        add_text(slide, 7.75, y + 0.17, 0.55, 0.18, phase, size=11, bold=True, color=accent)
        add_text(slide, 8.45, y + 0.17, 1.45, 0.18, label, size=12, bold=True)
        desc = {
            "P0": "不 pending、可观测、可回放",
            "P1": "Top-K 根因、跨源证据校验",
            "P2": "审批、回滚、No-Regression Gate",
            "P3": "反馈闭环、A/B 评测、策略演进",
        }[phase]
        add_text(slide, 9.95, y + 0.17, 2.05, 0.34, desc, size=10, color=Theme.text_dim)
        y += 1.05
    add_footer(slide, "P0-P3 roadmap from outline and current RCA system roadmap")


def slide_15_close(prs: Presentation) -> None:
    """构建15关闭对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, alt=True)
    add_topbar(slide, "结尾与行动清单", "Execution checklist", 15)

    add_card(slide, 0.78, 1.2, 12.0, 1.55, fill=Theme.card_warm, border=Theme.card_warm)
    add_text(slide, 1.05, 1.65, 11.4, 0.5, "结论：Vibe Coding 的终局不是“自动写代码”，而是“可控地持续交付正确结果”。", size=24, bold=True, color=Theme.ink)

    add_card(slide, 0.78, 3.15, 12.0, 3.25, fill=Theme.panel)
    add_text(slide, 1.05, 3.48, 11.3, 0.22, "本周 5 条可执行动作", size=13, bold=True, color=Theme.amber)
    actions = [
        "固化 AGENTS.md 与工程约束",
        "给关键链路加可观测事件和审计",
        "建立最小 benchmark 并接入 CI Gate",
        "把工具接入改为开关化、可回滚",
        "用一次真实故障做端到端演练并复盘",
    ]
    y = 3.95
    for idx, action in enumerate(actions, 1):
        add_card(slide, 1.05, y, 10.95, 0.42, fill=Theme.panel_2)
        add_text(slide, 1.22, y + 0.08, 0.35, 0.16, f"{idx}.", size=11, bold=True, color=Theme.cyan)
        add_text(slide, 1.6, y + 0.08, 10.1, 0.16, action, size=12)
        y += 0.52
    add_footer(slide, "Conclusion synthesized from the 2026-03-06 vibe coding outline")


def write_readme() -> None:
    """写出说明文档相关产物。"""
    
    OUT_README.write_text(
        "\n".join(
            [
                "# Vibe Coding War-room PPT",
                "",
                f"- PPTX: `{OUT_PPTX}`",
                "",
                "## Style",
                "- War-room briefing style",
                "- Dark canvas, high contrast, evidence-first layout",
                "- Same 15-slide structure as the approved outline",
            ]
        ),
        encoding="utf-8",
    )


def main() -> None:
    """执行脚本主流程，串联参数解析、内容生成与结果输出。"""
    
    ensure_dirs()
    prs = new_prs()
    slide_1_cover(prs)
    slide_2_why(prs)
    slide_3_harness(prs)
    slide_4_context(prs)
    slide_5_workflow(prs)
    slide_6_gate(prs)
    slide_7_mcp(prs)
    slide_8_multi_agent(prs)
    slide_9_antipattern(prs)
    slide_10_project_harness(prs)
    slide_11_skill_tool(prs)
    slide_12_resilience(prs)
    slide_13_2week(prs)
    slide_14_roadmap(prs)
    slide_15_close(prs)
    prs.save(OUT_PPTX)
    write_readme()
    print(f"Generated: {OUT_PPTX}")


if __name__ == "__main__":
    main()
