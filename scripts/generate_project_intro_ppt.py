#!/usr/bin/env python3
"""generate项目介绍PPT脚本。"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(
    "/Users/neochen/multi-agent-cli_v2/plans/2026-03-04-生产问题根因分析系统-项目介绍.pptx"
)


COLORS = {
    "bg": RGBColor(245, 248, 252),
    "header": RGBColor(20, 43, 76),
    "header_text": RGBColor(255, 255, 255),
    "title": RGBColor(20, 43, 76),
    "body": RGBColor(43, 43, 43),
    "accent": RGBColor(16, 100, 177),
    "card": RGBColor(255, 255, 255),
    "card_border": RGBColor(210, 223, 238),
}


def _set_background(slide, prs: Presentation) -> None:
    """执行setbackground相关逻辑。"""
    
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["bg"]
    shape.line.fill.background()


def _add_header(slide, title: str, subtitle: str = "") -> None:
    """执行add页眉相关逻辑。"""
    
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.95)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["header"]
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(8.9), Inches(0.5))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["header_text"]

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(9.1), Inches(0.25), Inches(3.9), Inches(0.4))
        tf = sub_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(223, 236, 252)
        p.alignment = PP_ALIGN.RIGHT


def _add_card(slide, left: float, top: float, width: float, height: float):
    """执行add卡片相关逻辑。"""
    
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["card"]
    card.line.color.rgb = COLORS["card_border"]
    return card


def _set_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: Iterable[str],
    *,
    font_size: int = 18,
    bold_first: bool = False,
    color: RGBColor | None = None,
    bullet: bool = False,
) -> None:
    """执行settextbox相关逻辑。"""
    
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size if idx > 0 or not bold_first else max(font_size, 22))
        p.font.bold = bool(bold_first and idx == 0)
        p.font.color.rgb = color or COLORS["body"]
        if bullet and idx > 0:
            p.level = 0
            p.text = f"• {line}"


def _slide_cover(prs: Presentation) -> None:
    """执行slide封面相关逻辑。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, prs)
    _add_header(slide, "生产问题根因分析系统", "LangGraph Multi-Agent")

    _set_textbox(
        slide,
        0.9,
        1.6,
        11.8,
        1.3,
        ["项目介绍", "面向生产故障的多 Agent 协同分析平台"],
        font_size=36,
        bold_first=True,
        color=COLORS["title"],
    )
    _set_textbox(
        slide,
        0.95,
        3.1,
        11.8,
        1.6,
        [
            "技术栈：FastAPI + LangGraph + React 18 + Ant Design",
            "核心能力：主 Agent 指挥 + 专家 Agent 协作 + 工具调用审计 + 结构化报告",
            "当前模型：kimi-k2.5（OpenAI 兼容接口）",
        ],
        font_size=20,
        color=COLORS["body"],
    )

    _set_textbox(
        slide,
        0.95,
        6.6,
        11.8,
        0.5,
        [f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}（Asia/Shanghai）"],
        font_size=14,
        color=RGBColor(90, 100, 115),
    )


def _slide_overview(prs: Presentation) -> None:
    """执行slideoverview相关逻辑。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, prs)
    _add_header(slide, "1. 项目定位与目标", "Why this system")

    _add_card(slide, 0.5, 1.25, 6.3, 5.8)
    _set_textbox(
        slide,
        0.85,
        1.55,
        5.65,
        5.2,
        [
            "系统定位",
            "从“人工排障”升级到“智能调查 + 证据化结论”。",
            "统一处理日志、代码、领域资产、监控现象，面向生产故障快速定位根因。",
            "",
            "业务价值",
            "缩短故障定位时间（MTTR）",
            "提升跨团队协作效率",
            "保留可复盘的决策链路",
        ],
        font_size=19,
        bold_first=True,
    )

    _add_card(slide, 6.95, 1.25, 5.9, 5.8)
    _set_textbox(
        slide,
        7.3,
        1.55,
        5.2,
        5.2,
        [
            "当前目标",
            "1. 主 Agent 统一调度专家 Agent",
            "2. 工具调用必须可审计、可回放",
            "3. 前端实时展示分析过程与证据链",
            "4. 输出 Top-K 根因候选 + 修复建议",
            "",
            "约束",
            "暂不依赖外部数据库，使用本地文件/内存存储。",
        ],
        font_size=18,
        bold_first=True,
    )


def _slide_architecture(prs: Presentation) -> None:
    """执行slide架构相关逻辑。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, prs)
    _add_header(slide, "2. 系统总体架构", "Frontend / Backend / Runtime")

    boxes = [
        ("用户", 0.35, 2.3, 1.5, 1.0, RGBColor(255, 250, 232)),
        ("前端\nReact + AntD", 2.2, 2.0, 2.3, 1.7, RGBColor(235, 246, 255)),
        ("后端 API\nFastAPI", 4.95, 2.0, 1.9, 1.7, RGBColor(229, 240, 255)),
        ("LangGraph Runtime\nOrchestrator", 7.15, 2.0, 2.65, 1.7, RGBColor(227, 238, 255)),
        ("专家 Agent 群\n(并行+协作)", 10.05, 1.6, 2.9, 1.4, RGBColor(229, 246, 240)),
        ("工具层\nGit/Log/Excel/Telemetry/CMDB", 10.05, 3.25, 2.9, 1.4, RGBColor(238, 252, 246)),
        ("本地存储\nSession/Report/Lineage", 7.15, 4.25, 2.65, 1.3, RGBColor(249, 244, 255)),
        ("WebSocket 事件流", 2.2, 4.4, 4.65, 1.15, RGBColor(245, 250, 255)),
    ]
    for text, left, top, width, height, color in boxes:
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.color.rgb = COLORS["card_border"]
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["title"]

    arrows = [
        (1.9, 2.65, 0.25, 0.25),
        (4.55, 2.65, 0.25, 0.25),
        (6.92, 2.65, 0.2, 0.25),
        (9.82, 2.3, 0.2, 0.25),
    ]
    for left, top, width, height in arrows:
        arr = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        arr.fill.solid()
        arr.fill.fore_color.rgb = COLORS["accent"]
        arr.line.fill.background()


def _slide_agents(prs: Presentation) -> None:
    """执行slideAgent相关逻辑。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, prs)
    _add_header(slide, "3. 多 Agent 协同机制", "Commander + Specialist")

    _add_card(slide, 0.5, 1.25, 12.3, 1.25)
    _set_textbox(
        slide,
        0.85,
        1.55,
        11.6,
        0.8,
        [
            "主 Agent（ProblemAnalysisAgent）先做问题拆解，再按命令分发任务；子 Agent 依据命令决定是否调用工具并回传结构化证据。"
        ],
        font_size=17,
        color=COLORS["title"],
    )

    cards = [
        ("LogAgent", "日志模式识别、错误链路与时序聚类"),
        ("DomainAgent", "接口/领域/聚合根/责任田映射"),
        ("CodeAgent", "代码路径定位、变更影响评估"),
        ("MetricsAgent", "资源瓶颈、指标异常关联"),
        ("ChangeAgent", "发布变更窗口关联分析"),
        ("RunbookAgent", "案例库检索与处置 SOP 建议"),
        ("CriticAgent", "反证审查，指出证据缺口"),
        ("RebuttalAgent", "针对质疑补强证据"),
        ("JudgeAgent", "综合裁决与置信度打分"),
        ("VerificationAgent", "验证计划与回归检查项"),
    ]
    left = 0.5
    top = 2.75
    width = 3.9
    height = 0.8
    for idx, (name, desc) in enumerate(cards):
        row = idx // 3
        col = idx % 3
        x = left + col * 4.1
        y = top + row * 0.9
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height)
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shp.line.color.rgb = RGBColor(207, 221, 238)
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{name}：{desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS["body"]


def _slide_flow(prs: Presentation) -> None:
    """执行slideflow相关逻辑。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, prs)
    _add_header(slide, "4. 端到端分析流程", "From Incident to Report")

    steps = [
        ("01", "创建 Incident", "输入报错日志、堆栈、现象与监控信息"),
        ("02", "资产映射", "定位接口 -> 领域 -> 聚合根 -> Owner"),
        ("03", "主 Agent 指挥", "拆解问题并向专家 Agent 下发命令"),
        ("04", "专家协作分析", "按需调用工具，回传证据与中间结论"),
        ("05", "质疑与反驳", "Critic/Rebuttal 循环收敛观点"),
        ("06", "裁决与报告", "Judge + Verification 输出 Top-K 根因与行动项"),
    ]
    for idx, (num, title, desc) in enumerate(steps):
        y = 1.3 + idx * 0.95
        badge = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.65), Inches(y), Inches(0.55), Inches(0.55)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLORS["accent"]
        badge.line.fill.background()
        tf = badge.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = num
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)

        _set_textbox(
            slide,
            1.35,
            y - 0.02,
            11.1,
            0.65,
            [f"{title}：{desc}"],
            font_size=16,
            color=COLORS["body"],
        )


def _slide_tools(prs: Presentation) -> None:
    """执行slide工具相关逻辑。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, prs)
    _add_header(slide, "5. 工具调用与审计机制", "Command-gated Tooling")

    _add_card(slide, 0.5, 1.25, 6.1, 5.8)
    _set_textbox(
        slide,
        0.85,
        1.55,
        5.45,
        5.2,
        [
            "当前工具能力",
            "CodeAgent：Git 仓库检索（本地/远程）",
            "LogAgent：本地日志文件读取",
            "DomainAgent：Excel/CSV 责任田查询",
            "Metrics/Change：Telemetry/CMDB 连接器入口",
            "",
            "调用约束",
            "主 Agent 先发命令后才允许调用",
            "工具配置支持开关控制",
            "未配置工具的 Agent 不展示调用记录",
        ],
        font_size=16,
        bold_first=True,
    )

    _add_card(slide, 6.9, 1.25, 5.95, 5.8)
    _set_textbox(
        slide,
        7.25,
        1.55,
        5.3,
        5.2,
        [
            "审计数据",
            "每次调用记录 command_gate 决策",
            "记录 I/O 轨迹（文件读取、Git 命令、HTTP 访问）",
            "记录状态、摘要、错误原因、截断引用ID",
            "",
            "前端可视化",
            "辩论过程可查看工具调用消息",
            "战情页同屏展示时间线、证据链、工具审计",
            "调查复盘台支持决策回放与报告比对",
        ],
        font_size=16,
        bold_first=True,
    )


def _slide_frontend(prs: Presentation) -> None:
    """执行slide前端相关逻辑。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, prs)
    _add_header(slide, "6. 前端体验与页面结构", "Product Surface")

    modules = [
        ("首页 /", "概览统计、快速创建与启动分析、Agent 角色说明"),
        ("故障分析 /incident", "资产映射 / 辩论过程 / 辩论结果 三标签联动"),
        ("战情页 /war-room", "实时态势：时间线 + 证据链 + 工具调用 + 结论"),
        ("调查复盘台 /workbench", "会话回放、关键决策路径、报告版本对比"),
        ("工具中心 /tools", "工具源管理、详情查看、参数试运行"),
        ("治理中心 /governance", "超时率、失败率、团队成功率与成本指标"),
        ("评测中心 /benchmark", "Top1/超时率/空结论率趋势与回归基线"),
    ]

    for idx, (title, desc) in enumerate(modules):
        y = 1.25 + idx * 0.8
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(y), Inches(12.0), Inches(0.62)
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shp.line.color.rgb = RGBColor(211, 224, 238)
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{title}：{desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["body"]


def _slide_data_model(prs: Presentation) -> None:
    """执行slidedatamodel相关逻辑。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, prs)
    _add_header(slide, "7. 数据与状态管理", "Local-first")

    _add_card(slide, 0.5, 1.3, 6.15, 5.7)
    _set_textbox(
        slide,
        0.82,
        1.55,
        5.5,
        5.1,
        [
            "核心状态对象",
            "Incident：故障输入与上下文",
            "DebateSession：执行状态机",
            "DebateRound：每轮 Agent 输出",
            "DebateResult：最终结论、Top-K、证据链",
            "LineageRecord：过程事件与工具审计",
            "",
            "存储策略",
            "LOCAL_STORE_BACKEND = file / memory",
            "默认文件持久化，便于复盘与导出",
        ],
        font_size=16,
        bold_first=True,
    )

    _add_card(slide, 6.95, 1.3, 5.9, 5.7)
    _set_textbox(
        slide,
        7.27,
        1.55,
        5.25,
        5.1,
        [
            "稳定性机制",
            "DebateStatus 状态迁移校验",
            "LLM 超时与重试策略（按阶段区分）",
            "长会话 compaction + prune",
            "输出截断与引用机制，防止上下文污染",
            "",
            "回放能力",
            "lineage replay：关键决策路径",
            "report compare：同 incident 版本对比",
        ],
        font_size=16,
        bold_first=True,
    )


def _slide_api(prs: Presentation) -> None:
    """执行slideAPI相关逻辑。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, prs)
    _add_header(slide, "8. 对外接口与部署", "API / WS / Startup")

    _add_card(slide, 0.5, 1.25, 7.75, 5.8)
    _set_textbox(
        slide,
        0.82,
        1.55,
        7.1,
        5.2,
        [
            "核心 API（/api/v1）",
            "Incident：创建、列表、详情",
            "Debate：创建会话、执行、状态、结果、取消",
            "Report：查询、重生成、版本对比",
            "Assets：接口与责任田定位、资产融合",
            "Settings/Tools：工具配置、审计、试运行",
            "",
            "WebSocket",
            "ws://host/ws/debates/{session_id}?auto_start=true",
            "实时推送 phase/agent_chat/tool_io/session 状态事件",
        ],
        font_size=15,
        bold_first=True,
    )
    _add_card(slide, 8.45, 1.25, 4.4, 5.8)
    _set_textbox(
        slide,
        8.75,
        1.55,
        3.75,
        5.2,
        [
            "部署与运行",
            "后端：FastAPI + Uvicorn",
            "前端：Vite Dev Server",
            "一键启动：npm run start:all",
            "一键停止：npm run stop:all",
            "",
            "关键配置",
            "LLM_BASE_URL",
            "LLM_API_KEY",
            "LLM_MODEL=kimi-k2.5",
            "DEBATE_MAX_ROUNDS",
        ],
        font_size=15,
        bold_first=True,
    )


def _slide_metrics(prs: Presentation) -> None:
    """执行slidemetrics相关逻辑。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, prs)
    _add_header(slide, "9. 质量保障与运营指标", "Benchmark + Governance")

    _add_card(slide, 0.5, 1.25, 6.1, 5.8)
    _set_textbox(
        slide,
        0.82,
        1.55,
        5.45,
        5.2,
        [
            "Benchmark Center",
            "Top1 命中率",
            "平均重叠分",
            "超时率",
            "空结论率",
            "",
            "支持手工运行与基线文件追踪",
            "可接入 CI 作为 Benchmark Gate",
        ],
        font_size=16,
        bold_first=True,
    )

    _add_card(slide, 6.9, 1.25, 5.95, 5.8)
    _set_textbox(
        slide,
        7.22,
        1.55,
        5.3,
        5.2,
        [
            "Governance Center",
            "团队级分析成功率与失败率",
            "Agent 调用耗时与超时分布",
            "工具调用成功率与错误类型",
            "会话可追踪审计链",
            "",
            "目标",
            "避免长期 pending",
            "降低无结论场景占比",
        ],
        font_size=16,
        bold_first=True,
    )


def _slide_roadmap(prs: Presentation) -> None:
    """执行slide路线图相关逻辑。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, prs)
    _add_header(slide, "10. 后续演进路线", "Near-term Roadmap")

    _set_textbox(
        slide,
        0.8,
        1.35,
        11.8,
        0.6,
        ["基于当前代码与计划文档，下一阶段优先级如下："],
        font_size=18,
        color=COLORS["title"],
    )

    items = [
        ("P0 可用性", "自动调查入口、超时切换、局部重试、避免 pending"),
        ("P1 准确率", "跨源证据约束、依赖拓扑推理、Top-K 置信区间"),
        ("P2 可控修复", "修复状态机、No-Regression Gate、自动回滚"),
        ("P3 平台治理", "A/B 策略、RBAC、多租户隔离、成本预算"),
    ]
    for idx, (title, desc) in enumerate(items):
        y = 2.05 + idx * 1.2
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(11.8), Inches(0.92)
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shp.line.color.rgb = RGBColor(209, 221, 238)
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{title}：{desc}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS["body"]


def _slide_summary(prs: Presentation) -> None:
    """执行slide摘要相关逻辑。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, prs)
    _add_header(slide, "11. 总结", "What is delivered now")
    _set_textbox(
        slide,
        0.9,
        1.6,
        11.6,
        4.6,
        [
            "当前系统已具备生产问题分析闭环：",
            "1) 主 Agent + 专家 Agent 协同分析",
            "2) 命令驱动工具调用与全链路审计",
            "3) 前端可视化展示资产映射、辩论过程、辩论结果",
            "4) 支持报告、回放、评测、治理等平台能力",
            "",
            "面向下一阶段，将继续提升：可靠性、准确率、自治修复能力。"
        ],
        font_size=20,
        color=COLORS["title"],
    )

    _set_textbox(
        slide,
        0.9,
        6.4,
        11.6,
        0.5,
        ["附：源代码与设计文档均已按当前实现更新，可直接用于内部汇报。"],
        font_size=14,
        color=RGBColor(95, 105, 120),
    )


def build_presentation(output_path: Path) -> Path:
    """构建presentation相关产物或页面内容。"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _slide_cover(prs)
    _slide_overview(prs)
    _slide_architecture(prs)
    _slide_agents(prs)
    _slide_flow(prs)
    _slide_tools(prs)
    _slide_frontend(prs)
    _slide_data_model(prs)
    _slide_api(prs)
    _slide_metrics(prs)
    _slide_roadmap(prs)
    _slide_summary(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def main() -> None:
    """执行脚本主流程，串联参数解析、内容生成与结果输出。"""
    
    result = build_presentation(OUTPUT_PATH)
    print(f"Generated: {result}")


if __name__ == "__main__":
    main()
