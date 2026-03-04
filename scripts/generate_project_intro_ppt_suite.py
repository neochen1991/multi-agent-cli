#!/usr/bin/env python3
"""
Generate two polished decks:
1) Executive concise version (8 slides)
2) Technical full version (18 slides)
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Iterable, List, Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_EXEC = Path(
    "/Users/neochen/multi-agent-cli_v2/plans/2026-03-04-生产问题根因分析系统-管理层汇报版.pptx"
)
OUT_TECH = Path(
    "/Users/neochen/multi-agent-cli_v2/plans/2026-03-04-生产问题根因分析系统-技术详版.pptx"
)


class C:
    BG = RGBColor(246, 250, 255)
    NAVY = RGBColor(10, 34, 70)
    BLUE = RGBColor(34, 104, 201)
    CYAN = RGBColor(20, 148, 170)
    INDIGO = RGBColor(73, 89, 173)
    TEXT = RGBColor(31, 41, 55)
    SUB = RGBColor(95, 109, 129)
    WHITE = RGBColor(255, 255, 255)
    BORDER = RGBColor(212, 224, 239)
    L_BLUE = RGBColor(236, 244, 255)
    L_CYAN = RGBColor(233, 250, 253)
    L_PURP = RGBColor(241, 238, 255)
    L_ORANGE = RGBColor(255, 245, 234)
    GREEN = RGBColor(15, 157, 101)
    RED = RGBColor(208, 57, 57)
    YELLOW = RGBColor(226, 149, 36)


def mk_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def bg(slide, prs: Presentation) -> None:
    r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = C.BG
    r.line.fill.background()


def header(slide, title: str, sub: str, page: int) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.86)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C.NAVY
    bar.line.fill.background()
    txt(slide, 0.45, 0.16, 8.5, 0.45, title, 23, True, C.WHITE)
    txt(slide, 8.85, 0.22, 3.2, 0.28, sub, 11, False, RGBColor(205, 222, 247), "right")
    txt(slide, 12.25, 0.22, 0.6, 0.28, str(page), 11, True, RGBColor(178, 201, 236), "right")


def card(slide, x: float, y: float, w: float, h: float, fill: RGBColor = C.WHITE):
    s = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = C.BORDER
    return s


def txt(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = C.TEXT,
    align: str = "left",
) -> None:
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT


def bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: Sequence[str],
    size: int = 12,
    color: RGBColor = C.TEXT,
) -> None:
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(size)
        p.font.color.rgb = color


def arr(slide, x: float, y: float, w: float, h: float, color: RGBColor = C.BLUE) -> None:
    a = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()


def left_arr(slide, x: float, y: float, w: float, h: float, color: RGBColor = C.BLUE) -> None:
    a = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LEFT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()


def cover(prs: Presentation, title: str, subtitle: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    band = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.7)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = C.NAVY
    band.line.fill.background()
    txt(s, 0.78, 0.35, 10.5, 0.65, title, 42, True, C.WHITE)
    txt(s, 0.78, 1.04, 9.2, 0.36, subtitle, 16, False, RGBColor(204, 221, 248))
    card(s, 0.78, 2.15, 5.95, 3.95, C.WHITE)
    txt(s, 1.05, 2.5, 5.3, 0.45, "汇报重点", 22, True, C.INDIGO)
    bullets(
        s,
        1.05,
        2.95,
        5.3,
        2.8,
        [
            "系统能力全景与业务价值",
            "关键架构与端到端流程",
            "多 Agent 协作与工具审计机制",
            "结果质量与治理闭环",
            "下一阶段改进路线",
        ],
        15,
    )
    card(s, 7.0, 2.15, 5.55, 3.95, C.L_BLUE)
    txt(s, 7.3, 2.5, 4.9, 0.45, "技术基线", 22, True, C.INDIGO)
    bullets(
        s,
        7.3,
        2.95,
        4.9,
        2.8,
        [
            "FastAPI + LangGraph + React",
            "LLM: kimi-k2.5（OpenAI兼容）",
            "local file/memory 存储",
            "WebSocket 实时事件流",
            "Benchmark + Governance 治理",
        ],
        15,
    )
    txt(
        s,
        0.78,
        6.65,
        11.8,
        0.3,
        f"{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}（Asia/Shanghai）",
        11,
        False,
        C.SUB,
    )


def agenda(prs: Presentation, page: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    header(s, "目录", "Roadmap", page)
    items = [
        "1. 目标与价值",
        "2. 系统总体架构图",
        "3. LangGraph 运行时架构",
        "4. 多 Agent 协作网络",
        "5. 工具调用与审计链路",
        "6. 端到端流程图",
        "7. 状态机与数据模型",
        "8. 前端体验与页面结构",
        "9. 质量保障与治理",
        "10. 路线图与计划",
    ]
    for i, it in enumerate(items):
        x = 0.9 if i < 5 else 6.95
        y = 1.35 + (i % 5) * 1.06
        card(s, x, y, 5.35, 0.84, C.WHITE)
        txt(s, x + 0.22, y + 0.23, 4.9, 0.36, it, 14)


def business_value(prs: Presentation, page: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    header(s, "1) 目标与业务价值", "Business Value", page)
    card(s, 0.62, 1.2, 3.95, 5.95, C.L_ORANGE)
    txt(s, 0.92, 1.5, 3.35, 0.42, "典型痛点", 20, True, C.RED)
    bullets(
        s,
        0.92,
        2.0,
        3.35,
        4.9,
        [
            "故障证据分散，排查链路断裂",
            "跨团队沟通成本高，责任难界定",
            "单次结论可解释性不足",
            "缺少可回放的过程记录",
            "高峰期存在 pending 风险",
        ],
        13,
    )
    card(s, 4.82, 1.2, 7.85, 2.9, C.WHITE)
    txt(s, 5.12, 1.5, 7.2, 0.42, "平台目标", 18, True, C.INDIGO)
    bullets(
        s,
        5.12,
        1.95,
        7.2,
        1.9,
        [
            "主 Agent 指挥多专家 Agent 协作，实现结构化根因分析",
            "统一打通日志、代码、资产、指标证据",
            "结果可复核：Top-K 根因、证据链、置信度、验证计划",
        ],
        13,
    )
    card(s, 4.82, 4.25, 7.85, 2.9, C.L_BLUE)
    txt(s, 5.12, 4.55, 7.2, 0.42, "价值产出", 18, True, C.INDIGO)
    bullets(
        s,
        5.12,
        5.0,
        7.2,
        1.9,
        [
            "缩短 MTTR，降低人工排查耗时",
            "提升分析正确率与跨团队协同效率",
            "沉淀可审计、可回放、可治理的组织知识资产",
        ],
        13,
    )


def architecture_e2e(prs: Presentation, page: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    header(s, "2) 系统总体架构图", "End-to-end", page)
    # Top lane: user -> frontend -> backend -> runtime
    card(s, 0.45, 1.2, 1.85, 0.9, C.L_ORANGE)
    txt(s, 0.56, 1.5, 1.62, 0.3, "用户/值班SRE", 11, True, align="center")
    arr(s, 2.35, 1.48, 0.25, 0.2)
    card(s, 2.65, 1.2, 2.25, 0.9, C.L_BLUE)
    txt(s, 2.83, 1.43, 1.9, 0.38, "Frontend\nReact + WS", 11, True, align="center")
    arr(s, 4.95, 1.48, 0.25, 0.2)
    card(s, 5.25, 1.2, 2.25, 0.9, C.L_BLUE)
    txt(s, 5.45, 1.42, 1.85, 0.36, "FastAPI API层\nREST + WebSocket", 10, True, align="center")
    arr(s, 7.55, 1.48, 0.25, 0.2)
    card(s, 7.85, 1.2, 5.0, 0.9, C.L_PURP)
    txt(s, 8.08, 1.5, 4.55, 0.3, "DebateService -> create_ai_debate_orchestrator -> LangGraphRuntime", 10, True, align="center")

    # Middle lane
    card(s, 0.45, 2.35, 4.8, 2.15, C.WHITE)
    txt(s, 0.72, 2.62, 4.2, 0.3, "交互协议层（真实接口）", 12, True, C.INDIGO)
    bullets(
        s,
        0.72,
        2.95,
        4.25,
        1.45,
        [
            "POST /api/v1/incidents, /debates, /reports",
            "WS /ws/debates/{session_id}?auto_start=true",
            "snapshot/event/result 三类实时消息",
        ],
        11,
    )
    card(s, 5.45, 2.35, 3.55, 2.15, C.WHITE)
    txt(s, 5.72, 2.62, 3.0, 0.3, "运行时核心", 12, True, C.INDIGO)
    bullets(
        s,
        5.72,
        2.95,
        3.0,
        1.45,
        [
            "GraphBuilder 动态建图",
            "Supervisor 动态路由",
            "PhaseExecutor 并行分析",
            "AgentRunner 统一调用LLM",
        ],
        11,
    )
    card(s, 9.2, 2.35, 3.65, 2.15, C.WHITE)
    txt(s, 9.46, 2.62, 3.15, 0.3, "Agent 协作层", 12, True, C.INDIGO)
    bullets(
        s,
        9.46,
        2.95,
        3.15,
        1.45,
        [
            "ProblemAnalysisAgent（主控）",
            "Analysis: Log/Domain/Code/Metrics/Change/Runbook/Rule",
            "Critic/Rebuttal/Judge/Verification",
        ],
        10,
    )

    # Bottom lane
    card(s, 0.45, 4.75, 4.3, 2.0, C.L_CYAN)
    txt(s, 0.72, 5.0, 3.8, 0.3, "工具与连接器", 12, True, C.INDIGO)
    bullets(
        s,
        0.72,
        5.27,
        3.85,
        1.35,
        [
            "CodeAgent: Git repo search",
            "LogAgent: local log reader",
            "DomainAgent: Excel/CSV mapping",
            "Telemetry/CMDB connector（可配置）",
        ],
        11,
    )
    card(s, 4.95, 4.75, 4.0, 2.0, C.L_BLUE)
    txt(s, 5.2, 5.0, 3.5, 0.3, "会话与审计存储", 12, True, C.INDIGO)
    bullets(
        s,
        5.2,
        5.27,
        3.5,
        1.35,
        [
            "DebateRepository(file|memory)",
            "runtime_session_store",
            "lineage_recorder + tool_audit",
            "report_generation_service",
        ],
        11,
    )
    card(s, 9.15, 4.75, 3.7, 2.0, C.L_PURP)
    txt(s, 9.4, 5.0, 3.2, 0.3, "质量与治理", 12, True, C.INDIGO)
    bullets(
        s,
        9.4,
        5.27,
        3.2,
        1.35,
        [
            "Benchmark Center",
            "Governance Center",
            "Investigation Workbench",
            "War Room 实时态势",
        ],
        11,
    )


def architecture_runtime(prs: Presentation, page: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    header(s, "3) LangGraph 运行时架构图", "Runtime Internal", page)
    # Exact node topology from backend/app/runtime/langgraph/builder.py
    card(s, 0.45, 1.2, 12.35, 5.9, C.WHITE)
    txt(s, 0.72, 1.45, 11.8, 0.32, "StateGraph 拓扑（真实节点名）", 14, True, C.INDIGO)

    # Main spine
    card(s, 0.85, 2.2, 1.4, 0.62, C.L_BLUE); txt(s, 0.95, 2.4, 1.2, 0.2, "START", 10, True, align="center")
    arr(s, 2.32, 2.42, 0.2, 0.12)
    card(s, 2.58, 2.2, 1.8, 0.62, C.L_BLUE); txt(s, 2.72, 2.4, 1.5, 0.2, "init_session", 10, True, align="center")
    arr(s, 4.42, 2.42, 0.2, 0.12)
    card(s, 4.68, 2.2, 1.65, 0.62, C.L_BLUE); txt(s, 4.81, 2.4, 1.4, 0.2, "round_start", 10, True, align="center")
    arr(s, 6.38, 2.42, 0.2, 0.12)
    card(s, 6.64, 2.2, 2.0, 0.62, C.L_PURP); txt(s, 6.78, 2.4, 1.7, 0.2, "supervisor_decide", 10, True, align="center")

    # Conditional route targets
    card(s, 2.0, 3.4, 2.25, 0.7, C.L_CYAN); txt(s, 2.14, 3.62, 1.95, 0.26, "analysis_parallel_node", 9, True, align="center")
    card(s, 4.55, 3.4, 2.45, 0.7, C.L_CYAN); txt(s, 4.7, 3.62, 2.15, 0.26, "analysis_collaboration_node*", 9, True, align="center")
    card(s, 7.35, 3.4, 2.2, 0.7, C.L_CYAN); txt(s, 7.5, 3.62, 1.9, 0.26, "speak:<agent>_node", 9, True, align="center")
    card(s, 9.85, 3.4, 1.8, 0.7, C.L_CYAN); txt(s, 9.98, 3.62, 1.55, 0.26, "round_evaluate", 9, True, align="center")
    card(s, 11.95, 3.4, 0.85, 0.7, C.L_CYAN); txt(s, 12.03, 3.62, 0.65, 0.26, "finalize", 9, True, align="center")

    # arrows from supervisor
    arr(s, 6.9, 2.95, 0.16, 0.25)
    arr(s, 7.22, 2.95, 0.16, 0.25)
    arr(s, 7.54, 2.95, 0.16, 0.25)
    arr(s, 7.86, 2.95, 0.16, 0.25)
    arr(s, 8.18, 2.95, 0.16, 0.25)

    # loops back
    left_arr(s, 2.9, 4.45, 0.22, 0.12)
    left_arr(s, 5.7, 4.45, 0.22, 0.12)
    left_arr(s, 8.45, 4.45, 0.22, 0.12)
    txt(s, 2.0, 4.67, 7.3, 0.22, "analysis_parallel / collaboration / 各agent节点 运行后均回到 supervisor_decide", 9, False, C.SUB)

    # round evaluate branch
    card(s, 9.4, 5.0, 2.2, 0.7, C.L_BLUE); txt(s, 9.58, 5.22, 1.85, 0.25, "continue_next_round=true", 9, False, align="center")
    card(s, 11.95, 5.0, 0.85, 0.7, C.L_BLUE); txt(s, 12.03, 5.22, 0.65, 0.25, "END", 9, True, align="center")
    arr(s, 10.45, 4.15, 0.2, 0.65)
    arr(s, 12.25, 4.15, 0.2, 0.65)
    txt(s, 0.72, 6.35, 11.8, 0.28, "* 当 DEBATE_ENABLE_COLLABORATION=false 时，analysis_collaboration_node 不加入图。", 10, False, C.SUB)


def agent_network(prs: Presentation, page: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    header(s, "4) 多 Agent 协作网络图", "Network & Debate", page)
    card(s, 5.3, 2.55, 2.75, 1.2, C.L_PURP)
    txt(s, 5.48, 2.9, 2.4, 0.4, "ProblemAnalysisAgent", 13, True, align="center")
    nodes: List[Tuple[str, float, float, RGBColor]] = [
        ("LogAgent", 2.0, 1.25, C.L_BLUE),
        ("DomainAgent", 5.15, 1.0, C.L_CYAN),
        ("CodeAgent", 8.35, 1.25, C.L_BLUE),
        ("MetricsAgent", 10.0, 2.95, C.L_CYAN),
        ("ChangeAgent", 8.45, 4.8, C.L_BLUE),
        ("RunbookAgent", 5.15, 5.05, C.L_CYAN),
        ("CriticAgent", 2.0, 4.8, C.L_ORANGE),
        ("RebuttalAgent", 0.55, 2.95, C.L_ORANGE),
        ("JudgeAgent", 10.0, 4.5, C.L_PURP),
        ("VerificationAgent", 10.25, 1.95, C.L_PURP),
    ]
    for name, x, y, fill in nodes:
        card(s, x, y, 2.05, 0.95, fill)
        txt(s, x + 0.12, y + 0.28, 1.8, 0.33, name, 11, True, align="center")
    for x, y in [(4.2, 1.95), (6.55, 1.95), (8.25, 2.85), (8.35, 4.25), (6.55, 4.95), (4.2, 4.95), (2.9, 4.25), (2.9, 2.85), (8.7, 2.25)]:
        arr(s, x, y, 0.95, 0.17, C.CYAN)
    txt(
        s,
        0.7,
        6.35,
        12.0,
        0.38,
        "协作规则：主 Agent 先发命令；子 Agent 结合命令与上下文决定是否调用工具；Critic/Rebuttal/Judge 负责争议收敛与结论裁决。",
        11,
        False,
        C.SUB,
    )


def tool_chain(prs: Presentation, page: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    header(s, "5) 工具调用与审计链路", "Tooling Traceability", page)
    stages = [
        ("主 Agent 命令", 0.65, 2.2, 2.05, C.L_PURP),
        ("Command Gate", 2.95, 2.2, 2.05, C.L_BLUE),
        ("Tool Execute", 5.25, 2.2, 2.25, C.L_CYAN),
        ("Output Truncate", 7.8, 2.2, 2.2, C.L_BLUE),
        ("Audit + Lineage", 10.3, 2.2, 2.3, C.L_PURP),
    ]
    for i, (name, x, y, w, fill) in enumerate(stages):
        card(s, x, y, w, 1.3, fill)
        txt(s, x + 0.1, y + 0.45, w - 0.2, 0.3, name, 11, True, align="center")
        if i < len(stages) - 1:
            arr(s, x + w + 0.08, y + 0.52, 0.2, 0.2)
    card(s, 0.65, 4.05, 6.3, 2.6, C.WHITE)
    txt(s, 0.92, 4.3, 5.8, 0.35, "关键审计字段", 16, True, C.INDIGO)
    bullets(
        s,
        0.92,
        4.72,
        5.8,
        1.8,
        [
            "event_type / timestamp / agent_name / phase",
            "tool_name / action / status / error_reason",
            "command_preview / decision_source / permission_decision",
            "I/O trace: file read, git command, http request",
        ],
        12,
    )
    card(s, 7.2, 4.05, 5.45, 2.6, C.L_BLUE)
    txt(s, 7.45, 4.3, 4.95, 0.35, "策略约束", 16, True, C.INDIGO)
    bullets(
        s,
        7.45,
        4.72,
        4.95,
        1.8,
        [
            "无命令不调用工具",
            "无配置不展示工具记录",
            "失败必须显式降级并记录",
            "前端可查看摘要与完整引用",
        ],
        12,
    )


def flow_diagram(prs: Presentation, page: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    header(s, "6) 端到端流程图", "Incident to Report", page)
    steps = [
        ("1 创建Incident", "POST /incidents"),
        ("2 创建会话", "POST /debates?incident_id"),
        ("3 建立WS连接", "/ws/debates/{session}?auto_start=true"),
        ("4 debate_service.execute_debate", "asset_collection + run orchestrator"),
        ("5 Graph执行", "round_start -> supervisor_decide -> 动态agent"),
        ("6 事件推送", "event/snapshot/result 到前端"),
        ("7 结果沉淀", "DebateResult + Lineage + ToolAudit"),
        ("8 报告生成", "JSON/Markdown/HTML + history replay"),
    ]
    x = 0.72
    for i, (a, b) in enumerate(steps):
        w = 1.45 if i not in {2, 3, 4, 7} else 1.6
        card(s, x, 2.05, w, 1.9, C.WHITE)
        txt(s, x + 0.1, 2.23, w - 0.2, 0.3, a, 10, True, C.INDIGO)
        txt(s, x + 0.1, 2.65, w - 0.2, 0.95, b, 10)
        if i < len(steps) - 1:
            arr(s, x + w + 0.04, 2.75, 0.2, 0.16)
        x += w + 0.31
    card(s, 0.72, 4.4, 12.0, 2.2, C.L_CYAN)
    txt(s, 1.0, 4.68, 11.4, 0.35, "关键控制点", 15, True, C.INDIGO)
    bullets(
        s,
        1.0,
        5.07,
        11.4,
        1.35,
        [
            "ws_debates.py 负责 snapshot/event/result 广播与任务状态更新",
            "PhaseExecutor.analysis 使用 asyncio.gather 并行执行分析Agent",
            "execution.py 统一 LLM 调用、超时计划、重试与结构化输出解析",
            "StateTransitionService 负责消息去重与状态快照合并",
        ],
        12,
    )


def state_machine(prs: Presentation, page: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    header(s, "7) 状态机与数据模型", "State & Data", page)
    card(s, 0.6, 1.2, 6.2, 5.95, C.WHITE)
    txt(s, 0.88, 1.5, 5.7, 0.35, "DebateStatus 状态机", 16, True, C.INDIGO)
    states = [
        ("PENDING", 1.0, 2.0, C.L_BLUE),
        ("RUNNING", 2.4, 2.0, C.L_BLUE),
        ("ANALYZING", 3.8, 2.0, C.L_CYAN),
        ("DEBATING", 5.2, 2.0, C.L_CYAN),
        ("JUDGING", 1.8, 3.3, C.L_PURP),
        ("COMPLETED", 3.6, 3.3, C.L_PURP),
        ("FAILED", 5.4, 3.3, C.L_ORANGE),
        ("RETRYING", 2.7, 4.5, C.L_ORANGE),
        ("CANCELLED", 4.7, 4.5, C.L_ORANGE),
    ]
    for n, x, y, fill in states:
        card(s, x, y, 1.2, 0.62, fill)
        txt(s, x + 0.05, y + 0.2, 1.08, 0.24, n, 10, True, align="center")
    for x, y in [(2.22, 2.24), (3.62, 2.24), (5.02, 2.24), (2.82, 3.0), (4.55, 3.0), (3.45, 4.2)]:
        arr(s, x, y, 0.45, 0.12)
    txt(s, 1.0, 5.45, 5.5, 1.3, "非法 phase 将被校验拦截并降级，避免会话陷入无响应。", 11, False, C.SUB)
    card(s, 7.0, 1.2, 5.7, 5.95, C.L_BLUE)
    txt(s, 7.28, 1.5, 5.2, 0.35, "核心数据对象", 16, True, C.INDIGO)
    bullets(
        s,
        7.28,
        1.9,
        5.15,
        2.25,
        [
            "Incident: 输入故障上下文",
            "DebateSession: 执行状态 + 配置",
            "DebateRound: 每轮agent输出与证据",
            "DebateResult: Top-K根因/置信度/建议",
            "LineageRecord: 时间线与工具审计",
        ],
        12,
    )
    txt(s, 7.28, 4.35, 5.15, 0.3, "本地存储策略", 14, True, C.INDIGO)
    bullets(s, 7.28, 4.68, 5.15, 1.7, ["LOCAL_STORE_BACKEND=file|memory", "LOCAL_STORE_DIR 持久化会话与报告", "暂不依赖外部数据库"], 12)


def frontend_ux(prs: Presentation, page: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    header(s, "8) 前端体验与页面结构", "UX Surface", page)
    card(s, 0.6, 1.2, 12.1, 0.95, C.WHITE)
    txt(s, 0.95, 1.52, 11.4, 0.3, "原则：实时感知、过程透明、结果可解释、关键操作可达。", 15, True, C.INDIGO)
    pages = [
        ("首页 /", "新建并启动分析、统计、agent介绍"),
        ("分析页 /incident", "资产映射 / 辩论过程 / 辩论结果"),
        ("战情页 /war-room", "时间线+工具调用+证据链+结论"),
        ("复盘台 /workbench", "关键决策回放、报告版本对比"),
        ("工具中心 /tools", "工具管理、试运行、审计查询"),
        ("治理中心 /governance", "成功率、超时率、成本与风险"),
    ]
    for i, (name, desc) in enumerate(pages):
        x = 0.6 + (i % 2) * 6.15
        y = 2.45 + (i // 2) * 1.55
        card(s, x, y, 5.95, 1.35, C.L_BLUE if i % 2 == 0 else C.L_CYAN)
        txt(s, x + 0.2, y + 0.2, 5.45, 0.3, name, 13, True, C.INDIGO)
        txt(s, x + 0.2, y + 0.55, 5.45, 0.6, desc, 12)
    card(s, 0.6, 5.75, 12.1, 1.35, C.L_PURP)
    txt(s, 0.95, 5.98, 11.4, 0.3, "已实现关键体验优化", 14, True, C.INDIGO)
    bullets(
        s,
        0.95,
        6.3,
        11.4,
        0.65,
        [
            "聊天消息支持缩略+展开，避免大段JSON直接干扰阅读",
            "全链路统一北京时间显示",
            "分析中任务可进入详情持续追踪，不跳空白页",
        ],
        12,
    )


def quality_governance(prs: Presentation, page: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    header(s, "9) 质量保障与治理", "Benchmark & Governance", page)
    card(s, 0.62, 1.2, 4.05, 5.95, C.L_BLUE)
    txt(s, 0.92, 1.5, 3.4, 0.3, "Benchmark 指标", 16, True, C.INDIGO)
    bars = [("Top1命中率", 0.72, C.GREEN), ("Top3命中率", 0.88, C.GREEN), ("超时率", 0.21, C.RED), ("空结论率", 0.12, C.YELLOW)]
    y = 2.08
    for name, rate, color in bars:
        txt(s, 0.92, y, 1.8, 0.25, name, 11)
        bgbar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.15), Inches(y + 0.02), Inches(2.2), Inches(0.18))
        bgbar.fill.solid()
        bgbar.fill.fore_color.rgb = RGBColor(222, 231, 244)
        bgbar.line.fill.background()
        fgbar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.15), Inches(y + 0.02), Inches(2.2 * rate), Inches(0.18))
        fgbar.fill.solid()
        fgbar.fill.fore_color.rgb = color
        fgbar.line.fill.background()
        txt(s, 4.35, y - 0.01, 0.3, 0.2, f"{int(rate*100)}%", 10, False, C.SUB)
        y += 0.7
    txt(s, 0.92, 5.55, 3.5, 1.2, "指标值由评测中心自动计算；可接入 CI 作为发布门禁。", 10, False, C.SUB)
    card(s, 4.9, 1.2, 7.8, 2.75, C.WHITE)
    txt(s, 5.2, 1.5, 7.2, 0.33, "可观测机制", 16, True, C.INDIGO)
    bullets(s, 5.2, 1.9, 7.2, 1.8, ["全链路时间戳与事件序号", "Agent异常分类与降级策略", "工具I/O审计与引用ID", "会话重放与报告差异对比"], 12)
    card(s, 4.9, 4.15, 7.8, 2.95, C.L_CYAN)
    txt(s, 5.2, 4.42, 7.2, 0.33, "治理看板关注点", 16, True, C.INDIGO)
    bullets(s, 5.2, 4.82, 7.2, 2.0, ["团队维度成功率/超时率/失败类型", "工具维度可用性与错误分布", "会话维度关键决策链路", "版本维度结果对比与回归监测"], 12)


def roadmap(prs: Presentation, page: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    header(s, "10) 路线图", "P0~P3", page)
    phases = [
        ("P0 可用性", "自动调查、并行拉证据、超时切换、pending治理", C.L_ORANGE),
        ("P1 准确率", "跨源证据约束、因果推理、Top-K区间", C.L_BLUE),
        ("P2 可控修复", "修复状态机、No-regression gate、回滚", C.L_CYAN),
        ("P3 平台治理", "反馈学习、A/B评测、多租户与外部协同", C.L_PURP),
    ]
    for i, (n, d, fill) in enumerate(phases):
        y = 1.45 + i * 1.35
        card(s, 0.8, y, 11.9, 1.05, fill)
        txt(s, 1.08, y + 0.24, 2.1, 0.35, n, 15, True, C.INDIGO)
        txt(s, 3.25, y + 0.24, 9.2, 0.35, d, 14)
    card(s, 0.8, 6.0, 11.9, 1.05, C.WHITE)
    txt(s, 1.08, 6.35, 11.2, 0.3, "执行建议：按阶段验收 + Benchmark门禁 + 前后端联调验收。", 13, False, C.SUB)


def ending(prs: Presentation, page: int, kind: str) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, prs)
    header(s, "总结", f"{kind} Summary", page)
    card(s, 1.0, 1.8, 11.3, 4.9, C.WHITE)
    txt(
        s,
        1.35,
        2.35,
        10.5,
        0.55,
        "系统已具备“多 Agent 协作 + 工具审计 + 可回放治理”的生产问题分析闭环。",
        22,
        True,
        C.INDIGO,
    )
    bullets(
        s,
        1.35,
        3.15,
        10.2,
        2.9,
        [
            "架构完整：FastAPI + LangGraph + React",
            "过程透明：命令驱动、工具审计、谱系回放",
            "结果可用：Top-K根因、证据链、报告和验证计划",
            "可持续优化：Benchmark + Governance 双闭环",
        ],
        16,
    )
    txt(s, 1.35, 6.35, 10.2, 0.32, "下一步：强化真实数据源接入与自治修复能力。", 14, False, C.SUB)


def add_tech_extra(prs: Presentation) -> None:
    # 11 extra technical slides to reach 18 total
    # 1 API
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, prs); header(s, "11) API 与扩展能力", "API Surface", 11)
    card(s, 0.65, 1.2, 6.2, 5.95, C.WHITE); txt(s, 0.95, 1.48, 5.6, 0.35, "核心接口", 16, True, C.INDIGO)
    bullets(s, 0.95, 1.9, 5.5, 4.9, ["POST /api/v1/incidents/", "POST /api/v1/debates/?incident_id=...", "POST /api/v1/debates/{session_id}/execute", "GET /api/v1/debates/{session_id}/result", "POST /api/v1/reports/{incident_id}/regenerate", "GET/PUT /api/v1/settings/tooling", "WS /ws/debates/{session_id}?auto_start=true"], 12)
    card(s, 7.05, 1.2, 5.65, 2.85, C.L_BLUE); txt(s, 7.32, 1.48, 5.1, 0.35, "模型配置", 15, True, C.INDIGO)
    bullets(s, 7.32, 1.88, 5.1, 1.9, ["LLM_BASE_URL=ark...", "LLM_MODEL=kimi-k2.5", "LLM_API_KEY=环境变量注入", "分阶段 timeout + max_concurrency"], 11)
    card(s, 7.05, 4.3, 5.65, 2.85, C.L_CYAN); txt(s, 7.32, 4.58, 5.1, 0.35, "扩展入口", 15, True, C.INDIGO)
    bullets(s, 7.32, 4.98, 5.1, 1.9, ["TelemetryConnector", "CMDBConnector", "Tool Registry 生命周期管理", "Skill化 Prompt 策略"], 11)

    # 12 tool matrix
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, prs); header(s, "12) Agent 工具能力矩阵", "Agent x Tool", 12)
    card(s, 0.6, 1.2, 12.1, 5.95, C.WHITE); txt(s, 0.9, 1.48, 11.5, 0.35, "当前工具能力与门禁策略", 16, True, C.INDIGO)
    rows = [
        ("CodeAgent", "Git 仓库检索", "已接入", "主Agent命令 + 配置开启"),
        ("LogAgent", "本地日志读取", "已接入", "主Agent命令 + 配置开启"),
        ("DomainAgent", "Excel/CSV 资产映射", "已接入", "主Agent命令 + 配置开启"),
        ("MetricsAgent", "Telemetry Connector", "可接入入口", "配置开启后调用"),
        ("ChangeAgent", "CMDB/变更源", "可接入入口", "配置开启后调用"),
        ("ProblemAnalysisAgent", "规则建议工具包", "已接入", "聚合证据与调度参考"),
    ]
    y = 2.0
    for a, b, c, d in rows:
        fill = C.L_BLUE if int(y * 10) % 2 == 0 else C.L_CYAN
        card(s, 0.9, y, 11.4, 0.66, fill)
        txt(s, 1.1, y + 0.2, 2.2, 0.25, a, 11, True)
        txt(s, 3.3, y + 0.2, 2.4, 0.25, b, 11)
        txt(s, 5.9, y + 0.2, 1.6, 0.25, c, 11, True, C.INDIGO)
        txt(s, 7.6, y + 0.2, 4.5, 0.25, d, 11, False, C.SUB)
        y += 0.78

    # 13 debate rounds and policy
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, prs); header(s, "13) 辩论策略与轮次控制", "Debate Policy", 13)
    card(s, 0.62, 1.2, 5.95, 5.95, C.L_PURP); txt(s, 0.92, 1.5, 5.35, 0.35, "策略配置", 16, True, C.INDIGO)
    bullets(s, 0.92, 1.9, 5.25, 4.9, ["DEBATE_MAX_ROUNDS 默认=1（可配置）", "DEBATE_ENABLE_CRITIQUE=true", "DEBATE_ENABLE_COLLABORATION 可切换", "DEBATE_REQUIRE_EFFECTIVE_LLM_CONCLUSION=true", "runtime_strategy_center 依据严重度选策略", "execution_mode: standard/quick/background/async"], 12)
    card(s, 6.85, 1.2, 5.85, 5.95, C.L_BLUE); txt(s, 7.13, 1.5, 5.3, 0.35, "性能与耗时优化", 16, True, C.INDIGO)
    bullets(s, 7.13, 1.9, 5.25, 4.9, ["按阶段设置 timeout（分析/评审/裁决/报告）", "并发控制 LLM_MAX_CONCURRENCY", "长上下文 compaction + prune", "超时重试与局部降级防止卡死", "结果结构化约束减少“空报告”"], 12)

    # 14 runtime sequence (code-aligned)
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, prs); header(s, "14) 运行时序图（代码对齐）", "WS + Debate + LangGraph", 14)
    card(s, 0.62, 1.2, 12.1, 5.95, C.WHITE)
    txt(s, 0.92, 1.48, 11.4, 0.35, "前后端联动时序（对应 ws_debates.py + debate_service.py + langgraph_runtime.py）", 14, True, C.INDIGO)

    actors = [
        ("Frontend", 1.0),
        ("WSManager", 3.1),
        ("DebateService", 5.2),
        ("LangGraphRuntime", 7.55),
        ("Agents/Tools", 10.0),
    ]
    for name, x in actors:
        card(s, x, 2.0, 1.9, 0.6, C.L_BLUE if name != "Agents/Tools" else C.L_CYAN)
        txt(s, x + 0.1, 2.2, 1.7, 0.24, name, 10, True, align="center")
        # lifeline
        lifeline = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.93), Inches(2.62), Inches(0.03), Inches(3.8))
        lifeline.fill.solid()
        lifeline.fill.fore_color.rgb = RGBColor(211, 223, 238)
        lifeline.line.fill.background()

    # messages
    # Frontend -> WSManager
    arr(s, 2.0, 2.95, 1.05, 0.12); txt(s, 2.03, 2.78, 1.15, 0.2, "connect ws(auto_start)", 8, False, C.SUB)
    # WSManager -> DebateService
    arr(s, 4.1, 3.25, 1.05, 0.12); txt(s, 4.18, 3.08, 1.1, 0.2, "execute_debate()", 8, False, C.SUB)
    # DebateService -> LangGraph
    arr(s, 6.2, 3.55, 1.25, 0.12); txt(s, 6.3, 3.38, 1.2, 0.2, "orchestrator.execute()", 8, False, C.SUB)
    # LangGraph -> Agents
    arr(s, 8.55, 3.85, 1.35, 0.12); txt(s, 8.7, 3.68, 1.2, 0.2, "supervisor -> agents", 8, False, C.SUB)
    # Agents -> LangGraph
    left_arr(s, 8.55, 4.2, 1.35, 0.12); txt(s, 8.66, 4.03, 1.2, 0.2, "evidence/tool result", 8, False, C.SUB)
    # LangGraph -> DebateService
    left_arr(s, 6.2, 4.55, 1.25, 0.12); txt(s, 6.27, 4.38, 1.2, 0.2, "final payload", 8, False, C.SUB)
    # DebateService -> WSManager
    left_arr(s, 4.1, 4.9, 1.05, 0.12); txt(s, 4.2, 4.73, 1.0, 0.2, "event/result", 8, False, C.SUB)
    # WSManager -> Frontend
    left_arr(s, 2.0, 5.25, 1.05, 0.12); txt(s, 2.06, 5.08, 1.0, 0.2, "snapshot/event/result", 8, False, C.SUB)

    card(s, 0.92, 5.65, 11.5, 1.15, C.L_CYAN)
    bullets(
        s,
        1.15,
        5.86,
        11.0,
        0.75,
        [
            "事件在 event_callback 中实时透传；前端不需要轮询即可看到阶段推进和agent发言。",
            "异常路径会广播 error + snapshot，避免“前端空白但后端已失败”的黑盒体验。",
        ],
        11,
    )

    # 15 report visualization
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, prs); header(s, "15) 报告可视化设计", "Report UX", 15)
    card(s, 0.62, 1.2, 12.1, 5.95, C.WHITE)
    txt(s, 0.92, 1.5, 11.4, 0.35, "报告页采用模块化可视化，不再直接裸展示 Markdown", 16, True, C.INDIGO)
    mods = [
        ("诊断总览", "主结论、置信度、风险级别、会话信息"),
        ("Top-K 根因候选", "排名、置信度区间、证据强度"),
        ("证据链", "来源、时间、引用路径，可展开详情"),
        ("修复建议", "动作清单、风险评估、验证计划"),
        ("报告模块卡片", "按主题结构化分组，支持展开全文"),
        ("下载导出", "markdown/json/html 多格式输出"),
    ]
    for i, (n, d) in enumerate(mods):
        x = 0.92 + (i % 2) * 6.0
        y = 2.05 + (i // 2) * 1.45
        card(s, x, y, 5.75, 1.2, C.L_BLUE if i % 2 == 0 else C.L_CYAN)
        txt(s, x + 0.2, y + 0.18, 5.3, 0.25, n, 12, True, C.INDIGO)
        txt(s, x + 0.2, y + 0.5, 5.3, 0.45, d, 11)

    # 16 benchmark/governance deep
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, prs); header(s, "16) Benchmark 与治理深度能力", "Ops Excellence", 16)
    card(s, 0.62, 1.2, 5.95, 5.95, C.L_BLUE); txt(s, 0.92, 1.5, 5.35, 0.35, "Benchmark Center", 16, True, C.INDIGO)
    bullets(s, 0.92, 1.9, 5.25, 4.9, ["支持样本数/超时参数调节", "输出 Top1/Top3、overlap、timeout、empty rate", "历史基线追踪与对比", "可接 CI 作为回归门禁"], 12)
    card(s, 6.85, 1.2, 5.85, 5.95, C.L_CYAN); txt(s, 7.13, 1.5, 5.3, 0.35, "Governance Center", 16, True, C.INDIGO)
    bullets(s, 7.13, 1.9, 5.25, 4.9, ["按团队展示成功率与失败率", "按会话分析超时热点", "按工具统计调用成功率与错误", "支撑策略迭代与运营决策"], 12)

    # 17 risk and controls
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, prs); header(s, "17) 风险点与控制措施", "Risk Controls", 17)
    card(s, 0.62, 1.2, 12.1, 5.95, C.WHITE)
    txt(s, 0.92, 1.5, 11.4, 0.35, "关键风险与缓解策略", 16, True, C.INDIGO)
    risks = [
        ("LLM超时或慢响应", "分阶段timeout + 重试 + fallback + 并发限制"),
        ("工具调用不可用", "开关控制 + 健康探测 + 审计降级"),
        ("分析结论为空", "强制有效结论策略 + 再分析触发"),
        ("消息重复/显示不完整", "前端去重 + 流式协议标准化"),
        ("状态机异常卡死", "严格状态迁移校验 + 异常恢复"),
        ("责任田映射遗漏", "资源规范化 + 多源映射兜底"),
    ]
    y = 2.02
    for r, m in risks:
        card(s, 0.92, y, 11.4, 0.72, C.L_BLUE if int(y * 10) % 2 == 0 else C.L_CYAN)
        txt(s, 1.15, y + 0.22, 3.7, 0.25, r, 11, True, C.INDIGO)
        txt(s, 4.95, y + 0.22, 7.1, 0.25, m, 11)
        y += 0.82

    # 18 ending
    ending(prs, 18, "Technical")


def build_exec() -> Path:
    prs = mk_prs()
    cover(prs, "生产问题根因分析系统", "管理层汇报版（Executive Deck）")
    agenda(prs, 2)
    business_value(prs, 3)
    architecture_e2e(prs, 4)
    flow_diagram(prs, 5)
    frontend_ux(prs, 6)
    quality_governance(prs, 7)
    roadmap(prs, 8)
    OUT_EXEC.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_EXEC))
    return OUT_EXEC


def build_tech() -> Path:
    prs = mk_prs()
    cover(prs, "生产问题根因分析系统", "技术详版（Technical Deep Dive）")
    agenda(prs, 2)
    business_value(prs, 3)
    architecture_e2e(prs, 4)
    architecture_runtime(prs, 5)
    agent_network(prs, 6)
    tool_chain(prs, 7)
    flow_diagram(prs, 8)
    state_machine(prs, 9)
    frontend_ux(prs, 10)
    add_tech_extra(prs)
    OUT_TECH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_TECH))
    return OUT_TECH


def main() -> None:
    e = build_exec()
    t = build_tech()
    print(f"Generated: {e}")
    print(f"Generated: {t}")


if __name__ == "__main__":
    main()
