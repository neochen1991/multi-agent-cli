#!/usr/bin/env python3
"""
Generate a visual PPTX deck for:
Vibe Coding in production SRE agent systems.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output" / "vibe-coding-ppt"
ASSET_DIR = OUT_DIR / "assets"
OUT_PPTX = OUT_DIR / "vibe-coding-best-practices-2026-03-06.pptx"
OUT_README = OUT_DIR / "README.md"


class Theme:
    white = RGBColor(255, 255, 255)
    bg = RGBColor(245, 248, 253)
    navy = RGBColor(18, 42, 76)
    blue = RGBColor(44, 110, 226)
    cyan = RGBColor(19, 149, 170)
    green = RGBColor(32, 151, 108)
    orange = RGBColor(230, 142, 37)
    red = RGBColor(210, 75, 75)
    text = RGBColor(33, 42, 53)
    subtext = RGBColor(95, 108, 123)
    border = RGBColor(206, 217, 233)
    light_blue = RGBColor(234, 242, 255)
    light_cyan = RGBColor(234, 250, 253)
    light_orange = RGBColor(255, 244, 232)
    light_green = RGBColor(232, 249, 242)
    light_red = RGBColor(254, 239, 239)


def ensure_dirs() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)

def add_chart(
    slide,
    *,
    chart_type: XL_CHART_TYPE,
    x: float,
    y: float,
    w: float,
    h: float,
    categories: list[str],
    series: list[tuple[str, list[float]]],
) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series:
        chart_data.add_series(name, values)

    chart = slide.shapes.add_chart(
        chart_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        chart_data,
    ).chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    if chart.value_axis is not None:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.font.size = Pt(10)
    if chart.category_axis is not None:
        chart.category_axis.tick_labels.font.size = Pt(10)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_bg(slide, prs: Presentation) -> None:
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = Theme.bg
    bg.line.fill.background()


def add_topbar(slide, title: str, subtitle: str, page: int) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.86)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = Theme.navy
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.42), Inches(0.16), Inches(8.2), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.color.rgb = Theme.white
    p.font.bold = True
    p.font.size = Pt(24)

    sbox = slide.shapes.add_textbox(Inches(8.8), Inches(0.22), Inches(3.6), Inches(0.35))
    sp = sbox.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.color.rgb = RGBColor(206, 221, 247)
    sp.font.size = Pt(11)
    sp.alignment = PP_ALIGN.RIGHT

    pbox = slide.shapes.add_textbox(Inches(12.55), Inches(0.2), Inches(0.45), Inches(0.3))
    pp = pbox.text_frame.paragraphs[0]
    pp.text = str(page)
    pp.font.color.rgb = RGBColor(175, 200, 236)
    pp.font.bold = True
    pp.font.size = Pt(12)
    pp.alignment = PP_ALIGN.RIGHT


def add_card(slide, x: float, y: float, w: float, h: float, fill: RGBColor = Theme.white):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = Theme.border
    return card


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: int = 16,
    bold: bool = False,
    color: RGBColor = Theme.text,
    align: str = "left",
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT


def add_bullets(
    slide, x: float, y: float, w: float, h: float, lines: Iterable[str], *, size: int = 14, color: RGBColor = Theme.text
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(size)
        p.font.color.rgb = color


def add_arrow(slide, x: float, y: float, w: float, h: float, color: RGBColor = Theme.blue) -> None:
    arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()


def slide_1_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.2)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = Theme.navy
    band.line.fill.background()

    add_text(slide, 0.65, 0.6, 9.8, 0.85, "Vibe Coding 实战方法论", size=44, bold=True, color=Theme.white)
    add_text(
        slide,
        0.65,
        1.45,
        10.8,
        0.45,
        "从 Prompt 驱动到 Harness 驱动（业界最佳实践 + 本项目经验）",
        size=17,
        color=RGBColor(209, 223, 246),
    )

    add_card(slide, 0.7, 2.7, 6.0, 3.85, Theme.white)
    add_text(slide, 1.0, 3.05, 5.4, 0.5, "分享目标", size=24, bold=True, color=Theme.blue)
    add_bullets(
        slide,
        1.0,
        3.55,
        5.5,
        2.7,
        [
            "借鉴业界优秀 Vibe Coding 团队的工程模式",
            "给出 Skill / MCP / 多 Agent 的落地方法",
            "用本项目作为样例展示如何避免“快而不稳”",
            "形成可执行的 2 周 + 90 天路线图",
        ],
        size=15,
    )

    add_card(slide, 7.0, 2.7, 5.6, 3.85, Theme.light_blue)
    add_text(slide, 7.3, 3.05, 4.9, 0.5, "核心结构", size=24, bold=True, color=Theme.blue)
    add_text(slide, 7.5, 3.7, 4.8, 0.35, "Harness Engineering", size=16, bold=True, color=Theme.text)
    add_text(slide, 7.5, 4.1, 4.8, 0.35, "Context Engineering", size=16, bold=True, color=Theme.text)
    add_text(slide, 7.5, 4.5, 4.8, 0.35, "Skill / MCP / Tooling", size=16, bold=True, color=Theme.text)
    add_text(slide, 7.5, 4.9, 4.8, 0.35, "Multi-Agent Governance", size=16, bold=True, color=Theme.text)
    add_text(slide, 7.5, 5.3, 4.8, 0.35, "Quality Gate & Eval", size=16, bold=True, color=Theme.text)

    add_text(
        slide,
        0.7,
        6.9,
        11.5,
        0.3,
        f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}（Asia/Shanghai）",
        size=11,
        color=Theme.subtext,
    )


def slide_2_why(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "为什么现在必须做 Vibe Coding", "Speed with quality", 2)

    add_card(slide, 0.6, 1.15, 5.6, 5.95, Theme.white)
    add_text(slide, 0.9, 1.45, 5.0, 0.55, "结论：价值在“缩短认知到交付闭环”", size=18, bold=True, color=Theme.blue)
    add_bullets(
        slide,
        0.9,
        2.05,
        5.0,
        2.3,
        [
            "不是替代工程师，而是提升工程杠杆率",
            "降低跨角色沟通损耗：需求、实现、验证同轴推进",
            "把高频重复工程动作自动化，让人聚焦决策",
            "前提：必须有约束、评测、审计，不然会放大错误",
        ],
        size=14,
    )
    add_card(slide, 0.9, 4.6, 4.9, 2.1, Theme.light_orange)
    add_text(slide, 1.1, 4.9, 4.4, 0.35, "常见误解", size=14, bold=True, color=Theme.orange)
    add_bullets(
        slide,
        1.1,
        5.25,
        4.5,
        1.3,
        ["“AI 会写代码 = 生产可用”", "“Prompt 写长一点就行”", "“先快后补质量”"],
        size=12,
        color=Theme.text,
    )

    add_card(slide, 6.4, 1.15, 6.3, 5.95, Theme.white)
    add_chart(
        slide,
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        x=6.65,
        y=1.65,
        w=5.8,
        h=3.7,
        categories=["需求澄清", "方案迭代", "实现速度", "缺陷回归", "上线稳定"],
        series=[
            ("传统流程", [42, 38, 26, 72, 55]),
            ("Vibe + Harness", [71, 74, 83, 34, 81]),
        ],
    )
    add_text(slide, 6.7, 5.65, 5.7, 0.35, "示意：Vibe + Harness 在效率与稳定性上可同时提升", size=12, color=Theme.subtext)


def slide_3_harness(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "业界共识 1：Harness Engineering", "System over ad-hoc prompts", 3)

    add_text(slide, 0.8, 1.15, 12.0, 0.45, "核心观点：Harness 决定生产上限，Prompt 只决定局部输出质量", size=20, bold=True, color=Theme.blue)

    levels = [
        ("策略层", "目标、约束、验收标准", Theme.light_blue, Theme.blue),
        ("编排层", "任务拆解、状态机、路由", Theme.light_cyan, Theme.cyan),
        ("执行层", "模型调用、工具调用、上下文注入", Theme.light_green, Theme.green),
        ("防护层", "测试、评测、审计、回滚", Theme.light_orange, Theme.orange),
    ]
    y = 2.0
    for name, desc, bg, c in levels:
        add_card(slide, 1.2, y, 10.9, 0.95, bg)
        add_text(slide, 1.45, y + 0.2, 1.8, 0.35, name, size=15, bold=True, color=c)
        add_text(slide, 3.35, y + 0.2, 8.3, 0.35, desc, size=14, color=Theme.text)
        y += 1.12

    add_card(slide, 1.2, 6.6, 10.9, 0.6, Theme.white)
    add_text(slide, 1.45, 6.77, 10.2, 0.25, "实践要点：让模型负责“思考”，让系统负责“确定性控制”。", size=13, color=Theme.subtext)


def slide_4_context(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "业界共识 2：上下文工程优先", "Context > long prompts", 4)

    add_card(slide, 0.7, 1.25, 5.6, 5.8, Theme.white)
    add_text(slide, 1.0, 1.55, 5.0, 0.4, "上下文漏斗", size=20, bold=True, color=Theme.blue)
    add_card(slide, 1.2, 2.1, 4.7, 0.85, Theme.light_blue)
    add_text(slide, 1.45, 2.38, 4.2, 0.3, "原始资料：日志、代码、监控、文档", size=13)
    add_arrow(slide, 2.85, 2.95, 1.4, 0.28)
    add_card(slide, 1.5, 3.4, 4.1, 0.85, Theme.light_cyan)
    add_text(slide, 1.75, 3.68, 3.6, 0.3, "清洗与结构化：去噪、切片、标注", size=13)
    add_arrow(slide, 2.85, 4.25, 1.4, 0.28)
    add_card(slide, 1.9, 4.7, 3.3, 0.85, Theme.light_green)
    add_text(slide, 2.1, 4.98, 2.9, 0.3, "任务相关上下文：当前步骤所需", size=13)
    add_arrow(slide, 2.85, 5.55, 1.4, 0.28)
    add_card(slide, 2.2, 6.0, 2.7, 0.72, Theme.light_orange)
    add_text(slide, 2.35, 6.22, 2.4, 0.25, "模型输入", size=13, bold=True, color=Theme.orange)

    add_card(slide, 6.55, 1.25, 6.1, 5.8, Theme.white)
    add_text(slide, 6.85, 1.6, 5.4, 0.35, "最佳实践（来自 Copilot/Cursor/Claude 实战）", size=16, bold=True, color=Theme.blue)
    add_bullets(
        slide,
        6.85,
        2.05,
        5.4,
        3.0,
        [
            "规则文件定义“如何做”，而不是每次重写提示词",
            "按任务最小化注入上下文，减少 token 浪费与漂移",
            "使用结构化输入（JSON/schema）降低解析不确定性",
            "分步链式执行：先计划，再执行，再验证",
            "每步都可回放，便于复盘与纠偏",
        ],
        size=14,
    )
    add_card(slide, 6.85, 5.6, 5.4, 1.2, Theme.light_orange)
    add_text(slide, 7.05, 5.9, 5.0, 0.28, "一句话：不要让模型“猜你的上下文”，要把上下文工程化交给系统。", size=12)


def slide_5_workflow(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "业界共识 3：标准工作流", "Explore -> Plan -> Build -> Verify -> Release", 5)

    add_text(slide, 0.8, 1.15, 12.0, 0.45, "固定节奏是高质量 Vibe Coding 的基础设施", size=20, bold=True, color=Theme.blue)

    stages = [
        ("Explore", "澄清目标/约束/成功标准"),
        ("Plan", "产出方案与任务清单"),
        ("Build", "按步骤实现，保持小步提交"),
        ("Verify", "单测/集成/评测门禁"),
        ("Release", "发布、回滚、复盘"),
    ]
    x = 0.8
    colors = [Theme.light_blue, Theme.light_cyan, Theme.light_green, Theme.light_orange, Theme.light_red]
    for i, (name, desc) in enumerate(stages):
        add_card(slide, x, 2.0, 2.3, 3.8, colors[i])
        add_text(slide, x + 0.2, 2.3, 1.9, 0.35, name, size=16, bold=True, color=Theme.navy, align="center")
        add_text(slide, x + 0.2, 2.8, 1.9, 2.6, desc, size=12, align="center")
        if i < len(stages) - 1:
            add_arrow(slide, x + 2.3, 3.5, 0.45, 0.28, Theme.blue)
        x += 2.55

    add_card(slide, 0.8, 6.2, 11.9, 0.85, Theme.white)
    add_text(
        slide,
        1.0,
        6.45,
        11.2,
        0.3,
        "落地建议：每一步都定义“完成标准”，任何一步失败都不允许带病进入下一步。",
        size=13,
        color=Theme.subtext,
    )


def slide_6_gate(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "业界共识 4：测试与评测门禁", "Quality gate as brakes", 6)

    add_card(slide, 0.6, 1.2, 6.1, 5.85, Theme.white)
    add_chart(
        slide,
        chart_type=XL_CHART_TYPE.BAR_CLUSTERED,
        x=0.9,
        y=1.65,
        w=5.5,
        h=3.7,
        categories=["代码生成", "静态检查", "单测", "集成测试", "Benchmark Gate", "发布"],
        series=[("通过样本数", [1000, 820, 690, 510, 340, 292])],
    )
    add_text(slide, 0.9, 5.75, 5.3, 0.35, "门禁越后越严格，成本越高但风险越低", size=12, color=Theme.subtext)

    add_card(slide, 6.95, 1.2, 5.75, 5.85, Theme.white)
    add_text(slide, 7.25, 1.55, 5.1, 0.35, "团队实施清单", size=17, bold=True, color=Theme.blue)
    add_bullets(
        slide,
        7.25,
        2.05,
        5.0,
        3.4,
        [
            "本地最小测试：语法/静态检查/关键单测",
            "PR 门禁：集成测试 + 关键路径端到端",
            "CI 评测：失败率、超时率、稳定性评分",
            "Benchmark 不达标直接阻断发布",
            "每周复盘：误报/漏报/回归根因",
        ],
        size=14,
    )
    add_card(slide, 7.25, 5.7, 5.0, 1.1, Theme.light_green)
    add_text(slide, 7.45, 6.0, 4.6, 0.28, "指标优先：Top1 命中率、超时率、长期 pending 占比", size=12)


def slide_7_mcp(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "业界共识 5：MCP/工具接入标准化", "Capability with governance", 7)

    add_text(slide, 0.8, 1.15, 12.0, 0.45, "统一接口 + 权限边界 + 审计追踪，是工具化成功的三要素", size=19, bold=True, color=Theme.blue)

    add_card(slide, 0.8, 1.85, 11.8, 4.9, Theme.white)

    # Layer 1
    add_card(slide, 1.2, 2.25, 10.9, 0.95, Theme.light_blue)
    add_text(slide, 1.45, 2.52, 2.4, 0.32, "Agent Layer", size=14, bold=True, color=Theme.blue)
    add_text(slide, 3.7, 2.52, 7.8, 0.32, "ProblemAnalysisAgent / LogAgent / CodeAgent / DatabaseAgent", size=13)

    # Layer 2
    add_arrow(slide, 6.2, 3.25, 0.8, 0.25, Theme.cyan)
    add_card(slide, 1.2, 3.55, 10.9, 0.95, Theme.light_cyan)
    add_text(slide, 1.45, 3.82, 2.4, 0.32, "MCP/Tool Gateway", size=14, bold=True, color=Theme.cyan)
    add_text(slide, 3.7, 3.82, 7.8, 0.32, "命令门禁 | 超时/重试 | 参数规范化 | 调用审计", size=13)

    # Layer 3
    add_arrow(slide, 6.2, 4.55, 0.8, 0.25, Theme.green)
    add_card(slide, 1.2, 4.85, 10.9, 0.95, Theme.light_green)
    add_text(slide, 1.45, 5.12, 2.4, 0.32, "Data Sources", size=14, bold=True, color=Theme.green)
    add_text(slide, 3.7, 5.12, 7.8, 0.32, "日志平台 | Git 仓库 | PostgreSQL | APM | 监控告警", size=13)

    add_card(slide, 1.2, 6.0, 10.9, 0.55, Theme.light_orange)
    add_text(slide, 1.45, 6.17, 10.3, 0.25, "安全基线：最小权限、禁用默认高风险动作、调用全链路可回放。", size=12, color=Theme.text)


def slide_8_multi_agent(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "业界共识 6：多 Agent 协作机制", "Protocol over parallel noise", 8)

    add_card(slide, 0.7, 1.2, 12.0, 5.95, Theme.white)
    add_text(slide, 0.95, 1.5, 11.4, 0.35, "关键：主控协调 + 专家分工 + 裁决收敛，而不是无序并行", size=18, bold=True, color=Theme.blue)

    # Center commander
    add_card(slide, 5.45, 3.15, 2.35, 1.0, Theme.light_blue)
    add_text(slide, 5.6, 3.47, 2.05, 0.3, "Main Agent", size=14, bold=True, align="center")

    nodes = [
        ("LogAgent", 1.4, 2.0, Theme.light_cyan),
        ("CodeAgent", 1.4, 4.3, Theme.light_cyan),
        ("DomainAgent", 4.1, 5.35, Theme.light_green),
        ("DatabaseAgent", 8.2, 5.35, Theme.light_green),
        ("CriticAgent", 10.0, 2.0, Theme.light_orange),
        ("JudgeAgent", 10.0, 4.3, Theme.light_orange),
    ]
    for name, x, y, fill in nodes:
        add_card(slide, x, y, 2.0, 0.85, fill)
        add_text(slide, x + 0.2, y + 0.26, 1.6, 0.25, name, size=12, bold=True, align="center")

    # arrows around
    add_arrow(slide, 3.45, 2.35, 1.85, 0.2)
    add_arrow(slide, 3.45, 4.65, 1.85, 0.2)
    add_arrow(slide, 7.85, 2.35, 1.85, 0.2)
    add_arrow(slide, 7.85, 4.65, 1.85, 0.2)
    add_arrow(slide, 6.5, 4.15, 0.3, 0.95)
    add_arrow(slide, 6.5, 2.2, 0.3, 0.95)

    add_text(slide, 0.95, 6.65, 11.4, 0.28, "协作协议：命令分发 -> 证据反馈 -> 交叉质疑 -> 主控收敛 -> 裁决输出", size=12, color=Theme.subtext)


def slide_9_antipattern(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "业界反模式（必须避开）", "Common failure patterns", 9)

    cards = [
        ("盲目 Accept-All", "无验证直接合并，回归概率高", 0.8, 1.55, Theme.light_red, Theme.red),
        ("超长 Prompt 堆砌", "上下文噪声增加，模型漂移", 6.9, 1.55, Theme.light_orange, Theme.orange),
        ("工具无审计", "出问题无法追责，安全风险高", 0.8, 4.0, Theme.light_orange, Theme.orange),
        ("异常无终态", "任务卡在 pending，体验崩坏", 6.9, 4.0, Theme.light_red, Theme.red),
    ]
    for title, desc, x, y, fill, color in cards:
        add_card(slide, x, y, 5.6, 2.15, fill)
        add_text(slide, x + 0.3, y + 0.35, 4.9, 0.35, title, size=19, bold=True, color=color)
        add_text(slide, x + 0.3, y + 0.9, 4.9, 0.95, desc, size=14, color=Theme.text)

    add_card(slide, 0.8, 6.45, 11.7, 0.65, Theme.white)
    add_text(slide, 1.0, 6.63, 11.2, 0.26, "修正策略：先约束流程，再放大生成速度。质量门禁必须前置。", size=12, color=Theme.subtext)


def slide_10_project_harness(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "本项目案例 1：Harness 化落地", "Runtime orchestration", 10)

    add_card(slide, 0.7, 1.2, 12.0, 5.95, Theme.white)
    add_text(slide, 1.0, 1.5, 11.2, 0.35, "关键实现：模型推理与系统控制分层，确保可控与可回放", size=18, bold=True, color=Theme.blue)

    # chain boxes
    steps = [
        ("用户输入故障信息", 1.0, 2.3, Theme.light_blue),
        ("Main Agent 分发命令", 3.15, 2.3, Theme.light_cyan),
        ("专家 Agent 分析", 5.3, 2.3, Theme.light_green),
        ("Judge 收敛结论", 7.45, 2.3, Theme.light_orange),
        ("报告生成与回放", 9.6, 2.3, Theme.light_red),
    ]
    for i, (label, x, y, fill) in enumerate(steps):
        add_card(slide, x, y, 1.95, 1.05, fill)
        add_text(slide, x + 0.1, y + 0.32, 1.75, 0.35, label, size=11, align="center")
        if i < len(steps) - 1:
            add_arrow(slide, x + 1.95, y + 0.38, 0.2, 0.25, Theme.blue)

    add_card(slide, 1.0, 4.0, 5.5, 2.75, Theme.white)
    add_text(slide, 1.25, 4.25, 5.0, 0.35, "核心代码入口", size=15, bold=True, color=Theme.blue)
    add_bullets(
        slide,
        1.25,
        4.7,
        5.0,
        1.9,
        [
            "runtime/langgraph_runtime.py",
            "runtime/langgraph/builder.py",
            "runtime/langgraph/nodes/agents.py",
            "runtime/langgraph/nodes/supervisor.py",
        ],
        size=12,
    )

    add_card(slide, 6.75, 4.0, 5.0, 2.75, Theme.light_blue)
    add_text(slide, 7.0, 4.25, 4.5, 0.35, "工程约束", size=15, bold=True, color=Theme.blue)
    add_bullets(
        slide,
        7.0,
        4.7,
        4.5,
        1.9,
        ["主 Agent 命令先行", "每个阶段有终态", "异常可降级可恢复", "全过程事件可追踪"],
        size=12,
    )


def slide_11_skill_tool(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "本项目案例 2：Skill + Tool + 审计闭环", "Controlled capabilities", 11)

    add_card(slide, 0.75, 1.25, 11.9, 5.8, Theme.white)

    add_text(slide, 1.0, 1.55, 11.2, 0.35, "调用链：Command -> Gate -> Tool/Skill -> LLM -> Audit", size=18, bold=True, color=Theme.blue)

    chain = [
        ("Command", 1.05, Theme.light_blue),
        ("Gate", 3.1, Theme.light_cyan),
        ("Tool/Skill", 5.15, Theme.light_green),
        ("LLM", 7.2, Theme.light_orange),
        ("Audit", 9.25, Theme.light_red),
    ]
    for i, (name, x, fill) in enumerate(chain):
        add_card(slide, x, 2.35, 1.8, 1.0, fill)
        add_text(slide, x + 0.15, 2.66, 1.5, 0.3, name, size=12, bold=True, align="center")
        if i < len(chain) - 1:
            add_arrow(slide, x + 1.8, 2.73, 0.2, 0.2)

    add_card(slide, 1.05, 3.8, 5.45, 2.95, Theme.white)
    add_text(slide, 1.25, 4.05, 5.0, 0.3, "可配置能力", size=15, bold=True, color=Theme.blue)
    add_bullets(
        slide,
        1.25,
        4.45,
        5.0,
        2.1,
        [
            "CodeAgent：Git 仓库检索（本地/远程）",
            "LogAgent：日志文件读取",
            "DomainAgent：责任田 Excel 查询",
            "DatabaseAgent：PostgreSQL 元数据查询",
        ],
        size=12,
    )

    add_card(slide, 6.75, 3.8, 5.2, 2.95, Theme.light_green)
    add_text(slide, 6.95, 4.05, 4.8, 0.3, "审计字段（前后端可见）", size=15, bold=True, color=Theme.green)
    add_bullets(
        slide,
        6.95,
        4.45,
        4.8,
        2.1,
        ["command_gate 判定", "请求摘要/返回摘要", "状态 + 耗时 + 错误", "skill 命中来源与注入摘要"],
        size=12,
    )


def slide_12_resilience(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "本项目案例 3：韧性设计", "Timeout, degrade, recover", 12)

    add_card(slide, 0.7, 1.2, 5.7, 5.9, Theme.white)
    add_text(slide, 0.95, 1.55, 5.2, 0.35, "状态机（会话不允许长期 pending）", size=16, bold=True, color=Theme.blue)
    states = [
        ("RUNNING", 1.0, 2.2, Theme.light_blue),
        ("RETRY", 2.4, 3.0, Theme.light_cyan),
        ("DEGRADED", 3.8, 3.8, Theme.light_orange),
        ("COMPLETED", 1.8, 5.0, Theme.light_green),
        ("FAILED", 4.2, 5.0, Theme.light_red),
    ]
    for text, x, y, fill in states:
        add_card(slide, x, y, 1.5, 0.75, fill)
        add_text(slide, x + 0.05, y + 0.24, 1.4, 0.25, text, size=11, bold=True, align="center")
    add_arrow(slide, 2.05, 2.52, 0.3, 0.2)
    add_arrow(slide, 3.45, 3.32, 0.3, 0.2)
    add_arrow(slide, 2.8, 5.02, 1.1, 0.2)
    add_arrow(slide, 4.55, 4.2, 0.2, 0.6)
    add_arrow(slide, 2.55, 3.75, 0.2, 1.0)

    add_card(slide, 0.95, 6.2, 5.2, 0.75, Theme.light_orange)
    add_text(slide, 1.15, 6.42, 4.8, 0.26, "策略：超时切换、局部重试、降级结论、恢复重跑。", size=12)

    add_card(slide, 6.65, 1.2, 6.0, 5.9, Theme.white)
    add_chart(
        slide,
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        x=6.95,
        y=1.65,
        w=5.4,
        h=3.7,
        categories=["连接池耗尽", "锁竞争", "网关重试风暴", "发布配置缺陷", "下游超时扩散"],
        series=[("置信度", [41, 24, 17, 11, 7])],
    )
    add_text(slide, 6.95, 5.78, 5.3, 0.35, "示意：输出 Top-K 候选，避免单点误判", size=12, color=Theme.subtext)


def slide_13_2week(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "团队落地模板（2 周）", "From idea to operational practice", 13)

    add_card(slide, 0.7, 1.2, 5.8, 5.95, Theme.white)
    add_text(slide, 1.0, 1.55, 5.2, 0.35, "Week 1：打基础", size=17, bold=True, color=Theme.blue)
    add_bullets(
        slide,
        1.0,
        2.0,
        5.2,
        2.7,
        [
            "定义 AGENTS.md：边界、流程、约束",
            "统一工作流：Explore -> Plan -> Build -> Verify",
            "接入最小测试与 PR 门禁",
            "建立关键日志与事件追踪",
        ],
        size=13,
    )
    add_card(slide, 1.0, 4.9, 5.2, 2.0, Theme.light_blue)
    add_text(slide, 1.2, 5.2, 4.8, 1.4, "验收：可重复执行一次端到端任务，\n全链路有日志、结果可回放。", size=12)

    add_card(slide, 6.85, 1.2, 5.8, 5.95, Theme.white)
    add_text(slide, 7.15, 1.55, 5.2, 0.35, "Week 2：补治理", size=17, bold=True, color=Theme.green)
    add_bullets(
        slide,
        7.15,
        2.0,
        5.2,
        2.7,
        [
            "工具调用审计与权限控制",
            "超时重试、失败降级、终态保证",
            "Benchmark 指标纳入 CI Gate",
            "一次真实故障演练与复盘",
        ],
        size=13,
    )
    add_card(slide, 7.15, 4.9, 5.2, 2.0, Theme.light_green)
    add_text(slide, 7.35, 5.2, 4.8, 1.4, "验收：无长期 pending；\n不达标指标可自动阻断发布。", size=12)


def slide_14_roadmap(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "90 天路线图", "P0 -> P3 capability maturity", 14)

    add_card(slide, 0.7, 1.2, 6.0, 5.95, Theme.white)
    add_chart(
        slide,
        chart_type=XL_CHART_TYPE.LINE_MARKERS,
        x=0.95,
        y=1.65,
        w=5.45,
        h=3.7,
        categories=["P0", "P1", "P2", "P3"],
        series=[("成熟度指数", [35, 56, 72, 85])],
    )
    add_text(slide, 0.95, 5.75, 5.3, 0.35, "成熟度曲线示意：先可用，再准确，再自治，再学习", size=12, color=Theme.subtext)

    add_card(slide, 6.95, 1.2, 5.75, 5.95, Theme.white)
    rows = [
        ("P0", "可用性", "不 pending，可观测，可回放", Theme.blue),
        ("P1", "准确性", "跨源证据 + Top-K 根因排序", Theme.cyan),
        ("P2", "可控修复", "审批、回滚、No-Regression Gate", Theme.green),
        ("P3", "持续学习", "反馈闭环、A/B 评测、策略中心", Theme.orange),
    ]
    y = 1.6
    for p, name, desc, c in rows:
        add_card(slide, 7.2, y, 5.2, 1.15, Theme.white)
        add_text(slide, 7.45, y + 0.18, 0.6, 0.3, p, size=12, bold=True, color=c)
        add_text(slide, 8.1, y + 0.18, 1.6, 0.3, name, size=13, bold=True, color=Theme.text)
        add_text(slide, 9.8, y + 0.18, 2.2, 0.5, desc, size=11, color=Theme.subtext)
        y += 1.35


def slide_15_close(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "结语与行动", "Execution checklist", 15)

    add_card(slide, 0.8, 1.4, 12.0, 2.0, Theme.light_blue)
    add_text(slide, 1.1, 1.85, 11.3, 0.55, "结论：Vibe Coding 的终局不是“自动写代码”，而是“可控地持续交付正确结果”。", size=21, bold=True, color=Theme.navy)

    add_card(slide, 0.8, 3.8, 12.0, 3.25, Theme.white)
    add_text(slide, 1.1, 4.1, 11.2, 0.35, "本周 5 条可执行动作", size=17, bold=True, color=Theme.blue)
    add_bullets(
        slide,
        1.1,
        4.55,
        11.0,
        2.2,
        [
            "固化 AGENTS.md 与流程约束（做什么/不做什么）",
            "关键链路接入审计：命令、工具、结论、耗时",
            "建立最小 benchmark，并加入 CI 阻断回归",
            "把工具接入改为开关化、可回滚、可权限控制",
            "选一个真实故障演练，做一次端到端复盘",
        ],
        size=14,
    )


def write_readme() -> None:
    OUT_README.write_text(
        "\n".join(
            [
                "# Vibe Coding PPT Deliverables",
                "",
                f"- PPTX: `{OUT_PPTX}`",
                f"- Assets: `{ASSET_DIR}`",
                "",
                "## Notes",
                "- Deck is generated from the 2026-03-06 outline.",
                "- Visual slides include diagrams, flow boxes, and data charts.",
                "- Focus ratio: industry practices (80%) + project case (20%).",
            ]
        ),
        encoding="utf-8",
    )


def main() -> None:
    ensure_dirs()

    prs = new_prs()
    slide_1_cover(prs)
    slide_2_why(prs)
    slide_3_harness(prs)
    slide_4_context(prs)
    slide_5_workflow(prs)
    slide_6_gate(prs)
    slide_7_mcp(prs)
    slide_8_multi_agent(prs)
    slide_9_antipattern(prs)
    slide_10_project_harness(prs)
    slide_11_skill_tool(prs)
    slide_12_resilience(prs)
    slide_13_2week(prs)
    slide_14_roadmap(prs)
    slide_15_close(prs)

    prs.save(OUT_PPTX)
    write_readme()
    print(f"Generated: {OUT_PPTX}")


if __name__ == "__main__":
    main()
