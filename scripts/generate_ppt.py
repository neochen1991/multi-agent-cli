#!/usr/bin/env python3
"""
SRE Debate Platform PPT Generator
生成面向领导的项目介绍PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    """创建PPT演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 定义颜色方案
    TITLE_COLOR = RGBColor(0, 51, 102)  # 深蓝色
    SUBTITLE_COLOR = RGBColor(51, 102, 153)  # 中蓝色
    ACCENT_COLOR = RGBColor(0, 102, 204)  # 亮蓝色
    TEXT_COLOR = RGBColor(51, 51, 51)  # 深灰色

    # ==================== 第1页：封面 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 添加背景色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 248, 255)  # 淡蓝色背景
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SRE Debate Platform"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "多模型辩论式 SRE 智能体平台"
    p.font.size = Pt(32)
    p.font.color.rgb = SUBTITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 底部信息
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "基于 AutoGen 多Agent编排构建"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # ==================== 第2页：项目概述 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "项目概述", TITLE_COLOR)

    # 左侧：项目定位
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "项目定位"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    p = tf.add_paragraph()
    p.text = "基于 AutoGen 多Agent编排构建的多模型辩论式SRE智能体平台"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_COLOR
    p.space_before = Pt(12)

    # 右侧：核心目标
    right_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.8), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "核心目标"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    goals = [
        ("三态资产融合", "打通设计态、开发态、运行态资产壁垒"),
        ("AI专家委员会", "多角色协作决策，提升分析准确性"),
        ("内部辩论机制", "质疑-反驳-裁决提升结论可信度"),
        ("智能修复建议", "自动生成修复方案与影响分析")
    ]

    for title, desc in goals:
        p = tf.add_paragraph()
        p.text = f"• {title}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = SUBTITLE_COLOR
        p.space_before = Pt(16)

        p = tf.add_paragraph()
        p.text = f"  {desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR

    # ==================== 第3页：核心价值与亮点 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "核心价值与亮点", TITLE_COLOR)

    # 创新亮点表格
    highlights_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
    tf = highlights_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "创新亮点"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    highlights = [
        ("三态资产融合", "统一建模运行态日志、开发态代码、设计态文档"),
        ("多模型专家协作", "6个专业Agent分工协作，模拟专家委员会决策"),
        ("AI辩论机制", "质疑-反驳-裁决四阶段辩论，避免单一模型偏见"),
        ("可追溯证据链", "构建完整的根因证据链，结论可解释、可追溯")
    ]

    for title, desc in highlights:
        p = tf.add_paragraph()
        p.text = f"● {title}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = SUBTITLE_COLOR
        p.space_before = Pt(14)

        p = tf.add_paragraph()
        p.text = f"   {desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR

    # 业务价值
    value_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.8), Inches(5.5))
    tf = value_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "业务价值"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    values = [
        ("缩短故障定位时间", "自动化分析替代人工排查"),
        ("提升分析准确性", "多角色辩论降低误判率"),
        ("知识沉淀复用", "案例库持续积累，经验可传承"),
        ("降低SRE技能门槛", "AI辅助决策，降低专家依赖")
    ]

    for title, desc in values:
        p = tf.add_paragraph()
        p.text = f"★ {title}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = SUBTITLE_COLOR
        p.space_before = Pt(14)

        p = tf.add_paragraph()
        p.text = f"   {desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR

    # ==================== 第4页：系统架构图（详细） ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "系统架构图", TITLE_COLOR)

    # 左侧用户层
    user_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.8),
        Inches(1.5), Inches(2.5)
    )
    user_box.fill.solid()
    user_box.fill.fore_color.rgb = RGBColor(255, 248, 220)
    user_box.line.color.rgb = RGBColor(200, 150, 50)
    tf = user_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "用户"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "SRE工程师\n开发人员\n运维人员"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 箭头1
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1.85), Inches(2.8), Inches(0.3), Inches(0.3))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = ACCENT_COLOR
    arrow1.line.fill.background()

    # 前端层
    frontend_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.2), Inches(1.6),
        Inches(2), Inches(3)
    )
    frontend_box.fill.solid()
    frontend_box.fill.fore_color.rgb = RGBColor(230, 247, 255)
    frontend_box.line.color.rgb = ACCENT_COLOR
    tf = frontend_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "前端层"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "React 18\nTypeScript\nAnt Design\nWebSocket"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 箭头2
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.25), Inches(2.8), Inches(0.3), Inches(0.3))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = ACCENT_COLOR
    arrow2.line.fill.background()

    # API网关层
    api_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.6), Inches(1.6),
        Inches(2), Inches(3)
    )
    api_box.fill.solid()
    api_box.fill.fore_color.rgb = RGBColor(200, 235, 255)
    api_box.line.color.rgb = ACCENT_COLOR
    tf = api_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "API网关"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "FastAPI\n鉴权/限流\n路由分发\nWebSocket"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 箭头3
    arrow3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.65), Inches(2.8), Inches(0.3), Inches(0.3))
    arrow3.fill.solid()
    arrow3.fill.fore_color.rgb = ACCENT_COLOR
    arrow3.line.fill.background()

    # 编排层
    flow_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7), Inches(1.6),
        Inches(2.2), Inches(3)
    )
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(170, 220, 255)
    flow_box.line.color.rgb = ACCENT_COLOR
    tf = flow_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Flow编排层"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "LangGraph\n辩论协调器\n上下文管理"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 箭头4
    arrow4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.25), Inches(2.8), Inches(0.3), Inches(0.3))
    arrow4.fill.solid()
    arrow4.fill.fore_color.rgb = ACCENT_COLOR
    arrow4.line.fill.background()

    # Agent层
    agent_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.6), Inches(1.6),
        Inches(3.2), Inches(3)
    )
    agent_box.fill.solid()
    agent_box.fill.fore_color.rgb = RGBColor(140, 210, 255)
    agent_box.line.color.rgb = ACCENT_COLOR
    tf = agent_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "多Agent协作层"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "LogAgent | DomainAgent\nCodeAgent | CriticAgent\nRebuttalAgent | JudgeAgent"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # LLM服务
    llm_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.3), Inches(4.8),
        Inches(2.2), Inches(1.5)
    )
    llm_box.fill.solid()
    llm_box.fill.fore_color.rgb = RGBColor(255, 230, 200)
    llm_box.line.color.rgb = RGBColor(200, 150, 50)
    tf = llm_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "LLM服务"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "glm-5\nOpenAI兼容网关"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 向下箭头到LLM
    arrow_down = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(11.1), Inches(4.65), Inches(0.25), Inches(0.25))
    arrow_down.fill.solid()
    arrow_down.fill.fore_color.rgb = RGBColor(200, 150, 50)
    arrow_down.line.fill.background()

    # 工具层
    tool_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(5.2),
        Inches(5.5), Inches(1.8)
    )
    tool_box.fill.solid()
    tool_box.fill.fore_color.rgb = RGBColor(230, 255, 230)
    tool_box.line.color.rgb = RGBColor(100, 180, 100)
    tf = tool_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "工具层"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "日志解析器 | Git工具 | DDD分析器 | 案例库检索 | 资产融合服务"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 存储层
    storage_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6), Inches(5.2),
        Inches(6.8), Inches(1.8)
    )
    storage_box.fill.solid()
    storage_box.fill.fore_color.rgb = RGBColor(240, 240, 255)
    storage_box.line.color.rgb = RGBColor(100, 100, 200)
    tf = storage_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "存储层"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "PostgreSQL | Redis | Neo4j图数据库 | 本地文件存储"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # ==================== 第5页：系统架构层次 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "系统架构层次", TITLE_COLOR)

    # 架构层次
    layers = [
        ("交互与接口层", "Web UI / REST API / CLI 工具", RGBColor(230, 247, 255)),
        ("Flow 编排层", "SRE Debate Flow (AutoGen/LangGraph)", RGBColor(200, 235, 255)),
        ("多模型专家协作层", "LogAgent | DomainAgent | CodeAgent | CriticAgent | RebuttalAgent | JudgeAgent", RGBColor(170, 220, 255)),
        ("工具层", "日志解析器 | Git工具 | DDD分析工具 | 案例库检索", RGBColor(140, 210, 255)),
        ("资产与存储层", "运行态资产 | 开发态资产 | 设计态资产 | 案例库", RGBColor(110, 195, 255))
    ]

    y_start = 1.8
    for i, (title, content, color) in enumerate(layers):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(y_start + i * 1.1),
            Inches(10.333), Inches(0.9)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0, 102, 204)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = content
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER

    # ==================== 第6页：系统运行示例 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "系统运行示例：订单服务故障分析", TITLE_COLOR)

    # 示例场景说明
    scenario_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.5),
        Inches(12.7), Inches(1.2)
    )
    scenario_box.fill.solid()
    scenario_box.fill.fore_color.rgb = RGBColor(255, 245, 230)
    scenario_box.line.color.rgb = RGBColor(200, 150, 50)
    tf = scenario_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "故障场景：用户反馈下单失败，日志显示 NullPointerException，涉及订单服务和库存服务"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(150, 100, 30)
    p.alignment = PP_ALIGN.CENTER

    # 步骤流程
    steps = [
        ("Step 1", "用户输入", "提交故障日志\n描述故障现象", RGBColor(230, 247, 255)),
        ("Step 2", "资产采集", "自动收集相关日志\n代码库、设计文档", RGBColor(200, 235, 255)),
        ("Step 3", "Agent辩论", "多Agent协作分析\n质疑-反驳-裁决", RGBColor(170, 220, 255)),
        ("Step 4", "生成报告", "根因分析+修复建议\n影响范围评估", RGBColor(140, 210, 255))
    ]

    x_start = 0.4
    for i, (step, name, desc, color) in enumerate(steps):
        x = x_start + i * 3.2
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2.9),
            Inches(2.9), Inches(1.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = ACCENT_COLOR
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER

        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.95), Inches(3.5), Inches(0.25), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_COLOR
            arrow.line.fill.background()

    # Agent辩论详情
    debate_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(4.7),
        Inches(6.2), Inches(2.5)
    )
    debate_box.fill.solid()
    debate_box.fill.fore_color.rgb = RGBColor(245, 250, 255)
    debate_box.line.color.rgb = ACCENT_COLOR
    tf = debate_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Agent辩论过程"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    debates = [
        ("LogAgent:", "识别NullPointerException，定位到OrderService.createOrder方法"),
        ("DomainAgent:", "映射到【订单域-订单聚合根】，影响库存服务"),
        ("CodeAgent:", "根因假设：库存扣减返回null未校验，置信度85%"),
        ("CriticAgent:", "质疑：未考虑并发场景，DDD边界可能误判"),
        ("RebuttalAgent:", "补充证据：日志显示单线程场景，边界判断正确"),
        ("JudgeAgent:", "最终裁决：确认根因，建议添加空值校验")
    ]
    for agent, content in debates:
        p = tf.add_paragraph()
        p.text = f"{agent} {content}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_COLOR

    # 输出结果
    result_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.7), Inches(4.7),
        Inches(6.2), Inches(2.5)
    )
    result_box.fill.solid()
    result_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    result_box.line.color.rgb = RGBColor(100, 180, 100)
    tf = result_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "分析报告输出"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(50, 150, 50)

    results = [
        "根因：InventoryService.deduct() 返回null未校验",
        "证据链：日志 → 代码 → 设计文档完整追溯",
        "影响范围：订单服务、库存服务",
        "修复建议：添加Optional处理或空值校验",
        "置信度：92%",
        "相似案例：INC-2024-0125（相似度89%）"
    ]
    for r in results:
        p = tf.add_paragraph()
        p.text = f"✓ {r}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_COLOR

    # ==================== 第7页：AI专家委员会 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "AI专家委员会分工", TITLE_COLOR)

    agents = [
        ("LogAgent", "日志分析专家", "解析运行态日志，提取异常栈、URL、类路径等关键信息"),
        ("DomainAgent", "领域映射专家", "将异常映射到领域模型，识别聚合根和限界上下文"),
        ("CodeAgent", "代码分析专家", "分析代码层面根因，构建证据链，提出修复建议"),
        ("CriticAgent", "架构质疑专家", "检查DDD原则违反，提出质疑意见"),
        ("RebuttalAgent", "技术反驳专家", "回应质疑，修正分析结论"),
        ("JudgeAgent", "技术委员会主席", "综合辩论过程，给出最终裁决")
    ]

    y_start = 1.6
    for i, (name, role, desc) in enumerate(agents):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 6.4
        y = y_start + row * 1.9

        # Agent卡片
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(6), Inches(1.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(245, 250, 255)
        shape.line.color.rgb = ACCENT_COLOR

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{name}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR

        p = tf.add_paragraph()
        p.text = role
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = SUBTITLE_COLOR

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR

    # ==================== 第8页：辩论流程 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "四阶段辩论流程", TITLE_COLOR)

    stages = [
        ("第一阶段", "独立分析", "CodeAgent", "提出根因假设\n构建证据链", RGBColor(200, 230, 255)),
        ("第二阶段", "交叉质疑", "CriticAgent", "检查DDD违反\n提出质疑意见", RGBColor(255, 230, 200)),
        ("第三阶段", "反驳修正", "RebuttalAgent", "回应质疑修正\n补充证据支持", RGBColor(200, 255, 200)),
        ("第四阶段", "最终裁决", "JudgeAgent", "综合裁决结论\n生成最终报告", RGBColor(230, 200, 255))
    ]

    x_start = 0.5
    for i, (phase, name, agent, desc, color) in enumerate(stages):
        x = x_start + i * 3.1

        # 阶段框
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2),
            Inches(2.9), Inches(3.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = ACCENT_COLOR

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = name
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(8)

        p = tf.add_paragraph()
        p.text = agent
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(12)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(8)

        # 箭头
        if i < 3:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + 2.95), Inches(3.5),
                Inches(0.25), Inches(0.5)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_COLOR
            arrow.line.fill.background()

    # 优势说明
    advantage_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.2))
    tf = advantage_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "辩论优势：避免单一偏见 | 提升结论可信度 | 可解释性强 | 完整辩论过程可追溯"
    p.font.size = Pt(18)
    p.font.color.rgb = SUBTITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # ==================== 第9页：三态资产融合 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "三态资产融合", TITLE_COLOR)

    assets = [
        ("运行态", "生产环境", "应用日志 | JVM监控 | Trace信息 | 慢SQL", RGBColor(255, 230, 230)),
        ("开发态", "代码仓库", "Git仓库 | 聚合根 | Controller | Repository", RGBColor(230, 255, 230)),
        ("设计态", "设计文档", "领域模型 | 接口设计 | 数据库设计 | 历史案例", RGBColor(230, 230, 255))
    ]

    x_start = 0.8
    for i, (name, source, content, color) in enumerate(assets):
        x = x_start + i * 4.1

        # 资产卡片
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2),
            Inches(3.8), Inches(3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = ACCENT_COLOR

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{name}资产"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = f"来源: {source}"
        p.font.size = Pt(16)
        p.font.color.rgb = SUBTITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(12)

        p = tf.add_paragraph()
        p.text = content
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(16)

    # 融合价值
    value_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1.5))
    tf = value_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "融合价值"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER

    values = ["跨态关联 - 从日志异常追溯到代码实现再到设计方案",
              "全景视角 - 打破信息孤岛，构建完整知识图谱",
              "精准定位 - 结合多态信息，提升根因分析准确性"]

    for v in values:
        p = tf.add_paragraph()
        p.text = f"✓ {v}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER

    # ==================== 第10页：技术选型 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "技术选型", TITLE_COLOR)

    # 后端技术
    backend_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(2.5))
    tf = backend_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "后端技术"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    backend_tech = ["Python 3.11+ / FastAPI", "AutoGen Runtime / LangGraph", "PostgreSQL + Neo4j", "Redis + Celery"]
    for tech in backend_tech:
        p = tf.add_paragraph()
        p.text = f"• {tech}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)

    # 前端技术
    frontend_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(2.5))
    tf = frontend_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "前端技术"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    frontend_tech = ["React 18 + TypeScript", "Ant Design 5", "Vite 构建工具", "Zustand 状态管理"]
    for tech in frontend_tech:
        p = tf.add_paragraph()
        p.text = f"• {tech}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)

    # AI模型
    ai_box = slide.shapes.add_textbox(Inches(9.5), Inches(1.5), Inches(3.5), Inches(2.5))
    tf = ai_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "AI模型"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    ai_tech = ["主力模型: glm-5", "模型服务: OpenAI兼容网关", "多Agent编排框架"]
    for tech in ai_tech:
        p = tf.add_paragraph()
        p.text = f"• {tech}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)

    # 技术亮点
    highlight_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(2.5))
    tf = highlight_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "技术亮点"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    highlights = [
        "微服务架构 - 前后端分离，支持水平扩展",
        "图数据库 - Neo4j存储资产关联关系",
        "实时通信 - WebSocket支持实时辩论流",
        "容器化部署 - Docker支持，K8s就绪"
    ]

    for h in highlights:
        p = tf.add_paragraph()
        p.text = f"★ {h}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)

    # ==================== 第11页：已实现能力 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "已实现能力", TITLE_COLOR)

    # 核心功能
    func_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
    tf = func_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "核心功能"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    features = [
        "Incident 全流程（创建 → 会话 → 辩论 → 报告）",
        "WebSocket 实时辩论流",
        "资产融合查询 API",
        "历史记录与资产图谱页面",
        "可选鉴权（JWT/RBAC）",
        "限流、熔断、指标监控"
    ]

    for f in features:
        p = tf.add_paragraph()
        p.text = f"✓ {f}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(12)

    # 部署能力
    deploy_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.8), Inches(5))
    tf = deploy_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "部署能力"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    deploy_features = [
        "一键启动脚本（前后端并行启动）",
        "Docker 容器化部署支持",
        "完善的API文档（Swagger/ReDoc）",
        "指标监控端点 (/metrics)"
    ]

    for f in deploy_features:
        p = tf.add_paragraph()
        p.text = f"✓ {f}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(12)

    # API文档地址
    api_box = slide.shapes.add_textbox(Inches(7), Inches(5), Inches(5.8), Inches(1.5))
    tf = api_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "API文档地址"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SUBTITLE_COLOR

    p = tf.add_paragraph()
    p.text = "Swagger: http://localhost:8000/docs"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_COLOR

    p = tf.add_paragraph()
    p.text = "ReDoc: http://localhost:8000/redoc"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_COLOR

    # ==================== 第12页：未来规划 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "未来规划", TITLE_COLOR)

    # 短期目标
    short_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
    tf = short_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "短期目标"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    short_goals = [
        "数据库持久化完善",
        "WebSocket实时通信优化",
        "案例库深度集成",
        "测试覆盖率提升"
    ]

    for g in short_goals:
        p = tf.add_paragraph()
        p.text = f"○ {g}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(14)

    # 中长期目标
    long_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.8), Inches(5))
    tf = long_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "中长期目标"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    long_goals = [
        "生产环境部署",
        "多模型路由支持",
        "自动修复能力（PR生成）",
        "灰度发布建议",
        "知识图谱可视化"
    ]

    for g in long_goals:
        p = tf.add_paragraph()
        p.text = f"○ {g}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(14)

    # ==================== 第13页：总结 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "项目价值总结", TITLE_COLOR)

    # 对比表格
    # 传统方式
    trad_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.8),
        Inches(5.8), Inches(4)
    )
    trad_box.fill.solid()
    trad_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
    trad_box.line.color.rgb = RGBColor(200, 100, 100)

    tf = trad_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "传统故障分析"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 50, 50)
    p.alignment = PP_ALIGN.CENTER

    trad_items = ["依赖专家经验", "单一视角分析", "知识难以沉淀", "分析周期长", "结论难追溯"]
    for item in trad_items:
        p = tf.add_paragraph()
        p.text = f"✗ {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(10)
        p.alignment = PP_ALIGN.CENTER

    # AI方式
    ai_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7), Inches(1.8),
        Inches(5.8), Inches(4)
    )
    ai_box.fill.solid()
    ai_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    ai_box.line.color.rgb = RGBColor(100, 200, 100)

    tf = ai_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI辩论式故障分析"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(50, 150, 50)
    p.alignment = PP_ALIGN.CENTER

    ai_items = ["AI辅助决策", "多角色协作辩论", "案例库持续积累", "自动化快速定位", "证据链完整可解释"]
    for item in ai_items:
        p = tf.add_paragraph()
        p.text = f"✓ {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(10)
        p.alignment = PP_ALIGN.CENTER

    # 核心竞争力
    comp_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12.333), Inches(1.2))
    tf = comp_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "核心竞争力"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "创新性：业界首创AI辩论式故障分析机制  |  实用性：覆盖故障分析全流程  |  扩展性：模块化设计，易于扩展"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    # ==================== 第14页：感谢页 ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加背景色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 248, 255)
    shape.line.fill.background()

    # 感谢文字
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢！"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 联系方式
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(2))
    tf = contact_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "项目仓库"
    p.font.size = Pt(20)
    p.font.color.rgb = SUBTITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "github.com/neochen1991/multi-agent-cli"
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = "技术文档: plans/sre-debate-platform-architecture.md"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(16)

    return prs


def add_title(slide, text, color):
    """添加标题"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color

    # 添加下划线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.3),
        Inches(12.333), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 102, 204)
    line.line.fill.background()


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/Users/neochen/multi-agent-cli_v2/plans/SRE_Debate_Platform_介绍.pptx"
    prs.save(output_path)
    print(f"PPT已生成: {output_path}")