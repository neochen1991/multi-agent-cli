#!/usr/bin/env python3
"""generate项目介绍WarRoomPPT脚本。"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output" / "project-intro-warroom-ppt"
OUT_PPTX = OUT_DIR / "project-intro-warroom.pptx"
OUT_SLIDES = OUT_DIR / "slides.md"
OUT_NOTES = OUT_DIR / "notes.md"
OUT_REFS = OUT_DIR / "refs.md"
OUT_README = OUT_DIR / "README.md"


class Theme:
    """封装Theme相关常量或数据结构。"""
    
    bg = RGBColor(8, 14, 24)
    panel = RGBColor(17, 27, 43)
    panel_2 = RGBColor(22, 35, 56)
    border = RGBColor(44, 62, 92)
    text = RGBColor(237, 241, 247)
    text_dim = RGBColor(160, 174, 196)
    cyan = RGBColor(78, 205, 196)
    amber = RGBColor(255, 191, 105)
    red = RGBColor(255, 107, 107)
    green = RGBColor(92, 220, 151)
    blue = RGBColor(96, 165, 250)
    white = RGBColor(255, 255, 255)
    ink = RGBColor(19, 28, 43)
    card_warm = RGBColor(255, 247, 237)
    card_mint = RGBColor(237, 251, 245)
    card_cyan = RGBColor(234, 249, 250)
    card_soft = RGBColor(242, 246, 252)


def ensure_dirs() -> None:
    """确保dirs相关前置条件已经满足。"""
    
    OUT_DIR.mkdir(parents=True, exist_ok=True)


def new_prs() -> Presentation:
    """执行新增prs相关逻辑。"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_bg(slide, prs: Presentation, *, alt: bool = False) -> None:
    """向当前页补充背景相关元素，并统一样式与布局。"""
    
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = Theme.panel if alt else Theme.bg
    bg.line.fill.background()


def add_topbar(slide, title: str, subtitle: str, page: int) -> None:
    """向当前页补充topbar相关元素，并统一样式与布局。"""
    
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.78)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = Theme.panel_2
    band.line.fill.background()
    add_text(slide, 0.45, 0.16, 8.7, 0.35, title, size=24, bold=True, color=Theme.text)
    add_text(slide, 9.0, 0.2, 2.8, 0.24, subtitle, size=10, color=Theme.text_dim, align="right")
    add_text(slide, 12.15, 0.18, 0.6, 0.24, str(page), size=12, bold=True, color=Theme.amber, align="right")


def add_footer(slide, source: str) -> None:
    """向当前页补充页脚相关元素，并统一样式与布局。"""
    
    add_text(slide, 0.45, 7.04, 12.1, 0.18, f"Source: {source}", size=8, color=Theme.text_dim)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: int = 16,
    bold: bool = False,
    color: RGBColor = Theme.text,
    align: str = "left",
) -> None:
    """向当前页补充文本相关元素，并统一样式与布局。"""
    
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: Iterable[str],
    *,
    size: int = 13,
    color: RGBColor = Theme.text,
) -> None:
    """向当前页补充bullets相关元素，并统一样式与布局。"""
    
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(size)
        p.font.color.rgb = color


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = Theme.panel,
    border: RGBColor = Theme.border,
) -> None:
    """向当前页补充卡片相关元素，并统一样式与布局。"""
    
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border


def add_pill(slide, x: float, y: float, w: float, text: str, *, fill: RGBColor, color: RGBColor) -> None:
    """向当前页补充标签相关元素，并统一样式与布局。"""
    
    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.42)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    add_text(slide, x + 0.08, y + 0.09, w - 0.16, 0.18, text, size=10, bold=True, color=color, align="center")


def add_metric(slide, x: float, y: float, w: float, title: str, value: str, *, accent: RGBColor) -> None:
    """向当前页补充指标相关元素，并统一样式与布局。"""
    
    add_card(slide, x, y, w, 1.02, fill=Theme.panel_2)
    add_text(slide, x + 0.16, y + 0.16, w - 0.32, 0.18, title, size=10, color=Theme.text_dim)
    add_text(slide, x + 0.16, y + 0.42, w - 0.32, 0.28, value, size=20, bold=True, color=accent)


def add_connector(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = Theme.text_dim) -> None:
    """向当前页补充连接器相关元素，并统一样式与布局。"""
    
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(1.4)


def add_box_label(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    *,
    fill: RGBColor,
    title_color: RGBColor = Theme.ink,
    body_color: RGBColor = Theme.ink,
) -> None:
    """向当前页补充boxlabel相关元素，并统一样式与布局。"""
    
    add_card(slide, x, y, w, h, fill=fill, border=Theme.border)
    add_text(slide, x + 0.16, y + 0.16, w - 0.32, 0.18, title, size=13, bold=True, color=title_color, align="center")
    add_text(slide, x + 0.16, y + 0.46, w - 0.32, h - 0.56, body, size=10, color=body_color, align="center")


def add_chart(
    slide,
    *,
    chart_type: XL_CHART_TYPE,
    x: float,
    y: float,
    w: float,
    h: float,
    categories: list[str],
    series: list[tuple[str, list[float]]],
) -> None:
    """向当前页补充图表相关元素，并统一样式与布局。"""
    
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    chart = slide.shapes.add_chart(chart_type, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)
        chart.legend.include_in_layout = False
    if chart.value_axis is not None:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.font.size = Pt(10)
    if chart.category_axis is not None:
        chart.category_axis.tick_labels.font.size = Pt(10)
    palette = [Theme.cyan, Theme.blue, Theme.green, Theme.amber, Theme.red]
    for idx, s in enumerate(chart.series):
        fill = s.format.fill
        fill.solid()
        fill.fore_color.rgb = palette[idx % len(palette)]


def slide_cover(prs: Presentation) -> None:
    """构建封面对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_text(slide, 0.58, 0.55, 5.2, 0.28, "WAR-ROOM ARCHITECTURE BRIEF", size=11, bold=True, color=Theme.amber)
    add_text(slide, 0.58, 1.0, 8.8, 0.9, "生产问题根因分析系统\n架构培训", size=28, bold=True, color=Theme.text)
    add_text(slide, 0.58, 2.18, 8.6, 0.3, "基于 code_wiki_v2 的团队内部架构讲解版", size=16, color=Theme.text_dim)
    add_pill(slide, 0.58, 3.0, 1.9, "LangGraph Runtime", fill=Theme.card_cyan, color=Theme.ink)
    add_pill(slide, 2.58, 3.0, 1.65, "Multi-Agent", fill=Theme.card_mint, color=Theme.ink)
    add_pill(slide, 4.35, 3.0, 1.55, "War-Room UX", fill=Theme.card_warm, color=Theme.ink)

    add_text(slide, 0.58, 3.85, 4.4, 0.2, "培训重点", size=12, bold=True, color=Theme.cyan)
    add_bullets(
        slide,
        0.58,
        4.15,
        5.1,
        1.8,
        [
            "系统目标、边界与总体分层",
            "端到端分析主流程和事件流",
            "LangGraph 多 Agent 运行时和工具治理",
            "前端工作台、可靠性和扩展路径",
        ],
        size=13,
    )

    add_card(slide, 7.35, 0.92, 5.3, 5.85, fill=Theme.panel_2)
    add_text(slide, 7.72, 1.18, 4.4, 0.24, "系统快照", size=12, bold=True, color=Theme.amber)
    add_metric(slide, 7.72, 1.55, 1.45, "UI", "10 pages", accent=Theme.cyan)
    add_metric(slide, 9.32, 1.55, 1.45, "Agents", "7+", accent=Theme.green)
    add_metric(slide, 10.92, 1.55, 1.45, "Mode", "WS", accent=Theme.blue)
    add_card(slide, 7.72, 2.82, 4.3, 1.0, fill=Theme.panel)
    add_text(slide, 7.95, 3.1, 3.8, 0.16, "ProblemAnalysisAgent", size=11, bold=True, color=Theme.cyan, align="center")
    for label, x, y, color in [
        ("Log", 7.82, 4.34, Theme.cyan),
        ("Code", 9.18, 4.34, Theme.blue),
        ("Domain", 10.54, 4.34, Theme.green),
        ("DB", 11.76, 4.34, Theme.amber),
    ]:
        add_card(slide, x, y, 0.98, 0.62, fill=Theme.panel)
        add_text(slide, x + 0.09, y + 0.18, 0.8, 0.14, label, size=10, bold=True, color=color, align="center")
        add_connector(slide, 9.88, 3.82, x + 0.49, y)
    add_text(slide, 7.82, 5.45, 4.2, 0.6, "用户输入 -> 资产映射 -> 命令分发 -> 专家分析 -> 裁决报告", size=12, color=Theme.text_dim)
    add_text(slide, 0.58, 6.82, 10.8, 0.18, f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')} Asia/Shanghai", size=9, color=Theme.text_dim)
    add_footer(slide, "docs/wiki/code_wiki_v2.md")


def slide_agenda(prs: Presentation, page: int) -> None:
    """构建议程对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "培训目录", "Architecture Roadmap", page)
    items = [
        "1. 系统目标与边界",
        "2. 总体分层架构",
        "3. 前端工作台架构",
        "4. 后端接入与服务流",
        "5. LangGraph 运行时",
        "6. 多 Agent 协作机制",
        "7. Tool / Skill / Connector",
        "8. 责任田资产与数据库链路",
        "9. 端到端分析流程",
        "10. 真实案例映射",
        "11. 事件流、结果与报告",
        "12. 报告结果视图",
        "13. 可靠性、治理与扩展",
        "14. 代码阅读地图",
    ]
    for i, item in enumerate(items):
        if i < 7:
            x = 0.9
            y = 1.18 + i * 0.76
            w = 5.35
        else:
            x = 6.95
            y = 1.18 + (i - 7) * 0.76
            w = 5.35
        add_card(slide, x, y, w, 0.62, fill=Theme.panel)
        add_text(slide, x + 0.18, y + 0.17, w - 0.36, 0.16, item, size=12, color=Theme.text)
    add_footer(slide, "docs/plans/2026-03-06-project-intro-ppt-design.md")


def slide_value(prs: Presentation, page: int) -> None:
    """构建价值对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "1) 系统目标与业务边界", "Value framing", page)
    add_text(slide, 0.55, 1.0, 8.5, 0.32, "结论：本项目不是聊天机器人，而是一套围绕生产故障调查流程构建的根因分析系统。", size=22, bold=True)
    add_card(slide, 0.55, 1.65, 3.75, 4.95, fill=Theme.panel)
    add_text(slide, 0.8, 1.95, 2.8, 0.2, "要解决的问题", size=14, bold=True, color=Theme.red)
    add_bullets(
        slide,
        0.8,
        2.25,
        3.05,
        3.8,
        [
            "故障证据分散在日志、代码、资产、数据库和监控",
            "跨团队协作成本高，责任归属和证据引用难统一",
            "传统排障过程不可回放，结论难复核",
            "高压值班场景要求快速拿到首批证据和可执行建议",
        ],
    )
    add_card(slide, 4.55, 1.65, 3.95, 4.95, fill=Theme.panel_2)
    add_text(slide, 4.8, 1.95, 2.9, 0.2, "系统边界", size=14, bold=True, color=Theme.cyan)
    add_bullets(
        slide,
        4.8,
        2.25,
        3.2,
        3.8,
        [
            "主 Agent 负责调度、收敛和报告归纳",
            "专家 Agent 负责日志、代码、领域、数据库等证据分析",
            "工具调用必须受命令门禁和开关控制",
            "默认使用本地文件 / memory / markdown，不强依赖外部存储",
        ],
    )
    add_card(slide, 8.8, 1.65, 4.0, 4.95, fill=Theme.panel)
    add_text(slide, 9.05, 1.95, 3.1, 0.2, "关键产出", size=14, bold=True, color=Theme.green)
    add_bullets(
        slide,
        9.05,
        2.25,
        3.2,
        3.8,
        [
            "责任田映射结果",
            "Agent 协作分析轨迹",
            "工具调用审计记录",
            "Top-K 根因候选与证据链",
            "验证建议、修复建议和回放能力",
        ],
    )
    add_footer(slide, "docs/wiki/code_wiki_v2.md#1-项目定位")


def slide_system_arch(prs: Presentation, page: int) -> None:
    """构建系统arch对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "2) 总体分层架构", "End-to-end system architecture", page)
    add_text(slide, 0.55, 1.0, 8.8, 0.26, "结论：系统采用“四层结构”，把用户体验、业务编排、Agent 运行时和工具治理明确分层。", size=22, bold=True)
    add_box_label(slide, 0.95, 1.75, 11.35, 0.82, "前端工作台层", "Home / Incident / WarRoom / Workbench / History / Assets", fill=Theme.card_cyan)
    add_box_label(slide, 1.55, 2.85, 10.15, 0.92, "API / Service / Flow 层", "FastAPI + DebateService + AssetService + ReportService + Flow", fill=Theme.card_soft)
    add_box_label(slide, 2.1, 4.05, 9.05, 1.0, "LangGraph 多 Agent 运行时", "GraphBuilder / Supervisor / Agent Nodes / State / EventDispatcher / Retry / Resume", fill=Theme.card_mint)
    add_box_label(slide, 2.75, 5.35, 7.75, 0.92, "Tool / Skill / Connector / Governance", "工具注册、Skill 注入、外部连接器、Benchmark、治理", fill=Theme.card_warm)
    for y1, y2 in [(2.57, 2.85), (3.77, 4.05), (5.05, 5.35)]:
        add_connector(slide, 6.6, y1, 6.6, y2, Theme.blue)
    add_text(slide, 0.98, 6.58, 11.0, 0.2, "这一分层的核心价值是：Model 与 System 分层，UI 与 Runtime 解耦，工具能力统一治理。", size=12, color=Theme.text_dim)
    add_footer(slide, "docs/wiki/code_wiki_v2.md#3-系统总体结构")


def slide_frontend(prs: Presentation, page: int) -> None:
    """构建前端对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "3) 前端工作台架构", "Frontend workspace", page)
    add_text(slide, 0.55, 1.0, 8.8, 0.26, "结论：前端不是单页表单，而是一套调查工作台，按“录入、观察、回放、治理”分场景组织。", size=22, bold=True)
    cards = [
        ("Home", 0.7, 1.8, 1.7, 0.92, Theme.panel),
        ("Incident", 2.6, 1.8, 1.7, 0.92, Theme.panel_2),
        ("WarRoom", 4.5, 1.8, 1.7, 0.92, Theme.panel),
        ("Workbench", 6.4, 1.8, 1.7, 0.92, Theme.panel_2),
        ("History", 8.3, 1.8, 1.7, 0.92, Theme.panel),
        ("Assets", 10.2, 1.8, 1.7, 0.92, Theme.panel_2),
    ]
    for title, x, y, w, h, fill in cards:
        add_card(slide, x, y, w, h, fill=fill)
        add_text(slide, x + 0.14, y + 0.28, w - 0.28, 0.16, title, size=14, bold=True, color=Theme.cyan, align="center")
    for x1, x2 in [(2.4, 2.6), (4.3, 4.5), (6.2, 6.4), (8.1, 8.3), (10.0, 10.2)]:
        add_connector(slide, x1, 2.25, x2, 2.25, Theme.blue)
    add_card(slide, 0.7, 3.05, 5.95, 3.0, fill=Theme.panel)
    add_text(slide, 0.95, 3.3, 2.8, 0.18, "工作台分工", size=14, bold=True, color=Theme.amber)
    add_bullets(
        slide,
        0.95,
        3.58,
        5.05,
        2.25,
        [
            "Home：能力总览、系统入口和 Agent 说明",
            "Incident：创建会话、连接 WebSocket、分标签展示资产映射/辩论过程/结果/报告",
            "WarRoom：聚焦时间线、证据链、关键结论和工具调用",
            "Workbench / History：调查回放、详情恢复、历史定位",
            "Assets / Settings / Benchmark / Governance：资产维护与平台治理",
        ],
    )
    add_card(slide, 6.9, 3.05, 5.55, 3.0, fill=Theme.panel_2)
    add_text(slide, 7.15, 3.3, 3.0, 0.18, "状态流与渲染", size=14, bold=True, color=Theme.green)
    add_box_label(slide, 7.1, 3.7, 1.45, 0.82, "WebSocket", "phase / chat / tool / result", fill=Theme.card_cyan)
    add_box_label(slide, 8.85, 3.7, 1.45, 0.82, "Event Store", "去重 / 排序 / 分组", fill=Theme.card_soft)
    add_box_label(slide, 10.6, 3.7, 1.45, 0.82, "UI Views", "资产 / 过程 / 结果 / 报告", fill=Theme.card_mint)
    add_connector(slide, 8.55, 4.1, 8.85, 4.1, Theme.blue)
    add_connector(slide, 10.3, 4.1, 10.6, 4.1, Theme.blue)
    add_bullets(
        slide,
        7.12,
        4.85,
        4.95,
        1.3,
        [
            "资产映射先于辩论可见，避免用户空等待",
            "Agent 对话、工具调用、阶段事件分轨渲染",
            "报告结果与实时过程分离，支持回放和重进详情页",
        ],
    )
    add_footer(slide, "frontend/src/pages/* + docs/wiki/code_wiki_v2.md#5-前端分层说明")


def slide_backend(prs: Presentation, page: int) -> None:
    """构建后端对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "4) 后端接入与服务流", "Backend service flow", page)
    add_text(slide, 0.55, 1.0, 8.7, 0.26, "结论：后端通过 API -> Service -> Flow -> Runtime 的链路，把业务上下文和运行时编排解耦。", size=22, bold=True)
    for label, x, fill in [
        ("FastAPI API", 0.9, Theme.card_cyan),
        ("DebateService", 3.3, Theme.card_mint),
        ("DebateFlow", 5.7, Theme.card_soft),
        ("LangGraphRuntime", 8.1, Theme.card_warm),
        ("Repositories / Reports", 10.45, Theme.card_cyan),
    ]:
        add_card(slide, x, 2.0, 2.0, 1.05, fill=fill)
        add_text(slide, x + 0.12, 2.35, 1.76, 0.18, label, size=13, bold=True, color=Theme.ink, align="center")
    for x1, x2 in [(2.9, 3.3), (5.3, 5.7), (7.7, 8.1), (10.1, 10.45)]:
        add_connector(slide, x1, 2.52, x2, 2.52, Theme.blue)
    add_card(slide, 0.95, 3.55, 5.6, 2.45, fill=Theme.panel)
    add_text(slide, 1.18, 3.82, 3.4, 0.18, "Service 层职责", size=14, bold=True, color=Theme.cyan)
    add_bullets(
        slide,
        1.18,
        4.1,
        4.95,
        1.7,
        [
            "读取 incident / session，整理业务上下文",
            "先完成责任田映射，再触发 runtime",
            "驱动结果持久化、报告生成和终态收敛",
        ],
    )
    add_card(slide, 6.85, 3.55, 5.5, 2.45, fill=Theme.panel_2)
    add_text(slide, 7.08, 3.82, 3.6, 0.18, "Flow 层意义", size=14, bold=True, color=Theme.amber)
    add_bullets(
        slide,
        7.08,
        4.1,
        4.8,
        1.7,
        [
            "把多个 Service 和 Runtime 串成一次完整调查旅程",
            "避免 API 层和 Runtime 直接耦合",
            "便于后续把自动告警触发、回放和评测纳入统一流程",
        ],
    )
    add_footer(slide, "backend/app/api router + services/debate_service.py + flows/debate_flow.py")


def slide_runtime(prs: Presentation, page: int) -> None:
    """构建运行时对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "5) LangGraph 运行时", "Runtime core", page)
    add_text(slide, 0.55, 1.0, 8.9, 0.26, "结论：运行时把图构建、状态、节点执行、事件分发和可靠性保护拆成了多个专用模块。", size=22, bold=True)
    add_box_label(slide, 5.0, 1.85, 3.3, 0.95, "Orchestrator", "langgraph_runtime.py\n统一协调图构建、状态初始化、终态收敛", fill=Theme.card_warm)
    add_box_label(slide, 1.0, 3.15, 2.55, 0.95, "GraphBuilder", "builder.py\n定义节点、边和路由", fill=Theme.panel)
    add_box_label(slide, 3.75, 3.15, 2.55, 0.95, "State / Messages", "state.py\n承载消息、路由、证据、输出", fill=Theme.panel_2, title_color=Theme.cyan, body_color=Theme.text_dim)
    add_box_label(slide, 6.5, 3.15, 2.55, 0.95, "AgentRunner", "agent_runner.py\n执行 Agent 节点与重试", fill=Theme.panel)
    add_box_label(slide, 9.25, 3.15, 2.55, 0.95, "EventDispatcher", "event_dispatcher.py\n把轨迹发给前端", fill=Theme.panel_2, title_color=Theme.cyan, body_color=Theme.text_dim)
    for x in [2.28, 5.03, 7.78, 10.53]:
        add_connector(slide, 6.65, 2.8, x, 3.15, Theme.blue)
    add_box_label(slide, 1.3, 4.65, 2.75, 0.95, "Phase Executor", "phase_executor.py\n控制阶段推进", fill=Theme.card_soft)
    add_box_label(slide, 4.35, 4.65, 2.75, 0.95, "Checkpointer", "checkpointer.py\n保存运行态，支持恢复", fill=Theme.card_mint)
    add_box_label(slide, 7.4, 4.65, 2.75, 0.95, "Doom Loop Guard", "doom_loop_guard.py\n防止循环提问与重复执行", fill=Theme.card_soft)
    add_box_label(slide, 10.45, 4.65, 1.55, 0.95, "Compaction", "session_compaction.py\n压缩历史上下文", fill=Theme.card_mint)
    add_card(slide, 0.9, 6.05, 11.4, 0.72, fill=Theme.panel)
    add_text(slide, 1.15, 6.28, 10.8, 0.14, "核心设计取舍：LLM 负责推理生成；系统负责状态机、路由、超时、重试、降级、审计和终态保证。", size=12, bold=True, color=Theme.green)
    add_footer(slide, "backend/app/runtime/langgraph/* + backend/app/runtime/langgraph_runtime.py")


def slide_agents(prs: Presentation, page: int) -> None:
    """构建Agent对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "6) 多 Agent 协作机制", "Main-agent driven collaboration", page)
    add_text(slide, 0.55, 1.0, 9.0, 0.26, "结论：系统不是自由聊天，而是“主 Agent 指挥 + 专家 Agent 执行 + Judge 收敛”的协作模式。", size=22, bold=True)
    add_card(slide, 5.15, 1.75, 3.0, 1.0, fill=Theme.card_warm)
    add_text(slide, 5.42, 2.05, 2.45, 0.18, "ProblemAnalysisAgent", size=15, bold=True, color=Theme.ink, align="center")
    add_text(slide, 5.42, 2.34, 2.45, 0.14, "任务拆解 / 命令分发 / 过程协调", size=10, color=Theme.ink, align="center")
    agents = [
        ("LogAgent", 1.2, 3.4, Theme.cyan),
        ("CodeAgent", 3.7, 3.4, Theme.blue),
        ("DomainAgent", 6.2, 3.4, Theme.green),
        ("DatabaseAgent", 8.7, 3.4, Theme.amber),
        ("CriticAgent", 2.45, 5.0, Theme.red),
        ("JudgeAgent", 7.45, 5.0, Theme.card_cyan),
    ]
    for name, x, y, color in agents:
        fill = Theme.panel if color != Theme.card_cyan else Theme.card_cyan
        txt_color = Theme.text if fill == Theme.panel else Theme.ink
        add_card(slide, x, y, 2.0, 0.88, fill=fill)
        add_text(slide, x + 0.12, y + 0.2, 1.76, 0.16, name, size=13, bold=True, color=txt_color, align="center")
    for x, y in [(2.2, 3.4), (4.7, 3.4), (7.2, 3.4), (9.7, 3.4), (3.45, 5.0), (8.45, 5.0)]:
        add_connector(slide, 6.65, 2.75, x, y, Theme.text_dim)
    add_text(slide, 1.18, 5.95, 10.6, 0.16, "命令流：ProblemAnalysisAgent -> 专家 Agent -> Critic 质疑 / Judge 收敛 -> 主 Agent 汇总", size=12, color=Theme.text_dim, align="center")
    add_card(slide, 0.95, 6.2, 11.3, 0.52, fill=Theme.panel_2)
    add_text(slide, 1.15, 6.36, 10.8, 0.14, "硬约束：先有 agent_command_issued，再有专家执行；工具调用与 Skill 命中都必须可审计。", size=12, color=Theme.text)
    add_footer(slide, "docs/agents/agent-catalog.md + docs/agents/protocol-contracts.md")


def slide_tooling(prs: Presentation, page: int) -> None:
    """构建tooling对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "7) Tool / Skill / Connector", "Capability and governance layer", page)
    add_text(slide, 0.55, 1.0, 9.2, 0.26, "结论：项目把“真正执行动作”“方法论注入”“外部平台适配”分成 Tool、Skill、Connector 三层。", size=22, bold=True)
    columns = [
        ("Tool", "真实执行动作", ["读取本地日志", "搜索代码仓", "查询 PostgreSQL", "解析责任田资产"], Theme.card_cyan),
        ("Skill", "分析套路注入", ["RCA 方法论", "角色专属能力提示", "命中审计", "字符预算与注入摘要"], Theme.card_mint),
        ("Connector", "外部系统适配", ["日志平台", "Git 平台", "Grafana / APM", "CMDB / 告警 / 工单"], Theme.card_warm),
    ]
    xs = [0.9, 4.45, 8.0]
    for (title, subtitle, items, fill), x in zip(columns, xs):
        add_card(slide, x, 1.85, 3.0, 4.45, fill=fill, border=Theme.border)
        add_text(slide, x + 0.18, 2.08, 2.65, 0.18, title, size=18, bold=True, color=Theme.ink, align="center")
        add_text(slide, x + 0.18, 2.4, 2.65, 0.14, subtitle, size=11, color=Theme.ink, align="center")
        add_bullets(slide, x + 0.22, 2.9, 2.5, 2.9, items, size=12, color=Theme.ink)
    add_footer(slide, "services/agent_tool_context_service.py + services/agent_skill_service.py + runtime/connectors")


def slide_asset_chain(prs: Presentation, page: int) -> None:
    """构建资产chain对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "8) 责任田资产与数据库链路", "Asset mapping to database reasoning", page)
    add_text(slide, 0.55, 1.0, 9.4, 0.26, "结论：责任田资产不是展示用资料，而是主 Agent 调度和 DatabaseAgent 深挖数据库证据的关键上下文。", size=22, bold=True)
    chain = [
        ("特性", 0.72, Theme.card_cyan),
        ("领域", 1.95, Theme.card_soft),
        ("聚合根", 3.18, Theme.card_mint),
        ("API 接口", 4.55, Theme.card_warm),
        ("代码清单", 5.95, Theme.card_cyan),
        ("数据库表", 7.35, Theme.card_soft),
        ("依赖服务", 8.78, Theme.card_mint),
        ("监控清单", 10.18, Theme.card_warm),
    ]
    for label, x, fill in chain:
        add_card(slide, x, 1.95, 1.05, 0.82, fill=fill, border=Theme.border)
        add_text(slide, x + 0.05, 2.23, 0.95, 0.14, label, size=11, bold=True, color=Theme.ink, align="center")
    for x1, x2 in [(1.77, 1.95), (3.0, 3.18), (4.23, 4.55), (5.6, 5.95), (7.0, 7.35), (8.4, 8.78), (9.83, 10.18)]:
        add_connector(slide, x1, 2.36, x2, 2.36, Theme.blue)
    add_card(slide, 0.85, 3.35, 5.35, 2.45, fill=Theme.panel)
    add_text(slide, 1.1, 3.62, 3.0, 0.18, "主 Agent 消费资产映射", size=14, bold=True, color=Theme.cyan)
    add_bullets(
        slide,
        1.1,
        3.92,
        4.6,
        1.7,
        [
            "接口 URL 命中领域、聚合根和 owner",
            "把代码清单、数据库表、依赖服务整理为命令上下文",
            "决定该把问题交给 Log / Code / Domain / Database 哪些 Agent",
        ],
    )
    add_card(slide, 6.55, 3.35, 5.75, 2.45, fill=Theme.panel_2)
    add_text(slide, 6.8, 3.62, 3.2, 0.18, "DatabaseAgent 深挖链路", size=14, bold=True, color=Theme.amber)
    add_bullets(
        slide,
        6.8,
        3.92,
        4.95,
        1.7,
        [
            "主 Agent 将映射到的表信息传入 DatabaseAgent",
            "查询表结构、字段、索引、慢 SQL、Top SQL、session、锁等待",
            "把数据库证据与代码事务边界、日志时序合并成证据链",
        ],
    )
    add_footer(slide, "docs/wiki/code_wiki_v2.md#14-责任田资产体系 + #15-数据库能力")


def slide_flow(prs: Presentation, page: int) -> None:
    """构建flow对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "9) 端到端分析主流程", "Incident to report", page)
    add_text(slide, 0.55, 1.0, 8.9, 0.26, "结论：完整流程遵循“输入 -> 资产映射 -> 命令分发 -> 专家取证 -> 裁决 -> 报告”的顺序。", size=22, bold=True)
    steps = [
        ("用户输入", 0.75, Theme.card_cyan),
        ("会话创建", 2.15, Theme.card_soft),
        ("资产映射", 3.55, Theme.card_mint),
        ("命令分发", 4.95, Theme.card_warm),
        ("专家取证", 6.35, Theme.card_cyan),
        ("Judge 收敛", 7.75, Theme.card_soft),
        ("报告输出", 9.25, Theme.card_mint),
        ("前端回放", 10.65, Theme.card_warm),
    ]
    for label, x, fill in steps:
        add_card(slide, x, 2.0, 1.1, 0.95, fill=fill, border=Theme.border)
        add_text(slide, x + 0.06, 2.33, 0.98, 0.16, label, size=11, bold=True, color=Theme.ink, align="center")
    for x1, x2 in [(1.85, 2.15), (3.25, 3.55), (4.65, 4.95), (6.05, 6.35), (7.45, 7.75), (8.85, 9.25), (10.35, 10.65)]:
        add_connector(slide, x1, 2.48, x2, 2.48, Theme.blue)
    add_card(slide, 1.0, 3.55, 11.1, 2.15, fill=Theme.panel)
    add_bullets(
        slide,
        1.28,
        3.9,
        10.5,
        1.4,
        [
            "资产映射是前置步骤，责任田命中后才把接口、领域、聚合根、代码、数据库表等上下文传给主 Agent。",
            "主 Agent 下达命令后，专家 Agent 才根据命令决定是否调用工具；无命令或开关关闭时走默认逻辑。",
            "前端通过 WebSocket 实时展示阶段变化、Agent 发言、工具 I/O、结论与报告刷新。",
        ],
    )
    add_footer(slide, "docs/wiki/code_wiki_v2.md#6-一次分析请求的端到端流程")


def slide_case(prs: Presentation, page: int) -> None:
    """构建case对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "10) 真实案例映射", "order-502-db-lock", page)
    add_text(slide, 0.55, 1.0, 8.8, 0.26, "结论：这套架构的价值，在于能把 `/orders 502 + DB lock` 这类故障拆成资产、代码、数据库、日志四条证据链。", size=22, bold=True)
    add_card(slide, 0.75, 1.8, 3.9, 4.8, fill=Theme.panel)
    add_text(slide, 1.0, 2.08, 2.8, 0.18, "输入现象", size=14, bold=True, color=Theme.red)
    add_bullets(
        slide,
        1.0,
        2.38,
        3.1,
        3.7,
        [
            "/api/v1/orders 返回 502",
            "order-service CPU 320%~380%",
            "Hikari pending threads 400+",
            "DB 活跃连接 100/100 打满",
            "MySQL lock wait timeout exceeded",
        ],
    )
    add_card(slide, 4.95, 1.8, 3.6, 4.8, fill=Theme.panel_2)
    add_text(slide, 5.2, 2.08, 2.8, 0.18, "资产映射命中", size=14, bold=True, color=Theme.cyan)
    add_bullets(
        slide,
        5.2,
        2.38,
        2.9,
        3.7,
        [
            "领域：订单域",
            "聚合根：OrderAggregate",
            "责任人：order-domain-team / alice",
            "代码：OrderController / OrderAppService",
            "表：t_order / t_order_item / t_order_snapshot",
        ],
    )
    add_card(slide, 8.8, 1.8, 3.8, 4.8, fill=Theme.panel)
    add_text(slide, 9.05, 2.08, 2.8, 0.18, "Agent 分工", size=14, bold=True, color=Theme.green)
    add_bullets(
        slide,
        9.05,
        2.38,
        3.1,
        3.7,
        [
            "LogAgent：确认 502、连接池超时和锁等待时序",
            "CodeAgent：定位 OrderAppService 事务边界与重试链路",
            "DatabaseAgent：关注锁等待、连接池打满、表热点",
            "Judge：收敛成“事务阻塞 -> 连接池耗尽 -> 网关 502”候选链",
        ],
    )
    add_footer(slide, "scripts/smoke-e2e.mjs + backend/examples/assets/domain-aggregate-responsibility.md")


def slide_events(prs: Presentation, page: int) -> None:
    """构建事件对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "11) 事件流、结果与报告", "Realtime visibility", page)
    add_text(slide, 0.55, 1.0, 8.8, 0.26, "结论：用户体验的核心不是最终报告，而是全过程可见、可追责、可回放。", size=22, bold=True)
    add_chart(
        slide,
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        x=0.75,
        y=1.8,
        w=5.5,
        h=3.2,
        categories=["资产", "命令", "分析", "工具", "裁决", "报告"],
        series=[("事件密度", [4, 6, 10, 8, 5, 3])],
    )
    add_card(slide, 6.7, 1.8, 5.7, 3.2, fill=Theme.panel)
    add_text(slide, 6.95, 2.08, 3.5, 0.18, "前端四类展示对象", size=14, bold=True, color=Theme.amber)
    add_bullets(
        slide,
        6.95,
        2.42,
        4.9,
        2.1,
        [
            "资产映射：只展示责任田、接口、领域、表、依赖服务等命中结果",
            "辩论过程：展示 Agent 对话流和工具调用记录",
            "辩论结果：展示主结论、Top-K、证据摘要",
            "报告结果：展示结构化 RCA 报告和验证/修复建议",
        ],
    )
    add_card(slide, 0.8, 5.3, 11.5, 1.0, fill=Theme.panel_2)
    add_text(slide, 1.0, 5.62, 10.9, 0.16, "关键事件：phase_changed / agent_chat_message / agent_tool_context_prepared / agent_tool_io / result / report", size=12, color=Theme.text)
    add_footer(slide, "api/ws_debates.py + core/event_schema.py + frontend/src/pages/Incident/index.tsx")


def slide_report_view(prs: Presentation, page: int) -> None:
    """构建报告view对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "12) 报告与结果视图", "Structured RCA output", page)
    add_text(slide, 0.55, 1.0, 9.2, 0.26, "结论：最终结果不应该是大段 Markdown，而要拆成可比较、可引用、可执行的结构化模块。", size=22, bold=True)
    add_card(slide, 0.75, 1.85, 3.7, 3.95, fill=Theme.panel)
    add_text(slide, 1.0, 2.12, 2.8, 0.18, "Top-K 根因候选", size=14, bold=True, color=Theme.cyan)
    add_metric(slide, 1.0, 2.46, 0.92, "Top1", "0.81", accent=Theme.cyan)
    add_metric(slide, 2.02, 2.46, 0.92, "Top3", "0.93", accent=Theme.green)
    add_metric(slide, 3.04, 2.46, 0.92, "NeedMore", "Low", accent=Theme.amber)
    add_bullets(
        slide,
        1.0,
        3.62,
        3.0,
        1.95,
        [
            "候选 1：事务阻塞导致连接池耗尽",
            "候选 2：库存表热点更新导致锁等待放大",
            "候选 3：应用层重试与线程堆积放大影响",
            "每个候选都应带置信度和支撑证据",
        ],
    )
    add_card(slide, 4.8, 1.85, 3.7, 3.95, fill=Theme.panel_2)
    add_text(slide, 5.05, 2.12, 2.8, 0.18, "证据链", size=14, bold=True, color=Theme.amber)
    add_text(slide, 5.08, 2.46, 2.9, 0.14, "时间顺序应可追溯，证据必须跨源。", size=10, color=Theme.text_dim)
    add_bullets(
        slide,
        5.05,
        2.82,
        3.0,
        2.75,
        [
            "日志证据：502 / Hikari timeout / lock wait timeout",
            "代码证据：OrderAppService 事务边界、重试链路",
            "数据库证据：t_order_item 热点更新、连接池打满",
            "资产证据：接口归属、聚合根、责任团队",
        ],
    )
    add_card(slide, 8.85, 1.85, 3.7, 3.95, fill=Theme.panel)
    add_text(slide, 9.1, 2.12, 2.8, 0.18, "行动建议", size=14, bold=True, color=Theme.green)
    add_bullets(
        slide,
        9.1,
        2.82,
        3.0,
        2.75,
        [
            "立即止血：限流、扩容连接池、终止阻塞事务",
            "修复建议：拆分热点事务、优化索引和 SQL",
            "验证计划：观察 Hikari pending、5xx、锁等待",
            "回滚与复盘：保留证据链和版本差异对比",
        ],
    )
    add_card(slide, 0.85, 6.0, 11.55, 0.62, fill=Theme.panel_2)
    add_text(slide, 1.05, 6.22, 11.0, 0.14, "前端落地形态：辩论结果页展示主结论与 Top-K，报告结果页展示证据链、建议动作、验证计划和对比回放入口。", size=12, color=Theme.text)
    add_footer(slide, "docs/wiki/code_wiki_v2.md#18-报告、证据链与结果组织")


def slide_code_map(prs: Presentation, page: int) -> None:
    """构建codemap对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "16) 代码阅读地图", "How to enter the codebase", page)
    add_text(slide, 0.55, 1.0, 9.2, 0.26, "结论：新人最有效的上手方式不是按目录遍历，而是沿着一次真实请求的调用链读代码。", size=22, bold=True)
    steps = [
        ("main.py", "应用启动与 router 装配", 0.8, Theme.card_cyan),
        ("api/*", "REST / WS 入口", 2.55, Theme.card_soft),
        ("debate_service.py", "业务主编排器", 4.3, Theme.card_mint),
        ("flows/*", "业务旅程组织", 6.25, Theme.card_warm),
        ("langgraph_runtime.py", "运行时主入口", 8.2, Theme.card_cyan),
        ("runtime/langgraph/*", "节点 / 状态 / 路由 / 事件", 10.15, Theme.card_soft),
    ]
    for title, body, x, fill in steps:
        add_box_label(slide, x, 2.0, 1.5, 1.08, title, body, fill=fill)
    for x1, x2 in [(2.3, 2.55), (4.05, 4.3), (6.0, 6.25), (7.95, 8.2), (9.9, 10.15)]:
        add_connector(slide, x1, 2.54, x2, 2.54, Theme.blue)
    add_card(slide, 0.9, 3.75, 5.55, 2.15, fill=Theme.panel)
    add_text(slide, 1.15, 4.02, 2.8, 0.18, "前端入口", size=14, bold=True, color=Theme.amber)
    add_bullets(
        slide,
        1.15,
        4.32,
        4.8,
        1.35,
        [
            "Incident/index.tsx：创建会话、连接 WS、组织事件流",
            "WarRoom/index.tsx：战情态势、证据链与关键结果",
            "Assets/index.tsx：责任田资产维护与查询",
        ],
    )
    add_card(slide, 6.85, 3.75, 5.45, 2.15, fill=Theme.panel_2)
    add_text(slide, 7.1, 4.02, 3.0, 0.18, "推荐阅读顺序", size=14, bold=True, color=Theme.green)
    add_bullets(
        slide,
        7.1,
        4.32,
        4.7,
        1.35,
        [
            "先读 code_wiki_v2，再跟一次 /orders 502 场景",
            "先看调用链和状态，再看 prompt 和 parser 细节",
            "先理解职责边界，再修改 Agent / Tool / UI",
        ],
    )
    add_footer(slide, "docs/wiki/code_wiki_v2.md#2-建议的阅读顺序")


def slide_reliability(prs: Presentation, page: int) -> None:
    """构建reliability对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "14) 可靠性、治理与 Benchmark", "Guardrails and quality", page)
    add_text(slide, 0.55, 1.0, 9.1, 0.26, "结论：要把多 Agent 系统用于生产场景，必须把超时、终态、审计和评测作为系统能力，而不是 Prompt 习惯。", size=22, bold=True)
    add_card(slide, 0.75, 1.85, 3.7, 4.2, fill=Theme.panel)
    add_text(slide, 1.0, 2.12, 2.8, 0.18, "可靠性守护", size=14, bold=True, color=Theme.cyan)
    add_bullets(
        slide,
        1.0,
        2.42,
        3.0,
        3.2,
        [
            "超时、重试、降级",
            "终态保证，避免长期 pending",
            "checkpoint / resume",
            "doom loop guard 防循环",
            "session compaction 控制上下文膨胀",
        ],
    )
    add_card(slide, 4.8, 1.85, 3.7, 4.2, fill=Theme.panel_2)
    add_text(slide, 5.05, 2.12, 2.8, 0.18, "治理与审计", size=14, bold=True, color=Theme.amber)
    add_bullets(
        slide,
        5.05,
        2.42,
        3.0,
        3.2,
        [
            "工具调用请求/返回摘要",
            "Skill 命中来源与注入摘要",
            "风险动作治理与 system card",
            "Agent 协议和命令门禁",
        ],
    )
    add_card(slide, 8.85, 1.85, 3.7, 4.2, fill=Theme.panel)
    add_text(slide, 9.1, 2.12, 2.8, 0.18, "Benchmark", size=14, bold=True, color=Theme.green)
    add_bullets(
        slide,
        9.1,
        2.42,
        3.0,
        3.2,
        [
            "标准 incident 样本",
            "Top-K / 超时率 / 报告质量打分",
            "可接入 CI Gate 阻断回归",
            "支撑提示词、路由和工具变更回归验证",
        ],
    )
    add_footer(slide, "AGENTS.md + docs/agents/reliability-governance.md + backend/app/benchmark")


def slide_extension(prs: Presentation, page: int) -> None:
    """构建extension对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "15) 扩展路径与阅读建议", "How to extend", page)
    add_text(slide, 0.55, 1.0, 8.7, 0.26, "结论：扩展新 Agent、新工具或新平台接入时，要先补协议，再改运行时，再补前端与评测。", size=22, bold=True)
    add_card(slide, 0.75, 1.9, 5.6, 4.45, fill=Theme.panel)
    add_text(slide, 1.0, 2.18, 2.8, 0.18, "新增 Agent 的顺序", size=14, bold=True, color=Theme.cyan)
    add_bullets(
        slide,
        1.0,
        2.48,
        4.9,
        3.1,
        [
            "先更新 docs/agents/agent-catalog.md 和 protocol-contracts.md",
            "再改 specs.py、prompts.py、parsers.py 和节点逻辑",
            "通过 agent_tool_context_service / agent_skill_service 注入能力",
            "最后补前端展示与 benchmark 用例",
        ],
    )
    add_card(slide, 6.75, 1.9, 5.55, 4.45, fill=Theme.panel_2)
    add_text(slide, 7.0, 2.18, 3.2, 0.18, "新同学阅读路径", size=14, bold=True, color=Theme.amber)
    add_bullets(
        slide,
        7.0,
        2.48,
        4.85,
        3.1,
        [
            "后端：main.py -> api -> debate_service -> runtime",
            "前端：Incident -> WarRoom -> Assets -> History",
            "Prompt / Agent：specs -> prompt_builder -> parsers",
            "平台能力：tools -> connectors -> benchmark -> governance",
        ],
    )
    add_footer(slide, "docs/wiki/code_wiki_v2.md#19-如何新增一个-agent")


def slide_summary(prs: Presentation, page: int) -> None:
    """构建摘要对应的幻灯片内容，并完成该页布局与文案写入。"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_topbar(slide, "18) 总结", "Training close", page)
    add_text(slide, 0.55, 1.0, 8.8, 0.26, "结论：这套系统的核心价值，不是“模型替代人”，而是把故障调查流程工程化、结构化和可审计化。", size=22, bold=True)
    add_card(slide, 0.8, 1.9, 3.7, 3.8, fill=Theme.panel)
    add_text(slide, 1.05, 2.18, 2.8, 0.18, "你需要记住的三件事", size=14, bold=True, color=Theme.cyan)
    add_bullets(
        slide,
        1.05,
        2.48,
        3.0,
        2.8,
        [
            "资产映射先行，主 Agent 再调度专家 Agent",
            "过程事件、工具审计、报告结果同等重要",
            "可靠性和治理是系统能力，不是补丁",
        ],
    )
    add_card(slide, 4.85, 1.9, 3.6, 3.8, fill=Theme.panel_2)
    add_text(slide, 5.1, 2.18, 2.8, 0.18, "建议继续深入的代码入口", size=14, bold=True, color=Theme.amber)
    add_bullets(
        slide,
        5.1,
        2.48,
        2.9,
        2.8,
        [
            "backend/app/services/debate_service.py",
            "backend/app/runtime/langgraph_runtime.py",
            "backend/app/runtime/langgraph/*",
            "frontend/src/pages/Incident/index.tsx",
        ],
    )
    add_card(slide, 8.8, 1.9, 3.6, 3.8, fill=Theme.panel)
    add_text(slide, 9.05, 2.18, 2.8, 0.18, "后续可继续优化", size=14, bold=True, color=Theme.green)
    add_bullets(
        slide,
        9.05,
        2.48,
        2.9,
        2.8,
        [
            "更多真实数据源 Connector",
            "更强的报告可视化和回放",
            "更细颗粒度的 Benchmark Gate",
            "更统一的工作台与战情页心智",
        ],
    )
    add_footer(slide, "docs/wiki/code_wiki_v2.md")


def build_slides_md() -> str:
    """构建幻灯片md相关产物或页面内容。"""
    
    sections = [
        ("封面", "生产问题根因分析系统架构培训", ["团队内部培训", "war-room briefing 风格", "基于 code_wiki_v2 组织内容"]),
        ("系统目标与业务边界", "本项目是根因分析系统，不是聊天机器人", ["目标：资产映射、证据收集、Agent 协作、可回放、可审计", "边界：主 Agent 调度，专家 Agent 分析，工具调用受门禁控制"]),
        ("总体分层架构", "系统分为前端工作台、API/Service/Flow、LangGraph Runtime、Tool/Skill/Connector 四层", ["Model 与 System 分层", "UI 与 Runtime 解耦"]),
        ("前端工作台架构", "前端按录入、观察、回放、治理组织", ["Incident / WarRoom / Workbench / History / Assets", "事件流、结果、报告分轨渲染"]),
        ("后端接入与服务流", "后端通过 API -> Service -> Flow -> Runtime 串起业务与执行", ["DebateService 是业务入口", "Flow 串联用户旅程"]),
        ("LangGraph 运行时", "运行时模块化拆分图构建、状态、节点执行、事件分发和可靠性", ["GraphBuilder / AgentRunner / State / EventDispatcher / Checkpointer"]),
        ("多 Agent 协作机制", "主 Agent 指挥，专家 Agent 执行，Judge 收敛", ["命令先行", "工具调用必须可审计"]),
        ("Tool / Skill / Connector", "把执行动作、方法论注入、外部平台适配分层治理", ["Tool 真执行", "Skill 注入套路", "Connector 接平台"]),
        ("责任田资产与数据库链路", "资产映射直接决定主 Agent 调度和 DatabaseAgent 的分析范围", ["从接口命中领域/聚合根/代码/表", "数据库表信息继续喂给 DatabaseAgent"]),
        ("端到端分析流程", "输入 -> 资产映射 -> 命令分发 -> 专家取证 -> 裁决 -> 报告", ["资产映射先行", "前端通过 WebSocket 实时展示全过程"]),
        ("真实案例映射", "架构要能落到真实故障：/orders 502 + DB lock", ["资产映射命中领域/聚合根/代码/表", "多 Agent 从日志、代码、数据库三个方向收敛"]),
        ("事件流、结果与报告", "过程可见比最终报告更关键", ["阶段事件、Agent 发言、工具 I/O、结果和报告分别展示"]),
        ("报告与结果视图", "RCA 结果要结构化，不应只是一段 Markdown", ["Top-K 根因、证据链、行动建议分模块呈现", "支持对比、验证和回放"]),
        ("可靠性、治理与 Benchmark", "超时、终态、审计、评测必须系统化", ["checkpoint/resume", "doom loop guard", "Benchmark Gate"]),
        ("扩展路径与阅读建议", "先补协议，再改实现，再补前端与评测", ["新增 Agent / Tool / Connector 的标准入口", "推荐阅读路径"]),
        ("代码阅读地图", "按真实请求调用链读代码，比按目录遍历更高效", ["main.py -> api -> debate_service -> runtime", "Incident / WarRoom / Assets 是前端主要入口"]),
        ("总结", "系统价值在于把故障调查流程工程化、结构化和可审计化", ["资产映射先行", "过程与结果同等重要", "可靠性和治理必须系统化"]),
    ]
    lines = ["# 生产问题根因分析系统架构培训", ""]
    for idx, (title, assertion, bullets) in enumerate(sections, start=1):
        lines.append(f"## {idx}. {title}")
        lines.append(f"- 结论：{assertion}")
        for item in bullets:
            lines.append(f"- {item}")
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def build_notes() -> str:
    """构建备注相关产物或页面内容。"""
    
    notes = [
        "# Speaker Notes",
        "",
        "## 使用方式",
        "- 总时长建议 15-20 分钟。",
        "- 每页控制在 45-60 秒，重点放在架构边界和调用链，不展开历史细节。",
        "",
        "## 讲解主线",
        "1. 先解释系统目标和边界，避免被理解成普通聊天应用。",
        "2. 再讲四层分层架构，让团队建立统一心智。",
        "3. 用端到端流程把前端、后端、运行时、多 Agent 串起来。",
        "4. 最后讲可靠性、治理和扩展路径，说明项目为什么能继续演进。",
        "",
        "## 页面提示",
        "- LangGraph 运行时页：强调 LLM 只负责推理，系统负责控制。",
        "- 多 Agent 页：强调命令先行，不是自由放任讨论。",
        "- Tool/Skill/Connector 页：重点解释三层差异，避免混淆。",
        "- 责任田资产页：强调资产不是展示数据，而是主 Agent 和 DatabaseAgent 的工作输入。",
        "- 真实案例页：把抽象架构拉回真实故障，帮助团队理解为什么要做资产映射和 DatabaseAgent。",
        "- 报告页：强调为什么报告要结构化、可比较，而不是简单 Markdown 渲染。",
        "- 可靠性页：强调 pending、timeout、loop 都是系统级问题，不应只靠 prompt 缓解。",
        "- 代码阅读地图页：用来带新人建立“从请求到代码”的入口心智。",
        "",
        "## 可选加讲",
        "- 如果培训对象偏前端，可在 Incident / WarRoom 页面多停留 2 分钟。",
        "- 如果培训对象偏后端，可展开 DebateService 和 runtime/langgraph 目录讲解。",
        "",
    ]
    return "\n".join(notes)


def build_refs() -> str:
    """构建参考资料相关产物或页面内容。"""
    
    refs = [
        "# References",
        "",
        "- `/Users/neochen/multi-agent-cli_v2/docs/wiki/code_wiki_v2.md`",
        "- `/Users/neochen/multi-agent-cli_v2/AGENTS.md`",
        "- `/Users/neochen/multi-agent-cli_v2/docs/agents/agent-catalog.md`",
        "- `/Users/neochen/multi-agent-cli_v2/docs/agents/protocol-contracts.md`",
        "- `/Users/neochen/multi-agent-cli_v2/docs/agents/reliability-governance.md`",
        "- `/Users/neochen/multi-agent-cli_v2/backend/app/services/debate_service.py`",
        "- `/Users/neochen/multi-agent-cli_v2/backend/app/runtime/langgraph_runtime.py`",
        "- `/Users/neochen/multi-agent-cli_v2/backend/app/runtime/langgraph/`",
        "- `/Users/neochen/multi-agent-cli_v2/frontend/src/pages/Incident/index.tsx`",
        "- `/Users/neochen/multi-agent-cli_v2/frontend/src/pages/WarRoom/index.tsx`",
        "",
    ]
    return "\n".join(refs)


def build_readme() -> str:
    """构建说明文档相关产物或页面内容。"""
    
    return "\n".join(
        [
            "# Project Intro War-room PPT",
            "",
            "- `project-intro-warroom.pptx`: 团队内部培训版项目介绍 PPT",
            "- `slides.md`: 页结构与核心断言",
            "- `notes.md`: 演讲备注",
            "- `refs.md`: 来源文件",
            "",
            "生成脚本：",
            "- `/Users/neochen/multi-agent-cli_v2/scripts/generate_project_intro_warroom_ppt.py`",
            "",
        ]
    )


def build_deck() -> Presentation:
    """构建deck相关产物或页面内容。"""
    
    prs = new_prs()
    slide_cover(prs)
    slide_agenda(prs, 2)
    slide_value(prs, 3)
    slide_system_arch(prs, 4)
    slide_frontend(prs, 5)
    slide_backend(prs, 6)
    slide_runtime(prs, 7)
    slide_agents(prs, 8)
    slide_tooling(prs, 9)
    slide_asset_chain(prs, 10)
    slide_flow(prs, 11)
    slide_case(prs, 12)
    slide_events(prs, 13)
    slide_report_view(prs, 14)
    slide_reliability(prs, 15)
    slide_extension(prs, 16)
    slide_code_map(prs, 17)
    slide_summary(prs, 18)
    return prs


def main() -> None:
    """执行脚本主流程，串联参数解析、内容生成与结果输出。"""
    
    ensure_dirs()
    prs = build_deck()
    prs.save(OUT_PPTX)
    OUT_SLIDES.write_text(build_slides_md(), encoding="utf-8")
    OUT_NOTES.write_text(build_notes(), encoding="utf-8")
    OUT_REFS.write_text(build_refs(), encoding="utf-8")
    OUT_README.write_text(build_readme(), encoding="utf-8")
    print(f"Generated: {OUT_PPTX}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
